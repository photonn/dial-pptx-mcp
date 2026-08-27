"""
Deterministic brand checks, driven entirely by a profile the deployment owns.

`deck_validation.py` asks whether the file is well-formed; `visual_qa.py` asks
whether it looks right. Neither can answer "is this our brand" — that is not a
property of the deck, it is a property of the organisation the deck is for, and
this repository has never known which organisation that is. It does not learn
here either: **the repo ships the engine, the deployment supplies the rules.**

So this module knows the *shape* of a brand rule — an allowed set of fonts, a
minimum text size, a palette, a safe area, a piece of chrome that must be
present, how often a background family may repeat — and nothing about any
particular brand's values.

Where the rules live
--------------------
In DIAL file storage, beside the template and the icons, and the server is told
only their **file names**:

- ``BRAND_PROFILE_FILE``        e.g. ``brand_profile.json``
- ``BRAND_REFERENCE_DECK_FILE`` e.g. ``brand_reference.pptx``

Not a bucket path, because buckets change: a file moves, a deployment is
rebuilt, storage is reorganised, and a server pinned to ``files/{bucket}/...``
stops finding its own rules. The name is the stable half of the contract. The
orchestrator holds the other half — it can see which files it has access to, so
it resolves the name to a URL and passes it to ``attach_brand_profile`` once
per deck (``tools/brand_tools.py``), which is also what makes the file readable
at all: DIAL Quick Apps grants the request's key access to files named by
``dial_url`` parameters, and nothing else.

That makes brand context a property of the *deck*, held in the store beside the
Presentation. It is per-caller by construction, which is the safe direction: a
profile's ``review_notes`` become prompt text for the vision reviewer, so one
caller's file must never be able to steer another caller's review. The file
name check is the integrity half — a profile is only accepted from a file
called what the operator said it would be called.

Unset ``BRAND_PROFILE_FILE`` means the feature does not exist: no tools, no
export summary, no reviewer focus.

Why deterministic checks at all, when a vision model reviews every slide?
Because these ones are cheap, need no model call, and produce an exact
"slide 4, shape 7" pointer. A reviewer can tell you text looks small; it cannot
tell you it is 12pt against a 14pt floor, and it will not notice that #0060BE
is not #005DB9. The two semantic rules that *are* model work — a headline must
carry a message, a slide must not be plain text on white — live in the
profile's ``review_notes`` and are fed to the reviewer as extra focus rather
than implemented here.

Findings use the same report shape as deck_validation, so an agent reads one
kind of problem dict from both passes.
"""
import io
import json
import os

# Same problem/report shape as the structural pass, on purpose: the consumer is
# one agent deciding what to call next, and two dialects would cost it.
from deck_validation import (ERROR, INFO, WARNING, EMU_PER_INCH, _Report,
                             _corner_pixels, _solid_fill_rgb)
from logging_utils import get_logger

logger = get_logger("brand_validation")

SEVERITIES = (ERROR, WARNING, INFO)
DEFAULT_SEVERITY = WARNING

# A background this dark is a "dark" slide: white type on a brand colour. The
# threshold is on perceived luminance, not on any particular colour, so it
# holds for a navy, a deep green or a black master alike.
_DARK_LUMINANCE = 0.5

# A shape covering this much of the slide at the bottom of the z-order is the
# background as far as the eye is concerned, whatever the slide's own fill says.
_BACKGROUND_COVERAGE = 0.9

# The brand files are named by the operator but fetched from a URL the caller
# supplies, so their size is worth bounding like any other caller-driven
# download. A rules file is a few kilobytes; a deck of example slides is not.
MAX_PROFILE_BYTES = 1024 * 1024
MAX_REFERENCE_BYTES = 50 * 1024 * 1024


class BrandConfigError(RuntimeError):
    """A brand file was named but could not be used.

    Distinct from "no brand configured", which is not an error at all: this
    one always means someone — the operator or the orchestrator — has to do
    something about it.
    """


# ---- Configuration ----

def profile_file_name():
    """The file name this deployment expects its brand rules in, or None."""
    return (os.environ.get("BRAND_PROFILE_FILE") or "").strip() or None


def reference_deck_file_name():
    """The file name of the brand's reference deck, or None."""
    return (os.environ.get("BRAND_REFERENCE_DECK_FILE") or "").strip() or None


def enabled():
    """Whether this deployment has brand rules at all.

    Keyed on the variable, never on a file having been supplied: the tools are
    registered at startup, long before any deck exists.
    """
    return profile_file_name() is not None


def file_name_of(url):
    """The bare file name a DIAL reference ends in.

    Query strings and trailing slashes are stripped; the bucket, folder and
    host are all deliberately ignored, because those are the parts that move.
    """
    ref = (url or "").strip().split("?", 1)[0].rstrip("/")
    return ref.rsplit("/", 1)[-1]


def check_file_name(url, expected, what):
    """Raise unless `url` names the file the operator configured.

    The orchestrator picks which accessible file to pass, so this is the one
    check standing between "the brand profile" and "any JSON file the caller
    happens to have". It compares names only — the path is the caller's
    business, the name is the operator's.
    """
    actual = file_name_of(url)
    if actual.lower() != expected.lower():
        raise BrandConfigError(
            f"'{actual}' is not this server's {what}: it expects a file named "
            f"'{expected}'. Pass the DIAL file URL of that file — the one you "
            f"have access to — rather than another file, and if you cannot "
            f"find it, say so instead of substituting a different one."
        )


# ---- Parsing ----

def parse_profile(blob, source=""):
    """A brand profile dict from the bytes of a profile file."""
    if len(blob) > MAX_PROFILE_BYTES:
        raise BrandConfigError(
            f"The brand profile {source} is {len(blob) / 1048576:.1f} MB, far "
            f"larger than a rules file should be. This is probably not the "
            f"right file."
        )
    try:
        profile = json.loads(blob.decode("utf-8"))
    except (UnicodeDecodeError, ValueError) as e:
        raise BrandConfigError(
            f"The brand profile {source} is not valid JSON ({e}). It must "
            f"match brand_profile.schema.json."
        )
    if not isinstance(profile, dict):
        raise BrandConfigError(
            f"The brand profile {source} is not a JSON object. It must match "
            f"brand_profile.schema.json."
        )
    return profile


def open_reference_deck(blob, source=""):
    """A Presentation from the bytes of the brand's reference deck."""
    from pptx import Presentation

    if len(blob) > MAX_REFERENCE_BYTES:
        raise BrandConfigError(
            f"The brand reference deck {source} is "
            f"{len(blob) / 1048576:.0f} MB, too large to hold open."
        )
    try:
        return Presentation(io.BytesIO(blob))
    except Exception as e:
        raise BrandConfigError(
            f"The brand reference deck {source} is not a readable .pptx "
            f"({e}). A .potx will not open — export it as .pptx."
        )


def review_focus(profile):
    """The profile's semantic rules, as a focus string for the vision reviewer.

    These are the checks that cannot be made deterministic — whether a headline
    states a message, whether a slide is a wall of plain text — so they are
    carried as prose in the profile and handed to the model that can judge them.
    """
    if not profile:
        return None
    notes = profile.get("review_notes") or []
    notes = [str(n).strip() for n in notes if str(n).strip()]
    return " ".join(notes) or None


# ---- Geometry and colour helpers ----

def _emu(inches):
    return int(round(float(inches) * EMU_PER_INCH))


def _inches(emu):
    return round(emu / EMU_PER_INCH, 2)


def _luminance(rgb):
    r, g, b = (c / 255 for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _near(rgb, other, tolerance):
    return all(abs(a - b) <= tolerance for a, b in zip(rgb, other))


def _on_palette(rgb, palette, tolerance):
    return any(_near(rgb, tuple(entry), tolerance) for entry in palette)


def _rgb_of(color):
    """An (r, g, b) tuple for an explicitly-set colour, else None.

    A theme colour is by definition on the brand's palette — the theme comes
    from the brand's own template — so only literal RGB is worth checking, and
    that is exactly what an agent passing `[r, g, b]` to a tool produces.
    """
    from pptx.enum.dml import MSO_COLOR_TYPE

    try:
        if color is None or color.type != MSO_COLOR_TYPE.RGB:
            return None
        return tuple(color.rgb)
    except Exception:
        return None


def _shape_rgbs(shape):
    """Every literal colour the shape sets: fill, outline and run colours.

    Yields (label, rgb) pairs; a shape that inherits everything yields nothing.
    """
    from pptx.enum.dml import MSO_FILL

    try:
        if shape.fill.type == MSO_FILL.SOLID:
            rgb = _rgb_of(shape.fill.fore_color)
            if rgb:
                yield "fill", rgb
    except Exception:
        pass
    try:
        rgb = _rgb_of(shape.line.color)
        if rgb:
            yield "outline", rgb
    except Exception:
        pass
    for run in _runs(shape):
        rgb = _rgb_of(run.font.color)
        if rgb:
            yield "text", rgb


def _runs(shape):
    """Every run of text in a shape, text frame and table cells alike."""
    frames = []
    if getattr(shape, "has_text_frame", False):
        frames.append(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                frames.append(cell.text_frame)
    for frame in frames:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                yield run


def _placeholder_idx(shape):
    try:
        if not shape.is_placeholder:
            return None
        return shape.placeholder_format.idx
    except Exception:
        return None


def _layout_name(slide):
    try:
        return (slide.slide_layout.name or "").strip()
    except Exception:
        return ""


def _background_rgb(slide):
    """The colour the audience sees behind the slide, or None if unresolvable.

    Checked in the order it actually paints: a full-bleed shape at the bottom
    of the z-order first (which is how most agents make a coloured slide), then
    the slide's own background, then the layout's and the master's.
    """
    slide_area = None
    try:
        pres_part = slide.part.package.presentation_part.presentation
        slide_area = pres_part.slide_width * pres_part.slide_height
    except Exception:
        pass

    for shape in slide.shapes:
        rgb = _solid_fill_rgb(shape)
        if rgb is None:
            break  # only shapes at the very bottom of the stack can be the bg
        if slide_area and shape.width and shape.height:
            if (shape.width * shape.height) / slide_area >= _BACKGROUND_COVERAGE:
                return rgb
        break

    from pptx.enum.dml import MSO_FILL

    for source in (slide, slide.slide_layout,
                   getattr(slide.slide_layout, "slide_master", None)):
        if source is None:
            continue
        try:
            fill = source.background.fill
            if fill.type == MSO_FILL.SOLID:
                rgb = _rgb_of(fill.fore_color)
                if rgb:
                    return rgb
        except Exception:
            continue
    return None


def classify_background(slide):
    """"light" or "dark" — which family of the profile a slide is judged by.

    An unresolvable background is treated as light: a white slide is the
    default everywhere, and the light family is the stricter of the two in
    every profile we have seen.
    """
    rgb = _background_rgb(slide)
    if rgb is None:
        return "light"
    return "dark" if _luminance(rgb) < _DARK_LUMINANCE else "light"


# ---- Individual checks ----

def _severity(profile, section, default=DEFAULT_SEVERITY):
    value = str((section or {}).get("severity")
                or profile.get("severity") or default).lower()
    return value if value in SEVERITIES else default


def check_font_face(shape, rules, report, slide_index, shape_index, profile):
    allowed = [str(f).strip().lower() for f in (rules or {}).get("allowed", [])]
    if not allowed:
        return
    severity = _severity(profile, rules)
    seen = set()
    for run in _runs(shape):
        name = (run.font.name or "").strip()
        if not name or name.lower() in allowed or name.lower() in seen:
            continue
        seen.add(name.lower())
        report.add(severity, "brand_font_not_allowed",
                   f"Slide {slide_index}, shape {shape_index} sets the font "
                   f"'{name}', which is not one of the brand's fonts "
                   f"({', '.join((rules or {}).get('allowed', []))}).",
                   "Set the allowed font with manage_fonts, or clear the "
                   "explicit font so the text inherits the template's.",
                   slide=slide_index, shape=shape_index)


def check_min_size(shape, minimum, exempt_idx, report, slide_index,
                   shape_index, profile):
    if not minimum:
        return
    if _placeholder_idx(shape) in set(exempt_idx or []):
        return
    smallest = None
    for run in _runs(shape):
        size = run.font.size
        if size is None or not run.text.strip():
            continue
        points = size.pt
        if points < minimum and (smallest is None or points < smallest):
            smallest = points
    if smallest is None:
        return
    report.add(_severity(profile, None), "brand_text_too_small",
               f"Slide {slide_index}, shape {shape_index} has text at "
               f"{smallest:g}pt, below the brand's {minimum}pt floor.",
               f"Raise it to at least {minimum}pt with manage_fonts, or cut "
               "the text so it fits at that size.",
               slide=slide_index, shape=shape_index)


def check_palette(shape, palette, tolerance, report, slide_index, shape_index,
                  profile):
    if not palette:
        return
    for label, rgb in _shape_rgbs(shape):
        if _on_palette(rgb, palette, tolerance):
            continue
        report.add(_severity(profile, None), "brand_colour_off_palette",
                   f"Slide {slide_index}, shape {shape_index} sets a {label} "
                   f"colour of {list(rgb)}, which is not on the brand palette.",
                   "Use a palette colour, or drop the explicit colour so the "
                   "shape inherits the template's theme colours — theme "
                   "colours are always on-palette.",
                   slide=slide_index, shape=shape_index)


def check_safe_area(shape, area, chrome_idx, report, slide_index, shape_index,
                    profile):
    """Content must stay inside the brand's frame, not merely on the slide.

    deck_validation already reports a shape past the slide edge; a brand's safe
    area sits well inside that, and crossing it is invisible to every other
    pass. Chrome — logo, page number, source line — is placed *by the template*
    outside the frame on purpose, so the profile names those placeholder idx
    values and they are skipped.
    """
    if not area:
        return
    if _placeholder_idx(shape) in set(chrome_idx or []):
        return
    if None in (shape.left, shape.top, shape.width, shape.height):
        return
    if shape.width <= 0 or shape.height <= 0:
        return

    left, top = _emu(area["x"]), _emu(area["y"])
    right, bottom = left + _emu(area["w"]), top + _emu(area["h"])
    tolerance = _emu(area.get("tolerance", 0.02))
    crossed = [name for name, hit in (
        ("left", shape.left < left - tolerance),
        ("top", shape.top < top - tolerance),
        ("right", shape.left + shape.width > right + tolerance),
        ("bottom", shape.top + shape.height > bottom + tolerance)) if hit]
    if not crossed:
        return
    report.add(_severity(profile, None), "brand_outside_safe_area",
               f"Slide {slide_index}, shape {shape_index} ('{shape.name}') "
               f"crosses the {', '.join(crossed)} edge of the brand safe area "
               f"({area['x']}, {area['y']} in, {area['w']}x{area['h']} in): it "
               f"sits at {_inches(shape.left)}, {_inches(shape.top)} in and is "
               f"{_inches(shape.width)}x{_inches(shape.height)} in.",
               "Move or resize it inside the safe area with move_shape / "
               "resize_shape, or with visual_repair_slides.",
               slide=slide_index, shape=shape_index)


def check_image_alpha(shape, required, report, slide_index, shape_index,
                      family, profile):
    """A flattened icon on a coloured slide.

    deck_validation catches the same defect against a coloured *shape* behind
    the picture; this covers the other half — a slide whose whole background is
    a brand colour, where there is no shape to overlap.
    """
    if not required or family != "dark":
        return
    if shape.shape_type != 13:  # PICTURE
        return
    try:
        corners = _corner_pixels(shape.image.blob)
    except Exception:
        return
    if not corners or not all(pixel[3] >= 250 for pixel in corners):
        return
    report.add(_severity(profile, None), "brand_opaque_image_on_fill",
               f"Slide {slide_index}, shape {shape_index} ('{shape.name}') has "
               f"an opaque background and sits on a coloured slide, so it "
               f"shows as a rectangle around the artwork.",
               "Use the transparent variant of the image.",
               slide=slide_index, shape=shape_index)


def _matches_chrome(shape, spec):
    kind = (spec or {}).get("type")
    if kind == "placeholder":
        return _placeholder_idx(shape) == spec.get("idx")
    if kind == "picture":
        if shape.shape_type != 13:
            return False
        ratios = spec.get("aspect_ratios") or []
        if not ratios:
            return True
        if not shape.width or not shape.height:
            return False
        actual = shape.width / shape.height
        tolerance = spec.get("aspect_tolerance", 0.03)
        return any(abs(actual - float(r)) <= tolerance for r in ratios)
    if kind == "line":
        if shape.top is None or shape.height is None:
            return False
        thin = shape.height <= _emu(spec.get("max_height", 0.06))
        y = spec.get("y_in")
        if y is None:
            return thin
        return thin and abs(shape.top - _emu(y)) <= _emu(
            spec.get("tolerance_in", 0.05))
    return False


def check_required_chrome(slide, required, chrome_shapes, report, slide_index,
                          profile):
    """Every piece of chrome the family requires must actually be there.

    Where it is looked for follows what PowerPoint renders. Drawn chrome —
    a separator rule, a logo — is inherited from the layout and the master and
    shows without existing on the slide, so all three are searched. A
    *placeholder* is the opposite: a slide-number placeholder defined on the
    layout renders nothing until the slide carries its own copy (which is what
    add_slide_numbers is for), so those are looked for on the slide alone.
    """
    if not required:
        return
    sources = {
        "slide": list(slide.shapes),
    }
    try:
        sources["layout"] = list(slide.slide_layout.shapes)
        sources["master"] = list(slide.slide_layout.slide_master.shapes)
    except Exception:
        pass

    for name in required:
        spec = (chrome_shapes or {}).get(name)
        if not spec:
            continue
        where = ["slide"] if spec.get("type") == "placeholder" \
            else ["slide", "layout", "master"]
        found = any(_matches_chrome(shape, spec)
                    for scope in where for shape in sources.get(scope, []))
        if found:
            continue
        fix = ("Call add_slide_numbers to copy the template's slide-number "
               "placeholder onto the slides."
               if spec.get("type") == "placeholder" else
               f"Build this slide from a layout that carries the {name}, or "
               "duplicate a slide that already has it with duplicate_slide.")
        report.add(_severity(profile, None), "brand_chrome_missing",
                   f"Slide {slide_index} is missing the brand's required "
                   f"'{name}'.",
                   fix, slide=slide_index)


def check_family_mix(families, limit, report, profile):
    """A deck that never changes background is the flattest way to be on-brand
    and off-brand at once: every rule passes and it reads as a wall."""
    if not limit or limit < 1:
        return
    run_family, run_start = None, 0
    for index, family in enumerate(families + [None]):
        if family == run_family:
            continue
        length = index - run_start
        if run_family is not None and length > limit:
            report.add(_severity(profile, None), "brand_monotonous_deck",
                       f"Slides {run_start} to {index - 1} are all "
                       f"'{run_family}' backgrounds — {length} in a row, above "
                       f"the brand's limit of {limit}.",
                       "Rebuild one of them on a layout from the other "
                       "background family with duplicate_slide, so the deck "
                       "has some rhythm.",
                       slide=run_start)
        run_family, run_start = family, index


# ---- Entry point ----

def validate_brand(pres, profile):
    """Check a deck against a caller-supplied brand profile.

    The profile is data, not code: this module knows the shape of a brand rule
    and nothing about any particular brand. Returns the same report dict shape
    as deck_validation.validate_presentation.
    """
    report = _Report()
    families = []

    for slide_index, slide in enumerate(pres.slides):
        family = classify_background(slide)
        families.append(family)
        rules = (profile.get("families") or {}).get(family) or {}
        exempt = {name.strip().lower()
                  for name in profile.get("exempt_layouts") or []}
        layout_exempt = _layout_name(slide).lower() in exempt

        for shape_index, shape in enumerate(slide.shapes):
            try:
                check_font_face(shape, profile.get("fonts"), report,
                                slide_index, shape_index, profile)
                check_min_size(shape, profile.get("min_font_pt"),
                               profile.get("min_font_exempt_idx"), report,
                               slide_index, shape_index, profile)
                check_palette(shape, profile.get("palette_rgb"),
                              profile.get("palette_tolerance", 6), report,
                              slide_index, shape_index, profile)
                check_safe_area(shape, rules.get("safe_area_in"),
                                profile.get("chrome_idx"), report,
                                slide_index, shape_index, profile)
                check_image_alpha(shape, profile.get("require_alpha_on_fills"),
                                  report, slide_index, shape_index, family,
                                  profile)
            except Exception as e:
                logger.debug("brand_shape_check_skipped slide=%d shape=%d "
                             "error=%s", slide_index, shape_index, e)
        if not layout_exempt:
            try:
                check_required_chrome(slide, rules.get("require"),
                                      profile.get("chrome_shapes"), report,
                                      slide_index, profile)
            except Exception as e:
                logger.debug("brand_chrome_check_skipped slide=%d error=%s",
                             slide_index, e)

    check_family_mix(families, profile.get("max_consecutive_same_family"),
                     report, profile)

    counts = report.counts()
    return {
        "ok": counts[ERROR] == 0,
        "brand": profile.get("name", "unnamed"),
        "slides": len(pres.slides),
        "backgrounds": families,
        "counts": counts,
        "problems": report.problems,
    }
