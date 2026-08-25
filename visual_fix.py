"""
Internal repair engine for the visual-QA loop (see visual_qa.py).

When inspection finds issues, the server — not the calling agent — repairs
the deck: the configured vision LLM is shown the issues, a structural
description of the affected slides, and their rendered images, and asked for
a plan made only of whitelisted, validated operations, which are applied
with python-pptx. The gate then re-inspects; see _visual_qa_gate in
tools/presentation_tools.py for the loop.

Whitelisted operations (anything else in a plan is skipped and reported):
- move_shape       {slide, shape_index, left_in, top_in}
- resize_shape     {slide, shape_index, width_in?, height_in?}
- set_font_size    {slide, shape_index, size_pt}  text box, table or chart
- fit_text         {slide, shape_index, min_pt?, max_pt?}  size text to its box
- set_autofit      {slide, shape_index, mode: shrink_text|grow_shape|none}
- set_text         {slide, shape_index, text}
- set_word_wrap    {slide, shape_index, wrap: true|false}
- delete_shape     {slide, shape_index}
- set_column_width {slide, shape_index, column, width_in}
- set_row_height   {slide, shape_index, row, height_in}
- set_cell_text    {slide, shape_index, row, column, text}
- set_chart_legend {slide, shape_index, show, position?}
- set_chart_data_labels {slide, shape_index, show}
- set_axis_title   {slide, shape_index, axis: category|value, text}

fit_text sizes text to the space it actually has — growing text that leaves
its box mostly empty as readily as shrinking text that overflows — with
growth anchored to the deck's existing typography so a two-word box does not
balloon to 96pt.

Text is not only in text boxes: the table and chart operations exist so crowded
table columns and colliding chart labels — the most common overlap the
reviewer reports — are actually fixable instead of being reported forever.

Slide numbers are 1-based (matching the inspector's issue reports); shape
indexes are the 0-based positions reported by describe_slides.
"""
import json
import logging
import math

from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Inches, Pt

import utils as ppt_utils
from logging_utils import get_logger, flatten

logger = get_logger("visual_fix")

MAX_TEXT_CHARS = 4000
FONT_PT_RANGE = (6, 96)
POSITION_IN_RANGE = (-5.0, 60.0)
SIZE_IN_RANGE = (0.05, 60.0)
AUTOFIT_MODES = {
    "shrink_text": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    "grow_shape": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
    "none": MSO_AUTO_SIZE.NONE,
}
# Text-fitting estimate. A rendered line of a proportional face averages
# roughly half the point size in width per character, and PowerPoint's
# single line spacing is about 1.2x the point size. Both are approximations:
# fit_text lands close, and the QA loop's re-render is what actually
# confirms it.
CHAR_WIDTH_RATIO = 0.5
LINE_HEIGHT_RATIO = 1.2
# Leave the text a little clear of its own border rather than filling the
# box to the pixel — text touching the frame reads as overcrowded.
FIT_SLACK = 0.92
LEGEND_POSITIONS = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "top": XL_LEGEND_POSITION.TOP,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
    "corner": XL_LEGEND_POSITION.CORNER,
}


def _emu_to_in(v):
    return round(Emu(v).inches, 3) if v is not None else None


def _usable_box_in(shape):
    """Inner width/height of a shape's text area, in inches (box minus the
    text frame's own margins)."""
    tf = shape.text_frame
    width = (Emu(shape.width).inches
             - Emu(tf.margin_left or 0).inches
             - Emu(tf.margin_right or 0).inches)
    height = (Emu(shape.height).inches
              - Emu(tf.margin_top or 0).inches
              - Emu(tf.margin_bottom or 0).inches)
    return max(width, 0.05), max(height, 0.05)


def _fits(paragraph_texts, size_pt, width_in, height_in, wrap):
    """Whether the text is estimated to fit the box at this point size."""
    char_width_in = size_pt / 72.0 * CHAR_WIDTH_RATIO
    line_height_in = size_pt / 72.0 * LINE_HEIGHT_RATIO
    if wrap:
        chars_per_line = max(1, int(width_in / char_width_in))
        lines = sum(max(1, math.ceil(len(t) / chars_per_line))
                    for t in paragraph_texts)
    else:
        # No wrapping: every paragraph is one line, but the longest one has
        # to fit the width on its own.
        lines = len(paragraph_texts)
        longest = max((len(t) for t in paragraph_texts), default=0)
        if longest * char_width_in > width_in:
            return False
    return lines * line_height_in <= height_in


def estimate_fit_font_size(shape, min_pt=None, max_pt=None):
    """Largest point size at which the shape's text is estimated to fit its
    box — used to grow text that is swimming in empty space as well as to
    shrink text that overflows.

    Returns None when the shape has no text to measure. The estimate is
    geometric (see CHAR_WIDTH_RATIO), not a real text layout: the inspect
    loop re-renders and re-reviews, which is what confirms the result.
    """
    if not shape.has_text_frame:
        return None
    texts = [p.text for p in shape.text_frame.paragraphs if p.text]
    if not texts:
        return None
    lo = int(min_pt if min_pt is not None else FONT_PT_RANGE[0])
    hi = int(max_pt if max_pt is not None else FONT_PT_RANGE[1])
    lo = max(lo, FONT_PT_RANGE[0])
    hi = min(hi, FONT_PT_RANGE[1])
    if hi < lo:
        return None
    width_in, height_in = _usable_box_in(shape)
    width_in *= FIT_SLACK
    height_in *= FIT_SLACK
    wrap = shape.text_frame.word_wrap is not False
    for size in range(hi, lo - 1, -1):
        if _fits(texts, size, width_in, height_in, wrap):
            return size
    return lo


def current_font_size_pt(shape):
    """Largest explicit run size in the shape, or None when the text
    inherits its size from the layout/master."""
    if not shape.has_text_frame:
        return None
    sizes = [run.font.size.pt
             for para in shape.text_frame.paragraphs
             for run in para.runs
             if run.font.size is not None]
    return max(sizes) if sizes else None


# Growing text is anchored to the deck's own typography: a box holding one
# short word would otherwise "fit" at 96pt and shout over the whole slide.
MAX_GROWTH_FACTOR = 1.5
DEFAULT_GROWTH_CEILING_PT = 44


def fit_growth_ceiling(shape):
    current = current_font_size_pt(shape)
    if current is None:
        return DEFAULT_GROWTH_CEILING_PT
    return max(current, min(FONT_PT_RANGE[1], current * MAX_GROWTH_FACTOR))


def _describe_table(table):
    """Column widths, row heights and cell text — what a planner needs to
    decide whether a column is too narrow or a row too short for its text."""
    info = {
        "rows": len(table.rows),
        "columns": len(table.columns),
        "column_widths_in": [_emu_to_in(c.width) for c in table.columns],
        "row_heights_in": [_emu_to_in(r.height) for r in table.rows],
    }
    cells, sizes = [], set()
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            text = cell.text_frame.text
            if text:
                cells.append({"row": r_idx, "column": c_idx,
                              "text": text[:80] + ("…" if len(text) > 80 else "")})
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        sizes.add(run.font.size.pt)
    info["cells"] = cells[:60]
    if sizes:
        info["font_sizes_pt"] = sorted(sizes)
    return info


def _describe_chart(chart):
    """Chart structure relevant to label crowding: how many categories and
    series compete for space, and which label layers are switched on."""
    info = {}
    try:
        info["chart_type"] = str(chart.chart_type)
    except Exception:  # some chart types raise on unknown enum values
        logger.debug("describe_chart_type_unavailable")
    try:
        info["series_count"] = len(chart.series)
        info["categories"] = [str(c)[:40] for c in chart.plots[0].categories][:20]
    except Exception:
        logger.debug("describe_chart_data_unavailable")
    try:
        info["has_legend"] = chart.has_legend
        if chart.has_legend:
            info["legend_position"] = str(chart.legend.position)
    except Exception:
        logger.debug("describe_chart_legend_unavailable")
    try:
        info["has_data_labels"] = chart.plots[0].has_data_labels
    except Exception:
        logger.debug("describe_chart_labels_unavailable")
    try:
        if chart.font.size is not None:
            info["font_size_pt"] = chart.font.size.pt
    except Exception:
        logger.debug("describe_chart_font_unavailable")
    for name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, name)
            # Reading axis_title materializes it, so ask has_title first.
            if axis.has_title:
                info[f"{name}_title"] = axis.axis_title.text_frame.text[:60]
        except Exception:  # pie/doughnut have no such axis
            logger.debug("describe_chart_axis_unavailable axis=%s", name)
    return info


def describe_slides(pres, slide_numbers=None):
    """Compact structural description of slides for the repair planner.
    slide_numbers: 1-based; None describes every slide.

    Text lives in more than text frames: tables and charts are described too,
    so the planner can act on a crowded axis or an over-narrow column instead
    of only seeing an opaque rectangle."""
    described = []
    for idx, slide in enumerate(pres.slides, start=1):
        if slide_numbers and idx not in slide_numbers:
            continue
        shapes = []
        for s_idx, shape in enumerate(slide.shapes):
            info = {
                "shape_index": s_idx,
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_in": _emu_to_in(shape.left),
                "top_in": _emu_to_in(shape.top),
                "width_in": _emu_to_in(shape.width),
                "height_in": _emu_to_in(shape.height),
                "is_placeholder": shape.is_placeholder,
            }
            if shape.has_text_frame:
                text = shape.text_frame.text
                info["text"] = text[:300] + ("…" if len(text) > 300 else "")
                sizes = sorted({
                    run.font.size.pt
                    for para in shape.text_frame.paragraphs
                    for run in para.runs
                    if run.font.size is not None
                })
                if sizes:
                    info["font_sizes_pt"] = sizes
                if shape.text_frame.word_wrap is not None:
                    info["word_wrap"] = shape.text_frame.word_wrap
            if getattr(shape, "has_table", False):
                info["table"] = _describe_table(shape.table)
            if getattr(shape, "has_chart", False):
                info["chart"] = _describe_chart(shape.chart)
            if shape.shape_type is not None and "GROUP" in str(shape.shape_type):
                # Group members are not individually addressable by the repair
                # ops; say so rather than describe shapes that cannot be fixed.
                info["group_members"] = len(shape.shapes)
            shapes.append(info)
        described.append({"slide": idx, "shapes": shapes})
    return described


REPAIR_PROMPT = """You are a presentation repair engine. A visual QA reviewer \
found the following issues in a PowerPoint deck:

{issues}

Below is the structural description of the affected slides (positions/sizes \
in inches, 1-based slide numbers, 0-based shape_index), followed by their \
rendered images.

{structure}

Slide dimensions: {width_in} x {height_in} inches.

Produce repair operations that fix the issues. Respond with ONLY a JSON \
object, no markdown fence:
{{"operations": [
  {{"op": "move_shape", "slide": N, "shape_index": N, "left_in": X, "top_in": Y}},
  {{"op": "resize_shape", "slide": N, "shape_index": N, "width_in": X, "height_in": Y}},
  {{"op": "set_font_size", "slide": N, "shape_index": N, "size_pt": X}},
  {{"op": "fit_text", "slide": N, "shape_index": N, "min_pt": X, "max_pt": X}},
  {{"op": "set_autofit", "slide": N, "shape_index": N, "mode": "shrink_text"}},
  {{"op": "set_text", "slide": N, "shape_index": N, "text": "..."}},
  {{"op": "set_word_wrap", "slide": N, "shape_index": N, "wrap": true}},
  {{"op": "delete_shape", "slide": N, "shape_index": N}},
  {{"op": "set_column_width", "slide": N, "shape_index": N, "column": N, "width_in": X}},
  {{"op": "set_row_height", "slide": N, "shape_index": N, "row": N, "height_in": X}},
  {{"op": "set_cell_text", "slide": N, "shape_index": N, "row": N, "column": N, "text": "..."}},
  {{"op": "set_chart_legend", "slide": N, "shape_index": N, "show": true, "position": "bottom"}},
  {{"op": "set_chart_data_labels", "slide": N, "shape_index": N, "show": false}},
  {{"op": "set_axis_title", "slide": N, "shape_index": N, "axis": "category", "text": "..."}}
]}}
Only these operation types are allowed. Keep every shape fully inside the \
slide bounds. Prefer minimal changes: resize, move or resize text before \
deleting anything. Do not touch shapes that have no reported issue.

Size text to the space it has. Text should fill its box comfortably — neither \
overflowing it nor floating in a mostly empty one — while the slide as a whole \
keeps its breathing room. Use fit_text (optionally bounded with min_pt/max_pt) \
to let the server compute the size that fits the box; it grows text that is too \
small for its container as well as shrinking text that overflows. Reach for \
set_font_size instead when you want one specific size, e.g. to match a sibling \
element. Growing text is not an improvement when it crowds neighbouring \
elements or leaves no margin: prefer consistency with comparable elements on \
the same slide, and never enlarge text just because there is room.

For text that is not in a text box, use the operation that matches the \
container. Table columns too narrow for their text: widen the column (and \
narrow another so the table still fits) or shrink the table font. Table rows \
clipping their text: raise the row height. Chart axis labels colliding or \
truncated: shrink the chart font, or widen/heighten the chart. Data labels \
overlapping their bars or each other: turn them off, or shrink the chart \
font. A legend covering the plot: move it to "bottom"/"right", or hide it \
when the categories are already labelled. set_font_size works on a text box, \
a whole table or a whole chart — shape_index alone picks the container. \
Members of a group cannot be addressed individually: move, resize or shrink \
the group as a whole.

Axis titles are named by role, not by screen position: "category" is the axis \
the categories sit on and "value" is the axis the numbers sit on. On a bar \
chart the category axis runs vertically and the value axis horizontally; on a \
column or line chart it is the other way round. Read the current titles in the \
structure above (category_axis_title / value_axis_title) and the image before \
deciding a title belongs on the other axis."""


def plan_repairs(llm, issues, pres, deck_images, image_slides=None):
    """Ask the vision LLM for a repair plan. Returns a list of operation
    dicts (possibly empty).

    image_slides: absolute 1-based slide number of each image in
    deck_images. Defaults to 1..N (a full-deck render); pass it whenever
    only a subset of slides was rendered, so the right images are attached.
    """
    slide_numbers = sorted({
        i.get("slide") for i in issues if isinstance(i.get("slide"), int)
    }) or None
    structure = describe_slides(pres, slide_numbers)
    prompt = REPAIR_PROMPT.format(
        issues=json.dumps(issues, indent=1),
        structure=json.dumps(structure, indent=1),
        width_in=_emu_to_in(pres.slide_width),
        height_in=_emu_to_in(pres.slide_height),
    )
    if image_slides is None:
        image_slides = list(range(1, len(deck_images) + 1))
    images = deck_images
    if slide_numbers:
        by_slide = dict(zip(image_slides, deck_images))
        images = [by_slide[n] for n in slide_numbers if n in by_slide]
    logger.debug("repair_plan_request issues=%d slides=%s images=%d",
                 len(issues), slide_numbers or "all", len(images))
    data = llm.ask_json(images, prompt)
    ops = data.get("operations", [])
    if not isinstance(ops, list):
        logger.warning("repair_plan_malformed keys=%s",
                       ",".join(sorted(data)) or "-")
        return []
    logger.info("repair_plan_received operations=%d issues=%d",
                len(ops), len(issues))
    return ops


def _in_range(v, lo_hi):
    return isinstance(v, (int, float)) and lo_hi[0] <= v <= lo_hi[1]


def skip_reason_summary(skipped):
    """Count the distinct reasons a plan's operations were rejected, e.g.
    {"bad shape_index": 2}. Reported to the caller so a round that applied
    nothing explains itself without a trip to the server logs; the parameter
    detail after ':' is dropped to keep it a small, stable summary."""
    counts = {}
    for entry in skipped:
        reason = str(entry.get("reason", "unknown")).split(":")[0]
        counts[reason] = counts.get(reason, 0) + 1
    return counts


def apply_repairs(pres, operations, allowed_slides=None):
    """Validate and apply whitelisted operations. Returns
    {"applied": [...], "skipped": [{"op":..., "reason":...}]}.

    allowed_slides: when set, operations targeting any other slide are
    skipped — a slide-scoped repair call must not touch slides the caller
    did not put in scope."""
    applied, skipped = [], []
    slides = list(pres.slides)

    for op in operations if isinstance(operations, list) else []:
        if not isinstance(op, dict):
            skipped.append({"op": op, "reason": "not an object"})
            continue
        kind = op.get("op")
        slide_no, shape_idx = op.get("slide"), op.get("shape_index")
        if not (isinstance(slide_no, int) and 1 <= slide_no <= len(slides)):
            skipped.append({"op": op, "reason": "bad slide number"})
            continue
        if allowed_slides and slide_no not in allowed_slides:
            skipped.append({"op": op, "reason": "slide out of scope"})
            continue
        shapes = list(slides[slide_no - 1].shapes)
        if not (isinstance(shape_idx, int) and 0 <= shape_idx < len(shapes)):
            skipped.append({"op": op, "reason": "bad shape_index"})
            continue
        shape = shapes[shape_idx]

        try:
            if kind == "move_shape":
                if not (_in_range(op.get("left_in"), POSITION_IN_RANGE)
                        and _in_range(op.get("top_in"), POSITION_IN_RANGE)):
                    skipped.append({"op": op, "reason": "position out of range"})
                    continue
                ppt_utils.pin_inherited_geometry(shape)
                shape.left = Inches(op["left_in"])
                shape.top = Inches(op["top_in"])
            elif kind == "resize_shape":
                w, h = op.get("width_in"), op.get("height_in")
                if w is None and h is None:
                    skipped.append({"op": op, "reason": "no dimensions"})
                    continue
                ppt_utils.pin_inherited_geometry(shape)
                if w is not None:
                    if not _in_range(w, SIZE_IN_RANGE):
                        skipped.append({"op": op, "reason": "width out of range"})
                        continue
                    shape.width = Inches(w)
                if h is not None:
                    if not _in_range(h, SIZE_IN_RANGE):
                        skipped.append({"op": op, "reason": "height out of range"})
                        continue
                    shape.height = Inches(h)
            elif kind == "set_font_size":
                if not _in_range(op.get("size_pt"), FONT_PT_RANGE):
                    skipped.append({"op": op, "reason": "font size out of range"})
                    continue
                size = Pt(op["size_pt"])
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = size
                elif getattr(shape, "has_table", False):
                    # Shrinking every cell is the usual answer to table text
                    # that wraps into an unreadable stack.
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                para.font.size = size
                                for run in para.runs:
                                    run.font.size = size
                elif getattr(shape, "has_chart", False):
                    # One chart-wide font: axis ticks, data labels and legend.
                    shape.chart.font.size = size
                else:
                    skipped.append({"op": op, "reason": "shape has no text"})
                    continue
            elif kind == "set_text":
                text = op.get("text")
                if not isinstance(text, str) or len(text) > MAX_TEXT_CHARS:
                    skipped.append({"op": op, "reason": "bad text"})
                    continue
                if not shape.has_text_frame:
                    skipped.append({"op": op, "reason": "shape has no text frame"})
                    continue
                shape.text_frame.text = text
            elif kind == "set_word_wrap":
                if not shape.has_text_frame:
                    skipped.append({"op": op, "reason": "shape has no text frame"})
                    continue
                shape.text_frame.word_wrap = bool(op.get("wrap", True))
            elif kind == "fit_text":
                if not shape.has_text_frame:
                    skipped.append({"op": op, "reason": "shape has no text frame"})
                    continue
                min_pt, max_pt = op.get("min_pt"), op.get("max_pt")
                for value, label in ((min_pt, "min_pt"), (max_pt, "max_pt")):
                    if value is not None and not _in_range(value, FONT_PT_RANGE):
                        skipped.append({"op": op,
                                        "reason": f"{label} out of range"})
                        break
                else:
                    if max_pt is None:
                        max_pt = fit_growth_ceiling(shape)
                    size = estimate_fit_font_size(shape, min_pt, max_pt)
                    if size is None:
                        skipped.append({"op": op, "reason": "shape has no text"})
                        continue
                    for para in shape.text_frame.paragraphs:
                        para.font.size = Pt(size)
                        for run in para.runs:
                            run.font.size = Pt(size)
                    op = dict(op, resolved_size_pt=size)
                    logger.debug("fit_text slide=%s shape_index=%s size_pt=%s",
                                 slide_no, shape_idx, size)
                    applied.append(op)
                continue
            elif kind == "set_autofit":
                if not shape.has_text_frame:
                    skipped.append({"op": op, "reason": "shape has no text frame"})
                    continue
                mode = op.get("mode")
                if mode not in AUTOFIT_MODES:
                    skipped.append({"op": op, "reason": "bad autofit mode"})
                    continue
                shape.text_frame.auto_size = AUTOFIT_MODES[mode]
            elif kind in ("set_column_width", "set_row_height",
                          "set_cell_text"):
                if not getattr(shape, "has_table", False):
                    skipped.append({"op": op, "reason": "shape is not a table"})
                    continue
                table = shape.table
                if kind == "set_column_width":
                    col = op.get("column")
                    if not (isinstance(col, int) and 0 <= col < len(table.columns)):
                        skipped.append({"op": op, "reason": "bad column index"})
                        continue
                    if not _in_range(op.get("width_in"), SIZE_IN_RANGE):
                        skipped.append({"op": op, "reason": "width out of range"})
                        continue
                    table.columns[col].width = Inches(op["width_in"])
                elif kind == "set_row_height":
                    row_idx = op.get("row")
                    if not (isinstance(row_idx, int) and 0 <= row_idx < len(table.rows)):
                        skipped.append({"op": op, "reason": "bad row index"})
                        continue
                    if not _in_range(op.get("height_in"), SIZE_IN_RANGE):
                        skipped.append({"op": op, "reason": "height out of range"})
                        continue
                    table.rows[row_idx].height = Inches(op["height_in"])
                else:  # set_cell_text
                    row_idx, col = op.get("row"), op.get("column")
                    text = op.get("text")
                    if not (isinstance(row_idx, int) and 0 <= row_idx < len(table.rows)
                            and isinstance(col, int) and 0 <= col < len(table.columns)):
                        skipped.append({"op": op, "reason": "bad cell reference"})
                        continue
                    if not isinstance(text, str) or len(text) > MAX_TEXT_CHARS:
                        skipped.append({"op": op, "reason": "bad text"})
                        continue
                    table.cell(row_idx, col).text_frame.text = text
            elif kind in ("set_chart_legend", "set_chart_data_labels"):
                if not getattr(shape, "has_chart", False):
                    skipped.append({"op": op, "reason": "shape is not a chart"})
                    continue
                show = op.get("show")
                if not isinstance(show, bool):
                    skipped.append({"op": op, "reason": "show must be true/false"})
                    continue
                chart = shape.chart
                if kind == "set_chart_legend":
                    position = op.get("position")
                    if position is not None and position not in LEGEND_POSITIONS:
                        skipped.append({"op": op, "reason": "bad legend position"})
                        continue
                    chart.has_legend = show
                    if show and position is not None:
                        chart.legend.position = LEGEND_POSITIONS[position]
                        chart.legend.include_in_layout = False
                else:  # set_chart_data_labels — the usual cure for labels
                       # stacked on top of their own bars
                    for plot in chart.plots:
                        plot.has_data_labels = show
            elif kind == "set_axis_title":
                if not getattr(shape, "has_chart", False):
                    skipped.append({"op": op, "reason": "shape is not a chart"})
                    continue
                which = op.get("axis")
                text = op.get("text")
                if which not in ("category", "value"):
                    skipped.append({"op": op,
                                    "reason": "axis must be category or value"})
                    continue
                if not isinstance(text, str) or len(text) > MAX_TEXT_CHARS:
                    skipped.append({"op": op, "reason": "bad text"})
                    continue
                axis = getattr(shape.chart, f"{which}_axis")
                axis.has_title = True
                axis.axis_title.text_frame.text = text
            elif kind == "delete_shape":
                shape._element.getparent().remove(shape._element)
            else:
                skipped.append({"op": op, "reason": f"unknown op '{kind}'"})
                continue
        except Exception as e:
            logger.debug("repair_op_failed op=%s slide=%s shape_index=%s error=%s",
                         kind, slide_no, shape_idx, flatten(str(e)))
            skipped.append({"op": op, "reason": f"apply failed: {e}"})
            continue
        logger.debug("repair_op_applied op=%s slide=%s shape_index=%s",
                     kind, slide_no, shape_idx)
        applied.append(op)

    level = logging.WARNING if skipped else logging.INFO
    logger.log(level, "repairs_applied applied=%d skipped=%d%s",
               len(applied), len(skipped),
               " reasons=" + ",".join(sorted({s["reason"].split(":")[0]
                                              for s in skipped})) if skipped else "")
    return {"applied": applied, "skipped": skipped}
