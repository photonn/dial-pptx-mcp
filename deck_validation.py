"""
Structural validation of a presentation, before it is handed to the user.

Visual QA looks at pixels and cannot see a broken package: a deck whose
relationships don't resolve still renders in LibreOffice, still opens in
python-pptx, and still fails in PowerPoint. This module covers the other axis —
whether the file is well-formed and whether its slides are structurally sane —
and it is deliberately cheap: no rendering, no model call, so it can run on
every export.

The checks are grouped by what they can tell you:

- ``error``   the deck is malformed. PowerPoint may refuse to open it, or will
              open it with content missing. Always worth fixing.
- ``warning`` the deck opens, but something on it is almost certainly a defect
              the user would notice — a shape off the canvas, an empty chart,
              template placeholder text nobody replaced.
- ``info``    worth knowing before you trust a visual QA verdict, e.g. fonts
              LibreOffice will substitute with different metrics.

Every problem carries a "fix" naming the tool that resolves it, because the
consumer is an agent deciding what to call next, not a person reading a report.

An *empty* placeholder is deliberately not reported: PowerPoint draws its prompt
text only in edit view, so it is invisible in a slideshow and in the exported
PDF, and a corporate template has dozens of them — reporting each one buries the
findings that matter. Placeholder text actually present in the deck is a
different matter and is reported.
"""
import io
import re
import zipfile

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from logging_utils import get_logger

logger = get_logger("deck_validation")

ERROR = "error"
WARNING = "warning"
INFO = "info"

# Attributes that carry a relationship id from slide XML into the part's rels.
_REL_ID_ATTRS = (qn("r:id"), qn("r:embed"), qn("r:link"),
                 qn("r:pict"), qn("r:dm"), qn("r:lo"), qn("r:qs"), qn("r:cs"))

# Text a template ships in its own placeholders, or that a half-finished draft
# leaves behind. Matched case-insensitively against every string on a slide.
_PLACEHOLDER_PATTERNS = (
    (r"\blorem ipsum\b", "lorem ipsum filler"),
    (r"\bclick to (add|edit)\b", "an unfilled PowerPoint placeholder prompt"),
    (r"\bx{3,}\b", "XXX filler"),
    (r"\bTODO\b", "a TODO marker"),
    (r"\[\s*insert[^\]]*\]", "an [insert ...] marker"),
    (r"\byour (title|text|company|logo) here\b", "a 'your ... here' marker"),
    (r"\bsample text\b", "sample text"),
    (r"\bplaceholder\b", "the word 'placeholder'"),
)
_PLACEHOLDER_RE = [(re.compile(p, re.IGNORECASE), label)
                   for p, label in _PLACEHOLDER_PATTERNS]

# How far outside the slide a shape may sit before it is reported. Templates
# routinely bleed decoration a hair past the edge on purpose.
_EDGE_TOLERANCE_EMU = 45720  # 0.05in
EMU_PER_INCH = 914400

# Slide ids live in this range; PowerPoint rejects a deck that breaks it.
_SLIDE_ID_MIN = 256
_SLIDE_ID_MAX = 2147483647

# A picture stretched away from its stored aspect ratio by more than this is
# visibly distorted. Nothing in the repair whitelist can un-distort a picture,
# so it is worth catching here rather than in the visual loop.
_ASPECT_TOLERANCE = 0.08

# An icon whose transparency was flattened shows as a rectangle once it is
# placed on a coloured card. Corner pixels at or above this alpha are opaque,
# and four corners within this per-channel distance of each other are a flat
# pane of one colour — the background the artwork was flattened against —
# rather than four unrelated corners of a photograph.
_OPAQUE_ALPHA = 250
_UNIFORM_CORNERS = 12
# Corners are sampled a little inside the edge: a flattened icon often keeps a
# single row of stray antialiasing at the very border.
_CORNER_INSET = 0.02


def _inches(emu):
    return round(emu / EMU_PER_INCH, 2)


class _Report:
    def __init__(self):
        self.problems = []

    def add(self, severity, code, message, fix, slide=None, shape=None):
        problem = {"severity": severity, "code": code, "message": message,
                   "fix": fix}
        if slide is not None:
            problem["slide_index"] = slide
        if shape is not None:
            problem["shape_index"] = shape
        self.problems.append(problem)

    def counts(self):
        return {level: sum(1 for p in self.problems if p["severity"] == level)
                for level in (ERROR, WARNING, INFO)}


# ---- Package-level checks ----

def _check_package(blob, report):
    """Round-trip the saved bytes: a package that cannot be reopened is the
    one failure that makes every other check meaningless."""
    try:
        archive = zipfile.ZipFile(io.BytesIO(blob))
    except zipfile.BadZipFile as e:
        report.add(ERROR, "package_unreadable",
                   f"The saved file is not a readable OPC package: {e}",
                   "Rebuild the presentation; this deck cannot be delivered.")
        return None

    broken = archive.testzip()
    if broken:
        report.add(ERROR, "package_corrupt",
                   f"Package entry '{broken}' has a bad CRC.",
                   "Rebuild the presentation; this deck cannot be delivered.")

    names = set(archive.namelist())
    if "[Content_Types].xml" not in names:
        report.add(ERROR, "content_types_missing",
                   "The package has no [Content_Types].xml part.",
                   "Rebuild the presentation.")
        return archive

    declared = archive.read("[Content_Types].xml").decode("utf-8", "replace")
    extensions = set(re.findall(r'Extension="([^"]+)"', declared))
    overrides = set(re.findall(r'PartName="([^"]+)"', declared))
    for name in sorted(names):
        if name == "[Content_Types].xml" or name.endswith("/"):
            continue
        if f"/{name}" in overrides:
            continue
        extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        if extension not in extensions:
            report.add(ERROR, "content_type_undeclared",
                       f"Part '{name}' has no declared content type.",
                       "Rebuild the presentation; a part was added to the "
                       "package without registering its content type.")

    try:
        Presentation(io.BytesIO(blob))
    except Exception as e:
        report.add(ERROR, "package_unopenable",
                   f"The saved deck cannot be reopened: {e}",
                   "Rebuild the presentation; this deck cannot be delivered.")
    return archive


def _check_relationships(pres, report):
    """Every rId referenced by a part's XML must exist in that part's rels.

    A dangling reference is what a hand-assembled or badly copied slide
    produces, and it is exactly the class of fault that renders fine in
    LibreOffice and fails in PowerPoint.
    """
    for slide_index, slide in enumerate(pres.slides):
        part = slide.part
        available = set(part.rels.keys())
        for element in slide._element.iter():
            for attr in _REL_ID_ATTRS:
                rId = element.get(attr)
                if rId and rId not in available:
                    report.add(ERROR, "dangling_relationship",
                               f"Slide {slide_index} refers to relationship "
                               f"'{rId}', which the slide does not have.",
                               "The slide references a part that is not "
                               "attached to it. Rebuild this slide, or "
                               "duplicate a working one with duplicate_slide.",
                               slide=slide_index)

        layouts = [r for r in part.rels.values()
                   if r.reltype == RT.SLIDE_LAYOUT]
        if not layouts:
            report.add(ERROR, "slide_without_layout",
                       f"Slide {slide_index} is not related to any slide "
                       f"layout.",
                       "Rebuild this slide with add_slide or duplicate_slide.",
                       slide=slide_index)


def _check_slide_ids(pres, report):
    seen = {}
    for position, sld_id in enumerate(pres.slides._sldIdLst):
        value = int(sld_id.get("id"))
        if value < _SLIDE_ID_MIN or value > _SLIDE_ID_MAX:
            report.add(ERROR, "slide_id_out_of_range",
                       f"Slide at position {position} has id {value}, outside "
                       f"the valid range {_SLIDE_ID_MIN}-{_SLIDE_ID_MAX}.",
                       "Rebuild the deck; PowerPoint rejects this id.",
                       slide=position)
        if value in seen:
            report.add(ERROR, "duplicate_slide_id",
                       f"Slides at positions {seen[value]} and {position} "
                       f"share slide id {value}.",
                       "Rebuild the deck; two slides cannot share an id.",
                       slide=position)
        seen[value] = position


def _check_shared_notes(pres, report):
    """Two slides pointing at one notesSlide part is an invalid package and a
    trap for anything that copies slides around."""
    owners = {}
    for slide_index, slide in enumerate(pres.slides):
        for rel in slide.part.rels.values():
            if rel.reltype != RT.NOTES_SLIDE or rel.is_external:
                continue
            partname = str(rel.target_part.partname)
            if partname in owners:
                report.add(ERROR, "shared_notes_slide",
                           f"Slides {owners[partname]} and {slide_index} share "
                           f"one notes slide ({partname}).",
                           "Rewrite the notes on one of them with "
                           "manage_speaker_notes(operation='set').",
                           slide=slide_index)
            owners[partname] = slide_index


def _check_orphan_parts(pres, report):
    """Media and chart parts nothing references bloat the file and, for charts,
    usually mean a slide lost its graphic frame."""
    package = pres.part.package
    referenced = set()
    for part in package.iter_parts():
        for rel in part.rels.values():
            if not rel.is_external:
                referenced.add(str(rel.target_part.partname))

    for part in package.iter_parts():
        partname = str(part.partname)
        if partname == "/ppt/presentation.xml" or partname in referenced:
            continue
        if partname.startswith(("/ppt/media/", "/ppt/charts/",
                                "/ppt/embeddings/", "/ppt/diagrams/")):
            report.add(WARNING, "orphan_part",
                       f"Part '{partname}' is in the package but nothing "
                       f"references it.",
                       "Harmless but wasteful; it inflates the exported file.")


# ---- Slide-level checks ----

def _iter_text(shape):
    """Yield every string a shape shows, text frames and table cells alike."""
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text
        if text:
            yield text
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame.text:
                    yield cell.text_frame.text


def _check_geometry(slide_index, shape_index, shape, width, height, report):
    if shape.width is None or shape.height is None:
        return
    if shape.width <= 0 or shape.height <= 0:
        report.add(ERROR, "zero_sized_shape",
                   f"Slide {slide_index}, shape {shape_index} "
                   f"('{shape.name}') has a zero or negative size.",
                   "Give it a size with visual_repair_slides, or remove it.",
                   slide=slide_index, shape=shape_index)
        return
    if shape.left is None or shape.top is None:
        return

    right, bottom = shape.left + shape.width, shape.top + shape.height
    off_left = shape.left < -_EDGE_TOLERANCE_EMU
    off_top = shape.top < -_EDGE_TOLERANCE_EMU
    off_right = right > width + _EDGE_TOLERANCE_EMU
    off_bottom = bottom > height + _EDGE_TOLERANCE_EMU
    if not (off_left or off_top or off_right or off_bottom):
        return

    fully_off = (right <= 0 or bottom <= 0 or shape.left >= width
                 or shape.top >= height)
    edges = ", ".join(name for name, hit in
                      (("left", off_left), ("top", off_top),
                       ("right", off_right), ("bottom", off_bottom)) if hit)
    report.add(
        ERROR if fully_off else WARNING,
        "shape_off_slide" if fully_off else "shape_past_edge",
        f"Slide {slide_index}, shape {shape_index} ('{shape.name}') "
        f"{'is entirely off the slide' if fully_off else 'extends past the ' + edges + ' edge'} "
        f"— it sits at {_inches(shape.left)}, {_inches(shape.top)} in and is "
        f"{_inches(shape.width)}x{_inches(shape.height)} in on a "
        f"{_inches(width)}x{_inches(height)} in slide.",
        "Move or resize it with visual_repair_slides, or delete it.",
        slide=slide_index, shape=shape_index)


def _check_transform(slide_index, shape_index, shape, report):
    """Report an `a:xfrm` that names a size but no offset, or the reverse.

    Neither visual QA nor _check_geometry can see this one. python-pptx reports
    the value the placeholder inherits for whichever half is missing, and
    LibreOffice renders from the same fallback, so the shape looks correctly
    placed everywhere we can look at it. PowerPoint reads the half-written
    transform literally: the missing offset is the slide's top-left corner and
    a missing extent is nothing at all.
    """
    xfrm = None
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = shape._element.find(qn("p:xfrm"))
    if xfrm is None:
        return
    missing = tuple(part for part in ("off", "ext")
                    if xfrm.find(qn("a:" + part)) is None)
    if not missing:
        return
    # What PowerPoint substitutes differs by half, and the fix differs with it.
    lost, consequence, fix = {
        ("off",): ("position",
                   "draw it at the top-left corner of the slide",
                   "Give it an explicit position with visual_repair_slides "
                   "(move_shape writes position and size together)."),
        ("ext",): ("size",
                   "draw it with no size, so it disappears",
                   "Give it an explicit size with visual_repair_slides "
                   "(resize_shape writes position and size together)."),
        ("off", "ext"): ("position or size",
                         "draw it at the top-left corner of the slide with "
                         "no size",
                         "Give it an explicit position and size with "
                         "visual_repair_slides (move_shape writes both)."),
    }[missing]
    report.add(
        WARNING, "partial_transform",
        f"Slide {slide_index}, shape {shape_index} ('{shape.name}') has a "
        f"transform with no {lost}. It renders in place in LibreOffice, which "
        f"falls back to the layout, but PowerPoint will {consequence}.",
        fix,
        slide=slide_index, shape=shape_index)


def _check_picture(slide_index, shape_index, shape, report):
    """Report a picture stretched off its native aspect ratio.

    Nothing in the repair whitelist can restore a distorted picture, so it has
    to be caught before the visual loop spends iterations on it.
    """
    try:
        native_w, native_h = shape.image.size
    except Exception:
        return
    if not native_w or not native_h or not shape.width or not shape.height:
        return

    # Cropping legitimately changes the visible aspect ratio.
    cropped = any(abs(getattr(shape, f"crop_{side}", 0) or 0) > 0.001
                  for side in ("left", "right", "top", "bottom"))
    if cropped:
        return

    native_aspect = native_w / native_h
    placed_aspect = shape.width / shape.height
    drift = abs(placed_aspect - native_aspect) / native_aspect
    if drift > _ASPECT_TOLERANCE:
        report.add(WARNING, "distorted_picture",
                   f"Slide {slide_index}, shape {shape_index} "
                   f"('{shape.name}') is stretched {round(drift * 100)}% away "
                   f"from the image's own proportions.",
                   "Re-place it with add_image_from_dial_url using "
                   "fit='contain' or fit='cover'; visual repair cannot undo "
                   "distortion.",
                   slide=slide_index, shape=shape_index)


def _solid_fill_rgb(shape):
    """The shape's solid fill as an (r, g, b) tuple, or None.

    A theme-coloured fill cannot be resolved to RGB without walking the theme,
    so it is reported as its role instead: the two light roles are treated as
    "white enough to hide nothing", everything else as a colour an opaque
    picture would sit on top of.
    """
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_THEME_COLOR

    try:
        fill = shape.fill
        if fill.type != MSO_FILL.SOLID:
            return None
        color = fill.fore_color
        if color.type == MSO_COLOR_TYPE.RGB:
            return tuple(color.rgb)
        if color.theme_color in (MSO_THEME_COLOR.BACKGROUND_1,
                                 MSO_THEME_COLOR.LIGHT_1):
            return (255, 255, 255)
        return None if color.theme_color is None else (0, 0, 0)
    except Exception:
        return None


def _overlaps(a, b):
    for attr in ("left", "top", "width", "height"):
        if getattr(a, attr, None) is None or getattr(b, attr, None) is None:
            return False
    return (a.left < b.left + b.width and b.left < a.left + a.width
            and a.top < b.top + b.height and b.top < a.top + a.height)


def _corner_pixels(blob):
    """The four corner pixels of an image, as RGBA tuples, or None."""
    from PIL import Image

    with Image.open(io.BytesIO(blob)) as image:
        image = image.convert("RGBA")
        width, height = image.size
        if width < 4 or height < 4:
            return None
        x = max(1, int(width * _CORNER_INSET))
        y = max(1, int(height * _CORNER_INSET))
        return [image.getpixel(p) for p in
                ((x, y), (width - 1 - x, y),
                 (x, height - 1 - y), (width - 1 - x, height - 1 - y))]


def _check_picture_backgrounds(slide_index, slide, report):
    """Report a picture with a baked-in white background sitting on a fill.

    This is the icon failure mode: artwork exported with its transparency
    flattened against a background looks fine on a slide of that colour and
    shows a hard rectangle the moment it is placed on a card or panel of
    another. Neither pass catches it today — the render shows it, but a vision
    reviewer reads it as part of the design, and nothing else inspects alpha
    channels.

    Reported when the four corners are opaque *and* all one colour, which is
    what a flattened background looks like and what the corners of a photo do
    not. A picture whose corners already match the fill beneath it is left
    alone: there is nothing to see.

    Only pictures that actually overlap a solid, non-white fill *behind* them
    are decoded, so the usual deck pays nothing for the check.
    """
    shapes = list(slide.shapes)
    fills = [(index, _solid_fill_rgb(shape))
             for index, shape in enumerate(shapes)]
    fills = [(index, rgb) for index, rgb in fills
             if rgb is not None and rgb != (255, 255, 255)]
    if not fills:
        return

    for shape_index, shape in enumerate(shapes):
        if shape.shape_type != 13:  # PICTURE
            continue
        behind = [index for index, _ in fills
                  if index < shape_index and _overlaps(shape, shapes[index])]
        if not behind:
            continue
        try:
            corners = _corner_pixels(shape.image.blob)
        except Exception as e:
            logger.debug("picture_alpha_skipped slide=%d shape=%d error=%s",
                          slide_index, shape_index, e)
            continue
        if not corners:
            continue
        if not all(pixel[3] >= _OPAQUE_ALPHA for pixel in corners):
            continue
        first = corners[0][:3]
        if not all(abs(channel - first[i]) <= _UNIFORM_CORNERS
                   for pixel in corners for i, channel in
                   enumerate(pixel[:3])):
            continue  # four different corners: artwork, not a flat pane
        against = dict(fills)[behind[0]]
        if all(abs(a - b) <= _UNIFORM_CORNERS for a, b in zip(first, against)):
            continue  # the same colour as the fill: invisible either way
        report.add(WARNING, "opaque_picture_background",
                   f"Slide {slide_index}, shape {shape_index} "
                   f"('{shape.name}') has an opaque {list(first)} background "
                   f"and sits on the differently coloured shape {behind[0]}, "
                   f"so it shows as a rectangle rather than blending into it.",
                   "Use a version of the image with a transparent background, "
                   "or move it onto a white area — visual repair can move and "
                   "delete a picture but cannot remove its background.",
                   slide=slide_index, shape=shape_index)


def _check_chart(slide_index, shape_index, shape, report):
    chart = shape.chart
    try:
        plots = list(chart.plots)
    except Exception as e:
        report.add(ERROR, "chart_unreadable",
                   f"Slide {slide_index}, shape {shape_index}: the chart "
                   f"cannot be read ({e}).",
                   "Rebuild it with add_chart.",
                   slide=slide_index, shape=shape_index)
        return

    series = [s for plot in plots for s in plot.series]
    if not series:
        report.add(ERROR, "chart_without_series",
                   f"Slide {slide_index}, shape {shape_index}: the chart has "
                   f"no data series and will render as an empty frame.",
                   "Give it data with update_chart_data, or delete it.",
                   slide=slide_index, shape=shape_index)
        return

    for plot in plots:
        categories = len(list(plot.categories))
        if categories == 0:
            report.add(WARNING, "chart_without_categories",
                       f"Slide {slide_index}, shape {shape_index}: a plot has "
                       f"no categories.",
                       "Supply categories with update_chart_data.",
                       slide=slide_index, shape=shape_index)
            continue
        for one in plot.series:
            values = len(list(one.values))
            if values != categories:
                report.add(WARNING, "chart_series_length_mismatch",
                           f"Slide {slide_index}, shape {shape_index}: series "
                           f"'{one.name}' has {values} value(s) for "
                           f"{categories} categories.",
                           "Re-send matching categories and values with "
                           "update_chart_data.",
                           slide=slide_index, shape=shape_index)


def _check_table(slide_index, shape_index, shape, report):
    table = shape.table
    if len(table.rows) == 0 or len(table.columns) == 0:
        report.add(ERROR, "empty_table",
                   f"Slide {slide_index}, shape {shape_index}: the table has "
                   f"no rows or no columns.",
                   "Rebuild it with add_table, or delete it.",
                   slide=slide_index, shape=shape_index)


def _check_placeholder_text(slide_index, shape_index, shape, report):
    for text in _iter_text(shape):
        for pattern, label in _PLACEHOLDER_RE:
            match = pattern.search(text)
            if match:
                snippet = text.strip().replace("\n", " ")[:60]
                report.add(WARNING, "leftover_placeholder_text",
                           f"Slide {slide_index}, shape {shape_index} contains "
                           f"{label}: \"{snippet}\".",
                           "Replace it with real content (manage_text or "
                           "populate_placeholder), or delete the shape.",
                           slide=slide_index, shape=shape_index)
                return


def _check_slides(pres, report):
    width, height = pres.slide_width, pres.slide_height
    if not width or not height:
        report.add(ERROR, "slide_size_missing",
                   "The presentation declares no slide size.",
                   "Rebuild the presentation.")
        return

    for slide_index, slide in enumerate(pres.slides):
        for shape_index, shape in enumerate(slide.shapes):
            try:
                _check_geometry(slide_index, shape_index, shape, width, height,
                                report)
                _check_transform(slide_index, shape_index, shape, report)
                _check_placeholder_text(slide_index, shape_index, shape, report)
                if getattr(shape, "has_chart", False):
                    _check_chart(slide_index, shape_index, shape, report)
                if getattr(shape, "has_table", False):
                    _check_table(slide_index, shape_index, shape, report)
                if shape.shape_type == 13:  # PICTURE
                    _check_picture(slide_index, shape_index, shape, report)
            except Exception as e:
                logger.debug("shape_check_skipped slide=%d shape=%d error=%s",
                             slide_index, shape_index, e)
        try:
            _check_picture_backgrounds(slide_index, slide, report)
        except Exception as e:
            logger.debug("picture_background_check_skipped slide=%d error=%s",
                         slide_index, e)


def _check_fonts(pres, report):
    """Report fonts LibreOffice renders with different metrics.

    This is not a defect in the deck — it is a caveat on the visual QA verdict,
    which is rendered through LibreOffice. See fonts.py.
    """
    from fonts import unreliable_fonts_in

    risky = unreliable_fonts_in(pres)
    if not risky:
        return
    report.add(INFO, "qa_unreliable_fonts",
               "This deck uses font(s) LibreOffice substitutes with different "
               "character widths: " + ", ".join(sorted(risky)) + ".",
               "Text-fit findings from visual_inspect_slides are approximate "
               "for text in these fonts — leave the containers ~10% slack, or "
               "switch body text to a metric-safe font with manage_fonts.")


# ---- Entry point ----

def validate_presentation(pres, blob=None):
    """Validate an open presentation; return a report dict.

    `blob` is the deck's saved bytes when the caller already has them; it is
    saved here otherwise. The package checks need the serialized form, since
    that is what the user actually receives.
    """
    report = _Report()

    if blob is None:
        buffer = io.BytesIO()
        pres.save(buffer)
        blob = buffer.getvalue()

    _check_package(blob, report)
    _check_relationships(pres, report)
    _check_slide_ids(pres, report)
    _check_shared_notes(pres, report)
    _check_orphan_parts(pres, report)
    _check_slides(pres, report)
    _check_fonts(pres, report)

    counts = report.counts()
    return {
        "ok": counts[ERROR] == 0,
        "slides": len(pres.slides),
        "size_bytes": len(blob),
        "counts": counts,
        "problems": report.problems,
    }
