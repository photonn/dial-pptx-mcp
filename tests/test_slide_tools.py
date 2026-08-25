"""
Tests for the slide-structure layer: delete, move, duplicate, cross-deck copy
and speaker notes.

Duplication is the one that can silently produce a broken package, so the
assertions go past "a slide appeared": the copy is re-read from a saved deck,
and its chart is edited to prove the two slides do not share one chart part.
"""
import io
import unittest

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches

import utils as ppt_utils

DEMO = "mcp_all_tools_templates_effects_demo.pptx"


def roundtrip(pres):
    """Save and reopen, so assertions see what a client would actually get."""
    buf = io.BytesIO()
    pres.save(buf)
    return Presentation(io.BytesIO(buf.getvalue()))


def slide_titles(pres):
    return [s.shapes.title.text if s.shapes.title else "" for s in pres.slides]


def build_deck(count=4):
    pres = Presentation()
    layout = pres.slide_layouts[5]
    for i in range(count):
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i}"
    return pres


class TestDeleteAndMove(unittest.TestCase):
    def test_delete_removes_the_named_slide(self):
        pres = build_deck(4)
        ppt_utils.delete_slide(pres, 1)
        self.assertEqual(slide_titles(pres), ["Slide 0", "Slide 2", "Slide 3"])
        self.assertEqual(slide_titles(roundtrip(pres)),
                         ["Slide 0", "Slide 2", "Slide 3"])

    def test_delete_drops_the_presentation_relationship(self):
        pres = build_deck(3)
        before = len(pres.part.rels)
        ppt_utils.delete_slide(pres, 0)
        self.assertEqual(len(pres.part.rels), before - 1)

    def test_move_forward_and_back(self):
        pres = build_deck(4)
        ppt_utils.move_slide(pres, 0, 2)
        self.assertEqual(slide_titles(pres),
                         ["Slide 1", "Slide 2", "Slide 0", "Slide 3"])
        ppt_utils.move_slide(pres, 2, 0)
        self.assertEqual(slide_titles(pres),
                         ["Slide 0", "Slide 1", "Slide 2", "Slide 3"])

    def test_move_to_last_position(self):
        pres = build_deck(3)
        ppt_utils.move_slide(pres, 0, 2)
        self.assertEqual(slide_titles(pres), ["Slide 1", "Slide 2", "Slide 0"])

    def test_move_survives_a_roundtrip(self):
        pres = build_deck(3)
        ppt_utils.move_slide(pres, 2, 0)
        self.assertEqual(slide_titles(roundtrip(pres)),
                         ["Slide 2", "Slide 0", "Slide 1"])

    def test_check_index_reports_the_valid_range(self):
        pres = build_deck(2)
        self.assertIsNone(ppt_utils.check_index(pres, 1))
        message = ppt_utils.check_index(pres, 5)
        self.assertIn("0-1", message)


class TestDuplicate(unittest.TestCase):
    def test_duplicate_appends_a_copy_of_the_content(self):
        pres = build_deck(2)
        new_index = ppt_utils.duplicate_slide(pres, 0)
        self.assertEqual(new_index, 2)
        self.assertEqual(slide_titles(roundtrip(pres)),
                         ["Slide 0", "Slide 1", "Slide 0"])

    def test_duplicate_inserts_after_the_given_index(self):
        pres = build_deck(3)
        new_index = ppt_utils.duplicate_slide(pres, 2, insert_after=0)
        self.assertEqual(new_index, 1)
        self.assertEqual(slide_titles(pres),
                         ["Slide 0", "Slide 2", "Slide 1", "Slide 2"])

    def test_the_copy_is_editable_and_independent(self):
        pres = build_deck(2)
        new_index = ppt_utils.duplicate_slide(pres, 0)
        pres.slides[new_index].shapes.title.text = "Rewritten"
        self.assertEqual(slide_titles(roundtrip(pres)),
                         ["Slide 0", "Slide 1", "Rewritten"])

    def test_the_copy_can_take_new_shapes(self):
        """The copy's shape collection must be bound to the tree that was
        actually kept — the bug an in-place element swap avoids."""
        pres = build_deck(1)
        new_index = ppt_utils.duplicate_slide(pres, 0)
        slide = pres.slides[new_index]
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        box.text_frame.text = "added later"
        reopened = roundtrip(pres)
        texts = [s.text_frame.text for s in reopened.slides[new_index].shapes
                 if s.has_text_frame]
        self.assertIn("added later", texts)

    def test_speaker_notes_travel_with_the_copy(self):
        pres = build_deck(1)
        ppt_utils.set_speaker_notes(pres.slides[0], "say this")
        new_index = ppt_utils.duplicate_slide(pres, 0)
        reopened = roundtrip(pres)
        self.assertEqual(
            ppt_utils.get_speaker_notes(reopened.slides[new_index]), "say this")

    def test_notes_are_not_shared_between_the_copies(self):
        pres = build_deck(1)
        ppt_utils.set_speaker_notes(pres.slides[0], "original")
        new_index = ppt_utils.duplicate_slide(pres, 0)
        ppt_utils.set_speaker_notes(pres.slides[new_index], "the copy")
        reopened = roundtrip(pres)
        self.assertEqual(ppt_utils.get_speaker_notes(reopened.slides[0]),
                         "original")
        self.assertEqual(
            ppt_utils.get_speaker_notes(reopened.slides[new_index]), "the copy")


class TestDuplicateRichSlides(unittest.TestCase):
    """The demo deck carries the parts that make duplication interesting:
    pictures (shared by design), charts and tables (cloned)."""

    @classmethod
    def setUpClass(cls):
        cls.source = Presentation(DEMO)
        cls.chart_slide = cls._first(cls.source, "has_chart")
        cls.table_slide = cls._first(cls.source, "has_table")
        cls.picture_slide = cls._first_picture(cls.source)

    @staticmethod
    def _first(pres, attr):
        for i, slide in enumerate(pres.slides):
            if any(getattr(shape, attr, False) for shape in slide.shapes):
                return i
        raise unittest.SkipTest(f"no slide with {attr} in {DEMO}")

    @staticmethod
    def _first_picture(pres):
        for i, slide in enumerate(pres.slides):
            if any(shape.shape_type == 13 for shape in slide.shapes):
                return i
        raise unittest.SkipTest(f"no picture slide in {DEMO}")

    def test_chart_part_is_cloned_not_shared(self):
        pres = Presentation(DEMO)
        new_index = ppt_utils.duplicate_slide(pres, self.chart_slide)
        original = self._chart(pres.slides[self.chart_slide])
        copy = self._chart(pres.slides[new_index])
        self.assertNotEqual(original.part.partname, copy.part.partname)

    def test_editing_the_copys_chart_leaves_the_original_alone(self):
        pres = Presentation(DEMO)
        new_index = ppt_utils.duplicate_slide(pres, self.chart_slide)
        original = self._chart(pres.slides[self.chart_slide])
        before = list(original.plots[0].categories)

        data = CategoryChartData()
        data.categories = ["only", "these"]
        data.add_series("Replaced", (1, 2))
        self._chart(pres.slides[new_index]).replace_data(data)

        self.assertEqual(list(original.plots[0].categories), before)
        reopened = roundtrip(pres)
        self.assertEqual(
            list(self._chart(reopened.slides[new_index]).plots[0].categories),
            ["only", "these"])

    def test_picture_bytes_are_shared_rather_than_duplicated(self):
        pres = Presentation(DEMO)
        new_index = ppt_utils.duplicate_slide(pres, self.picture_slide)
        original = self._picture(pres.slides[self.picture_slide])
        copy = self._picture(pres.slides[new_index])
        self.assertEqual(original.image.blob, copy.image.blob)
        self.assertEqual(original.image.sha1, copy.image.sha1)

    def test_table_content_survives_duplication(self):
        pres = Presentation(DEMO)
        new_index = ppt_utils.duplicate_slide(pres, self.table_slide)
        reopened = roundtrip(pres)
        original = self._table(reopened.slides[self.table_slide])
        copy = self._table(reopened.slides[new_index])
        self.assertEqual(len(original.rows), len(copy.rows))
        self.assertEqual(original.cell(0, 0).text, copy.cell(0, 0).text)

    def test_shape_count_matches_the_source(self):
        pres = Presentation(DEMO)
        for index in (self.chart_slide, self.table_slide, self.picture_slide):
            with self.subTest(slide=index):
                new_index = ppt_utils.duplicate_slide(pres, index)
                self.assertEqual(len(pres.slides[index].shapes),
                                 len(pres.slides[new_index].shapes))

    @staticmethod
    def _chart(slide):
        return next(s for s in slide.shapes if getattr(s, "has_chart", False)).chart

    @staticmethod
    def _table(slide):
        return next(s for s in slide.shapes if getattr(s, "has_table", False)).table

    @staticmethod
    def _picture(slide):
        return next(s for s in slide.shapes if s.shape_type == 13)


class TestCopyBetweenPresentations(unittest.TestCase):
    def test_copy_lands_in_the_target_deck(self):
        source = build_deck(2)
        target = build_deck(1)
        new_index, _, _ = ppt_utils.copy_slide_to_presentation(source, 1, target)
        self.assertEqual(new_index, 1)
        self.assertEqual(slide_titles(roundtrip(target)),
                         ["Slide 0", "Slide 1"])
        self.assertEqual(len(source.slides), 2, "source must not be modified")

    def test_layout_is_matched_by_name(self):
        source = build_deck(1)
        target = build_deck(1)
        _, layout_name, matched = ppt_utils.copy_slide_to_presentation(
            source, 0, target)
        self.assertTrue(matched)
        self.assertEqual(layout_name, source.slides[0].slide_layout.name)

    def test_pictures_are_cloned_into_the_target_package(self):
        source = Presentation(DEMO)
        picture_slide = TestDuplicateRichSlides._first_picture(source)
        target = Presentation()
        new_index, _, _ = ppt_utils.copy_slide_to_presentation(
            source, picture_slide, target)
        reopened = roundtrip(target)
        pictures = [s for s in reopened.slides[new_index].shapes
                    if s.shape_type == 13]
        self.assertTrue(pictures, "the picture did not survive the copy")
        self.assertTrue(pictures[0].image.blob)

    def test_charts_are_cloned_into_the_target_package(self):
        source = Presentation(DEMO)
        chart_slide = TestDuplicateRichSlides._first(source, "has_chart")
        target = Presentation()
        new_index, _, _ = ppt_utils.copy_slide_to_presentation(
            source, chart_slide, target)
        reopened = roundtrip(target)
        charts = [s for s in reopened.slides[new_index].shapes
                  if getattr(s, "has_chart", False)]
        self.assertTrue(charts, "the chart did not survive the copy")
        self.assertTrue(list(charts[0].chart.plots[0].categories))


class TestSpeakerNotes(unittest.TestCase):
    def test_get_returns_empty_for_a_slide_without_notes(self):
        pres = build_deck(1)
        self.assertEqual(ppt_utils.get_speaker_notes(pres.slides[0]), "")

    def test_set_then_get_roundtrips(self):
        pres = build_deck(1)
        ppt_utils.set_speaker_notes(pres.slides[0], "remember the numbers")
        reopened = roundtrip(pres)
        self.assertEqual(ppt_utils.get_speaker_notes(reopened.slides[0]),
                         "remember the numbers")

    def test_clearing_notes(self):
        pres = build_deck(1)
        ppt_utils.set_speaker_notes(pres.slides[0], "temporary")
        ppt_utils.set_speaker_notes(pres.slides[0], "")
        self.assertEqual(ppt_utils.get_speaker_notes(roundtrip(pres).slides[0]),
                         "")

    def test_notes_are_extracted_but_kept_out_of_the_slide_text(self):
        pres = build_deck(1)
        pres.slides[0].shapes.title.text = "On the slide"
        ppt_utils.set_speaker_notes(pres.slides[0], "In the notes")
        result = ppt_utils.extract_slide_text_content(pres.slides[0])
        self.assertTrue(result["has_notes"])
        self.assertEqual(result["text_content"]["speaker_notes"], "In the notes")
        self.assertNotIn("In the notes",
                         result["text_content"]["all_text_combined"])
        self.assertIn("On the slide",
                      result["text_content"]["all_text_combined"])


class TestDirtyMarking(unittest.TestCase):
    """Speaker notes are not rendered, so they must not cost a re-inspection."""

    def test_speaker_notes_do_not_invalidate_a_visual_pass(self):
        import state
        self.assertIn("manage_speaker_notes", state._NON_EDITING_TOOLS)

    def test_slide_structure_tools_do_invalidate_it(self):
        import state
        for name in ("duplicate_slide", "delete_slide", "move_slide",
                     "copy_slide_between_presentations"):
            with self.subTest(tool=name):
                self.assertNotIn(name, state._NON_EDITING_TOOLS)


if __name__ == "__main__":
    unittest.main()
