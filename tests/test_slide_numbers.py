"""Tests for add_slide_numbers.

The point of the tool is that python-pptx will not do this: a template can
style and position its page numbers on every layout and the slides still come
out blank. So the assertions are about the *slide's* XML, re-read from a saved
deck, not about what the layout says.
"""
import io
import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

from state import PresentationStore
from tools.slide_number_tools import register_slide_number_tools


class _FakeApp:
    def __init__(self):
        self.tools = {}
        self.annotations = {}

    def tool(self, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            self.annotations[fn.__name__] = kwargs.get("annotations")
            return fn
        return decorator


def _sldnum_count(slide):
    total = 0
    for shape in slide.shapes:
        nvSpPr = shape.element.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        ph = None if nvPr is None else nvPr.find(qn("p:ph"))
        if ph is not None and ph.get("type") == "sldNum":
            total += 1
    return total


def _layout_sldnum(layout):
    """The layout's own slide-number placeholder. The stock template puts one
    on every layout — which is the whole point: it renders nothing."""
    for shape in layout.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            return shape
    return None


def deck(slides=3, numbered_layouts=True):
    pres = Presentation()
    layout = pres.slide_layouts[6]  # blank
    if not numbered_layouts:
        # A template says "no number on this layout" by leaving it out.
        placeholder = _layout_sldnum(layout)
        placeholder.element.getparent().remove(placeholder.element)
    for _ in range(slides):
        pres.slides.add_slide(layout)
    return pres


class SlideNumberToolTestCase(unittest.TestCase):
    def setUp(self):
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        self.app = _FakeApp()
        register_slide_number_tools(self.app, self.store)
        self.tool = self.app.tools["add_slide_numbers"]

    def _hold(self, pres):
        pid = self.store.new_id()
        self.store[pid] = pres
        return pid


class TestAddSlideNumbers(SlideNumberToolTestCase):
    def test_the_placeholder_lands_on_every_slide(self):
        pres = deck(3)
        self.assertEqual([_sldnum_count(s) for s in pres.slides], [0, 0, 0])

        result = self.tool(self._hold(pres))

        self.assertEqual(result["added"], 3)
        self.assertEqual([_sldnum_count(s) for s in pres.slides], [1, 1, 1])

    def test_it_survives_a_save(self):
        pres = deck(2)
        self.tool(self._hold(pres))
        buf = io.BytesIO()
        pres.save(buf)
        reopened = Presentation(io.BytesIO(buf.getvalue()))
        self.assertEqual([_sldnum_count(s) for s in reopened.slides], [1, 1])

    def test_calling_twice_adds_nothing(self):
        pres = deck(2)
        pid = self._hold(pres)
        self.tool(pid)
        again = self.tool(pid)
        self.assertEqual(again["added"], 0)
        self.assertEqual(again["already_present"], 2)
        self.assertEqual([_sldnum_count(s) for s in pres.slides], [1, 1])

    def test_skip_slides_leaves_them_unnumbered(self):
        pres = deck(3)
        result = self.tool(self._hold(pres), skip_slides=[0, 2])
        self.assertEqual(result["added"], 1)
        self.assertEqual([_sldnum_count(s) for s in pres.slides], [0, 1, 0])

    def test_a_layout_without_the_placeholder_is_how_a_template_says_no(self):
        pres = deck(2, numbered_layouts=False)
        result = self.tool(self._hold(pres))
        self.assertEqual(result["added"], 0)
        self.assertEqual(result["layout_has_no_slide_number_placeholder"], 2)
        self.assertIn("does not carry page numbers", result["note"])

    def test_the_copy_keeps_the_layouts_geometry(self):
        pres = deck(1)
        source = _layout_sldnum(pres.slide_layouts[6])
        self.tool(self._hold(pres))
        placed = list(pres.slides[0].shapes)[-1]
        self.assertEqual((placed.left, placed.top), (source.left, source.top))

    def test_unknown_presentation_id(self):
        self.assertIn("error", self.tool("nope"))

    def test_it_is_not_marked_read_only(self):
        """The dirty-marking wrapper keys off this: the tool edits the deck."""
        annotations = self.app.annotations["add_slide_numbers"]
        self.assertNotEqual(getattr(annotations, "readOnlyHint", None), True)


if __name__ == "__main__":
    unittest.main()
