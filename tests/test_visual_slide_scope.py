"""Tests for slide-scoped visual QA: the standalone inspect/repair tools,
slide-selection validation, and confinement of repairs to the slides asked for.
"""
import os
import sys
import unittest
from unittest.mock import patch
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches

import visual_fix
import visual_qa
from state import PresentationStore
from tools.visual_tools import register_visual_tools

VISION_ENV = {
    "VISION_LLM_ENDPOINT": "https://example.invalid/responses",
    "VISION_LLM_API_KEY": "k",
    "VISION_LLM_MODEL": "m",
}
ENV_VARS = list(VISION_ENV) + ["VISUAL_QA_ENFORCE", "VISUAL_QA_EXPORT_GATE"]


def make_deck(slides=3):
    pres = Presentation()
    for i in range(slides):
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
    return pres


class _FakeApp:
    """Collects the functions registered with @app.tool()."""

    def __init__(self):
        self.tools = {}

    def tool(self, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn
        return decorator


class ToolTestCase(unittest.TestCase):
    def setUp(self):
        self._saved = {k: os.environ.get(k) for k in ENV_VARS}
        for k in ENV_VARS:
            os.environ.pop(k, None)
        os.environ.update(VISION_ENV)
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        self.pid = self.store.new_id()
        self.store[self.pid] = make_deck()
        self.app = _FakeApp()
        register_visual_tools(self.app, self.store)

    def tearDown(self):
        for k, v in self._saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v


class TestRegistration(ToolTestCase):
    def test_tools_registered_when_llm_configured(self):
        self.assertIn("visual_inspect_slides", self.app.tools)
        self.assertIn("visual_repair_slides", self.app.tools)

    def test_tools_hidden_when_llm_unconfigured(self):
        os.environ.pop("VISION_LLM_MODEL")
        app = _FakeApp()
        register_visual_tools(app, self.store)
        self.assertEqual(app.tools, {})


class TestSlideSelection(ToolTestCase):
    def test_normalize_sorts_and_dedupes(self):
        pres = self.store[self.pid]
        self.assertEqual(visual_qa.normalize_slides(pres, [3, 1, 3]), [1, 3])
        self.assertEqual(visual_qa.normalize_slides(pres, 2), [2])
        self.assertIsNone(visual_qa.normalize_slides(pres, None))
        self.assertIsNone(visual_qa.normalize_slides(pres, []))

    def test_out_of_range_is_reported(self):
        pres = self.store[self.pid]
        with self.assertRaises(ValueError) as ctx:
            visual_qa.normalize_slides(pres, [1, 9])
        self.assertIn("3 slide(s)", str(ctx.exception))

    def test_tool_rejects_bad_slide_numbers(self):
        result = self.app.tools["visual_inspect_slides"](self.pid, slides=[7])
        self.assertIn("error", result)
        self.assertIn("numbered 1-3", result["error"])

    def test_tool_rejects_unknown_presentation(self):
        result = self.app.tools["visual_inspect_slides"]("nope")
        self.assertIn("Unknown or expired", result["error"])


class TestInspectTool(ToolTestCase):
    def test_inspects_only_requested_slides(self):
        seen = {}

        def fake_render(pres, max_slides=None, slides=None):
            seen["slides"] = slides
            return [b"png"] * (len(slides) if slides else 3)

        verdict = {"passed": True, "issues": []}
        with patch.object(visual_qa, "_render_deck", fake_render), \
             patch.object(visual_qa.VisionLLM, "review",
                          lambda *a, **k: dict(verdict)):
            result = self.app.tools["visual_inspect_slides"](self.pid, slides=[2])
        self.assertEqual(seen["slides"], [2])
        self.assertEqual(result["scope"], [2])
        self.assertTrue(result["passed"])

    def test_partial_pass_does_not_clear_the_deck(self):
        verdict = {"passed": True, "issues": []}
        with patch.object(visual_qa, "_render_deck", return_value=[b"png"]), \
             patch.object(visual_qa.VisionLLM, "review",
                          lambda *a, **k: dict(verdict)):
            self.app.tools["visual_inspect_slides"](self.pid, slides=[2])
            self.assertTrue(self.store.is_dirty(self.pid))
            self.app.tools["visual_inspect_slides"](self.pid)
        self.assertFalse(self.store.is_dirty(self.pid))

    def test_prompt_maps_images_to_absolute_slide_numbers(self):
        prompt = visual_qa.review_prompt(False, None, [2, 5])
        self.assertIn("image 1 = slide 2", prompt)
        self.assertIn("image 2 = slide 5", prompt)


class TestRepairScope(ToolTestCase):
    def test_out_of_scope_operations_are_skipped(self):
        pres = self.store[self.pid]
        ops = [{"op": "move_shape", "slide": 1, "shape_index": 0,
                "left_in": 2.0, "top_in": 2.0},
               {"op": "move_shape", "slide": 3, "shape_index": 0,
                "left_in": 2.0, "top_in": 2.0}]
        result = visual_fix.apply_repairs(pres, ops, allowed_slides=[1])
        self.assertEqual(len(result["applied"]), 1)
        self.assertEqual(result["skipped"][0]["reason"], "slide out of scope")

    def test_plan_attaches_the_image_of_the_reported_slide(self):
        pres = self.store[self.pid]
        captured = {}

        class FakeLLM:
            def ask_json(self, images, prompt):
                captured["images"] = images
                return {"operations": []}

        # Only slides 2 and 3 were rendered, in that order.
        visual_fix.plan_repairs(FakeLLM(), [{"slide": 3, "description": "x"}],
                                pres, [b"img-2", b"img-3"], [2, 3])
        self.assertEqual(captured["images"], [b"img-3"])

    def test_issues_outside_scope_are_ignored(self):
        verdicts = [{"passed": False,
                     "issues": [{"slide": 1, "description": "not in scope"}]}]

        with patch.object(visual_qa, "_render_deck", return_value=[b"png"]), \
             patch.object(visual_qa.VisionLLM, "review",
                          lambda *a, **k: dict(verdicts[0])), \
             patch.object(visual_fix, "plan_repairs") as plan:
            outcome = visual_qa.inspect_and_repair(self.store[self.pid],
                                                   slides=[2])
        plan.assert_not_called()
        self.assertTrue(outcome["passed"])

    def test_repair_tool_confines_operations_to_scope(self):
        seen = {}

        def fake_inspect(pres, slides=None, focus=None, max_iterations=None):
            seen.update(slides=slides, max_iterations=max_iterations)
            return {"passed": True, "iterations": 1, "repair_rounds": []}

        with patch.object(visual_qa, "inspect_and_repair", fake_inspect):
            result = self.app.tools["visual_repair_slides"](
                self.pid, slides=[3, 1], max_iterations=2)
        self.assertEqual(seen["slides"], [1, 3])
        self.assertEqual(seen["max_iterations"], 2)
        self.assertEqual(result["scope"], [1, 3])
        # A scoped pass does not certify the whole deck.
        self.assertTrue(self.store.is_dirty(self.pid))


if __name__ == "__main__":
    unittest.main()
