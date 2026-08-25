"""
Render-backed checks for the XML this project assembles by hand.

Duplicated slides, cloned chart parts and combo-chart plot groups are all built
at the OOXML level, and the failure mode that matters is a package python-pptx
will happily re-read but a real renderer rejects or silently drops. Every other
test in the suite reads the XML back with the library that wrote it, which
cannot catch that; these put the deck through LibreOffice.

Assertions are on **text extracted from the rendered PDF**, not on image size.
That distinction is the point of this module: a combo chart whose line group
names an axis that was never declared still renders its columns, so it is only
~45% smaller than a correct one and comfortably passes any "is it blank" check
— but the secondary axis' tick labels are missing from the render, which is
unambiguous. The same text-level comparison is what makes the duplicated-slide
check meaningful rather than a size heuristic.

Skipped when LibreOffice is absent, like the rest of the renderer-dependent
suite; CI installs libreoffice-impress, so they run there. Each soffice
invocation costs a few seconds, so the deck is built and rendered once.
"""
import io
import unittest

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

import utils as ppt_utils
import visual_qa

DEMO = "mcp_all_tools_templates_effects_demo.pptx"


def renderer_available():
    try:
        visual_qa._soffice_binary()
    except visual_qa.VisualQAError:
        return False
    return True


def page_texts(blob):
    """Render to PDF and return the text of every page."""
    import pymupdf

    with pymupdf.open(stream=visual_qa.render_pptx_bytes_to_pdf(blob),
                      filetype="pdf") as doc:
        return [page.get_text() for page in doc]


def combo_deck(break_secondary_axis=False):
    """A one-slide deck holding a column+line combo chart on two axes.

    `break_secondary_axis` points the line group at an axis id nothing
    declares — the exact corruption apply_combo_layout exists to avoid, used
    to prove the assertions below can actually detect it.
    """
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4"]
    data.add_series("Revenue", (10, 12, 15, 18))
    data.add_series("Margin", (0.21, 0.23, 0.22, 0.25))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1), Inches(8),
        Inches(4.5), data).chart
    ppt_utils.apply_combo_layout(
        chart, [{"type": "column"},
                {"type": "line_markers", "secondary_axis": True}])
    chart.has_legend = True

    if break_secondary_axis:
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        for group in plot_area:
            if group.tag.endswith("lineChart"):
                for axis in group.findall(qn("c:axId")):
                    axis.set("val", "123456789")

    buffer = io.BytesIO()
    pres.save(buffer)
    return buffer.getvalue()


@unittest.skipUnless(renderer_available(), "LibreOffice not installed")
class TestComboChartRenders(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.text = page_texts(combo_deck())[0]

    def test_both_series_reach_the_render(self):
        self.assertIn("Revenue", self.text)
        self.assertIn("Margin", self.text)

    def test_the_primary_axis_is_scaled_to_the_column_series(self):
        self.assertIn("20", self.text)

    def test_the_secondary_axis_is_drawn_with_its_own_scale(self):
        """Tick labels in the 0.19-0.26 range can only come from a value axis
        declared for the line series — the primary axis runs 0-20."""
        self.assertIn("0.26", self.text)
        self.assertIn("0.21", self.text)

    def test_the_assertions_detect_an_undeclared_axis(self):
        """Guards against this whole class passing vacuously: with the line
        group pointed at an axis nothing declares, the renderer drops the
        secondary scale (while still drawing the columns)."""
        broken = page_texts(combo_deck(break_secondary_axis=True))[0]
        self.assertNotIn("0.26", broken)
        self.assertIn("Revenue", broken, "the columns should still render")


@unittest.skipUnless(renderer_available(), "LibreOffice not installed")
class TestDuplicatedSlideRenders(unittest.TestCase):
    """A duplicated template slide must render exactly like its source: same
    text, from a cloned chart part and a rebuilt relationship set."""

    @classmethod
    def setUpClass(cls):
        pres = Presentation(DEMO)
        cls.source_index = next(
            i for i, slide in enumerate(pres.slides)
            if any(getattr(s, "has_chart", False) for s in slide.shapes))
        cls.copy_index = ppt_utils.duplicate_slide(pres, cls.source_index)
        ppt_utils.set_speaker_notes(pres.slides[cls.copy_index], "notes")

        buffer = io.BytesIO()
        pres.save(buffer)
        cls.blob = buffer.getvalue()
        cls.pages = page_texts(cls.blob)

    def test_the_whole_deck_rendered(self):
        self.assertEqual(len(self.pages), len(Presentation(DEMO).slides) + 1)

    def test_the_copy_renders_the_same_text_as_its_source(self):
        self.assertEqual(self.pages[self.copy_index],
                         self.pages[self.source_index])

    def test_that_text_is_not_empty(self):
        """Two blank pages would also compare equal."""
        self.assertGreater(len(self.pages[self.copy_index].strip()), 20)

    def test_the_deck_converts_to_pdf(self):
        pdf = visual_qa.render_pptx_bytes_to_pdf(self.blob)
        self.assertTrue(pdf.startswith(b"%PDF-"))
        self.assertGreater(len(pdf), 10000)

    def test_selected_slides_rasterize(self):
        wanted = [self.source_index + 1, self.copy_index + 1]
        images = visual_qa.render_pptx_bytes_to_pngs(
            self.blob, dpi=80, slides=wanted)
        self.assertEqual(len(images), 2)
        for png in images:
            self.assertTrue(png.startswith(b"\x89PNG"))


@unittest.skipUnless(renderer_available(), "LibreOffice not installed")
class TestLegacyPptRoundTrip(unittest.TestCase):
    def test_a_real_binary_ppt_can_be_imported(self):
        with open(DEMO, "rb") as handle:
            source = handle.read()
        legacy = visual_qa.convert_with_soffice(source, ".pptx", "ppt")
        self.assertTrue(legacy.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"),
                        "soffice did not produce a binary PowerPoint file")

        converted = visual_qa.convert_legacy_ppt(legacy)
        self.assertTrue(converted.startswith(b"PK"))
        reopened = Presentation(io.BytesIO(converted))
        self.assertGreater(len(reopened.slides), 0)


if __name__ == "__main__":
    unittest.main()
