"""Tests for the brand profile pass.

Two things are being defended here. One is that each rule actually fires on a
deck that breaks it and stays quiet on one that doesn't — a profile that
reports nothing is worse than no profile, because the agent believes it.

The other is the plumbing around them: the deployment names a *file*, the
orchestrator supplies a URL for it, and the rules end up bound to one deck.
Unset means no tools at all; a deck with nothing attached is told to attach
rather than reported clean; and one deck's profile never reaches another's.
"""
import json
import os
import sys
import unittest
from pathlib import Path
from unittest.mock import patch

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import brand_validation
from state import PresentationStore
from tools.brand_tools import register_brand_tools

PROFILE = {
    "name": "Example Corp",
    "fonts": {"allowed": ["Example Sans"]},
    "min_font_pt": 14,
    "min_font_exempt_idx": [12],
    "palette_rgb": [[0, 93, 185], [255, 255, 255], [0, 0, 0]],
    "palette_tolerance": 6,
    "families": {
        "light": {"safe_area_in": {"x": 0.5, "y": 1.5, "w": 9.0, "h": 5.0},
                  "require": ["slide_number"]},
        "dark": {"safe_area_in": {"x": 0.5, "y": 1.0, "w": 9.0, "h": 5.5},
                 "require": []},
    },
    "chrome_idx": [12],
    "chrome_shapes": {"slide_number": {"type": "placeholder", "idx": 12}},
    "require_alpha_on_fills": True,
    "max_consecutive_same_family": 2,
    "exempt_layouts": ["Title Slide"],
    "review_notes": ["Headlines must be messages.", "No plain text on white."],
}


class _FakeApp:
    def __init__(self):
        self.tools = {}
        self.annotations = {}

    def tool(self, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            self.annotations[fn.__name__] = kwargs.get("annotations")
            return fn
        return decorator


def blank_deck(slides=1):
    pres = Presentation()
    for _ in range(slides):
        pres.slides.add_slide(pres.slide_layouts[6])
    return pres


def textbox(slide, text="Hello", size=None, font=None, color=None,
            left=1.0, top=2.0, width=3.0, height=1.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                   Inches(height))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    if size:
        run.font.size = Pt(size)
    if font:
        run.font.name = font
    if color:
        run.font.color.rgb = RGBColor(*color)
    return box


def codes(report, severity=None):
    return {p["code"] for p in report["problems"]
            if severity is None or p["severity"] == severity}


def check(pres, **overrides):
    profile = dict(PROFILE, **overrides)
    return brand_validation.validate_brand(pres, profile)


class TestFontsAndSizes(unittest.TestCase):
    def test_a_font_outside_the_brand_is_reported(self):
        pres = blank_deck()
        textbox(pres.slides[0], font="Comic Sans MS", size=18)
        self.assertIn("brand_font_not_allowed", codes(check(pres)))

    def test_the_brand_font_is_not_reported(self):
        pres = blank_deck()
        textbox(pres.slides[0], font="Example Sans", size=18)
        self.assertNotIn("brand_font_not_allowed", codes(check(pres)))

    def test_inherited_text_is_on_brand_by_definition(self):
        """Text that names no font renders in the template's own face."""
        pres = blank_deck()
        textbox(pres.slides[0], size=18)
        self.assertNotIn("brand_font_not_allowed", codes(check(pres)))

    def test_text_below_the_floor_is_reported_with_its_size(self):
        pres = blank_deck()
        textbox(pres.slides[0], size=11)
        problems = [p for p in check(pres)["problems"]
                    if p["code"] == "brand_text_too_small"]
        self.assertEqual(len(problems), 1)
        self.assertIn("11pt", problems[0]["message"])
        self.assertEqual(problems[0]["slide_index"], 0)

    def test_text_at_the_floor_is_fine(self):
        pres = blank_deck()
        textbox(pres.slides[0], size=14)
        self.assertNotIn("brand_text_too_small", codes(check(pres)))


class TestPalette(unittest.TestCase):
    def test_an_off_palette_fill_is_reported(self):
        pres = blank_deck()
        shape = pres.slides[0].shapes.add_shape(1, Inches(1), Inches(2),
                                                Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 12, 200)
        self.assertIn("brand_colour_off_palette", codes(check(pres)))

    def test_a_palette_colour_within_tolerance_is_fine(self):
        pres = blank_deck()
        shape = pres.slides[0].shapes.add_shape(1, Inches(1), Inches(2),
                                                Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(3, 95, 187)  # ~[0, 93, 185]
        self.assertNotIn("brand_colour_off_palette", codes(check(pres)))

    def test_a_text_colour_is_checked_too(self):
        pres = blank_deck()
        textbox(pres.slides[0], size=18, color=(200, 30, 30))
        problems = [p for p in check(pres)["problems"]
                    if p["code"] == "brand_colour_off_palette"]
        self.assertTrue(any("text" in p["message"] for p in problems))


class TestSafeArea(unittest.TestCase):
    def test_a_shape_inside_the_frame_is_fine(self):
        pres = blank_deck()
        textbox(pres.slides[0], size=18, left=1.0, top=2.0)
        self.assertNotIn("brand_outside_safe_area", codes(check(pres)))

    def test_a_shape_above_the_frame_is_reported(self):
        """The slide's own edge is 0in; the brand's frame starts at 1.5in."""
        pres = blank_deck()
        textbox(pres.slides[0], size=18, top=0.6)
        self.assertIn("brand_outside_safe_area", codes(check(pres)))

    def test_chrome_placeholders_are_exempt(self):
        """The template puts the page number outside the frame on purpose."""
        pres = blank_deck()
        pid_store = PresentationStore(ttl_seconds=60, max_items=5)
        pid = pid_store.new_id()
        pid_store[pid] = pres
        app = _FakeApp()
        from tools.slide_number_tools import register_slide_number_tools

        register_slide_number_tools(app, pid_store)
        app.tools["add_slide_numbers"](pid)
        number = list(pres.slides[0].shapes)[-1]
        self.assertEqual(number.placeholder_format.idx, 12)
        self.assertGreater(number.top, Inches(6.5))  # below the frame

        textbox(pres.slides[0], size=18, top=0.6)  # content, and outside it
        report = check(pres)
        offenders = {p["shape_index"] for p in report["problems"]
                     if p["code"] == "brand_outside_safe_area"}
        self.assertEqual(offenders, {1})


class TestBackgroundFamily(unittest.TestCase):
    def _full_bleed(self, slide, rgb):
        shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        return shape

    def test_a_white_slide_is_light(self):
        pres = blank_deck()
        self.assertEqual(brand_validation.classify_background(pres.slides[0]),
                         "light")

    def test_a_full_bleed_dark_shape_makes_the_slide_dark(self):
        pres = blank_deck()
        self._full_bleed(pres.slides[0], (0, 61, 122))
        self.assertEqual(brand_validation.classify_background(pres.slides[0]),
                         "dark")

    def test_a_small_dark_shape_does_not(self):
        pres = blank_deck()
        shape = pres.slides[0].shapes.add_shape(1, Inches(1), Inches(1),
                                                Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        self.assertEqual(brand_validation.classify_background(pres.slides[0]),
                         "light")

    def test_too_many_slides_in_a_row_on_one_family(self):
        pres = blank_deck(4)  # four white slides, limit is 2
        report = check(pres)
        problems = [p for p in report["problems"]
                    if p["code"] == "brand_monotonous_deck"]
        self.assertEqual(len(problems), 1)
        self.assertEqual(report["backgrounds"], ["light"] * 4)

    def test_an_alternating_deck_is_fine(self):
        pres = blank_deck(4)
        for index in (1, 3):
            self._full_bleed(pres.slides[index], (0, 61, 122))
        self.assertNotIn("brand_monotonous_deck", codes(check(pres)))


class TestRequiredChrome(unittest.TestCase):
    def test_a_missing_slide_number_is_reported(self):
        pres = blank_deck()
        report = check(pres)
        problems = [p for p in report["problems"]
                    if p["code"] == "brand_chrome_missing"]
        self.assertEqual(len(problems), 1)
        self.assertIn("add_slide_numbers", problems[0]["fix"])

    def test_the_placeholder_on_the_slide_satisfies_it(self):
        """An inherited placeholder renders nothing, so only the slide counts."""
        pres = blank_deck()
        from tools.slide_number_tools import register_slide_number_tools

        store = PresentationStore(ttl_seconds=60, max_items=5)
        pid = store.new_id()
        store[pid] = pres
        app = _FakeApp()
        register_slide_number_tools(app, store)
        app.tools["add_slide_numbers"](pid)
        self.assertNotIn("brand_chrome_missing", codes(check(pres)))

    def test_an_exempt_layout_is_not_checked(self):
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[0])  # "Title Slide"
        self.assertNotIn("brand_chrome_missing", codes(check(pres)))


class BrandToolsTestCase(unittest.TestCase):
    """The brand files live in DIAL storage and arrive as caller-supplied
    URLs, so the client is stubbed: what is under test is which files this
    server accepts, what it does with the bytes, and what it says when they
    do not arrive."""

    PROFILE_FILE = "brand_profile.json"
    REFERENCE_FILE = "brand_reference.pptx"

    def setUp(self):
        self._saved = {k: os.environ.get(k) for k in
                       ("BRAND_PROFILE_FILE", "BRAND_REFERENCE_DECK_FILE")}
        os.environ["BRAND_PROFILE_FILE"] = self.PROFILE_FILE
        os.environ["BRAND_REFERENCE_DECK_FILE"] = self.REFERENCE_FILE
        self.addCleanup(self._restore)
        self.store = PresentationStore(ttl_seconds=60, max_items=5)
        self.pid = self.store.new_id()
        self.store[self.pid] = blank_deck()
        self.downloads = []
        self.files = {}
        self.app = _FakeApp()
        register_brand_tools(self.app, self.store)

    def _restore(self):
        for key, value in self._saved.items():
            if value is None:
                os.environ.pop(key, None)
            else:
                os.environ[key] = value

    def _serving(self):
        recorded, files = self.downloads, self.files

        class FakeClient:
            def __init__(self, *args, **kwargs):
                pass

            def download(self, ref):
                recorded.append(ref)
                if ref not in files:
                    from dial_client import DialConfigError
                    raise DialConfigError("DIAL Core refused this file (404).")
                return files[ref]

        return patch("dial_client.DialFileClient", FakeClient)

    def _deck_bytes(self):
        import io

        buffer = io.BytesIO()
        blank_deck().save(buffer)
        return buffer.getvalue()

    def attach(self, profile_url=None, **kwargs):
        with self._serving():
            return self.app.tools["attach_brand_profile"](
                self.pid, profile_url or f"files/b/anywhere/{self.PROFILE_FILE}",
                **kwargs)


class TestAttachingTheProfile(BrandToolsTestCase):
    def setUp(self):
        super().setUp()
        self.files[f"files/b/anywhere/{self.PROFILE_FILE}"] = \
            json.dumps(PROFILE).encode()

    def test_the_rules_are_read_and_bound_to_the_deck(self):
        result = self.attach()
        self.assertEqual(result["brand"], "Example Corp")
        self.assertEqual(
            self.store.brand_for(self.pid)["profile"]["name"], "Example Corp")

    def test_any_bucket_or_folder_is_accepted_only_the_name_matters(self):
        """Buckets move; the file name is the stable half of the contract."""
        for url in (f"files/other-bucket/{self.PROFILE_FILE}",
                    f"https://dial.example.com/v1/files/b/deep/path/"
                    f"{self.PROFILE_FILE}"):
            with self.subTest(url=url):
                self.files[url] = json.dumps(PROFILE).encode()
                self.assertNotIn("error", self.attach(url))

    def test_a_differently_named_file_is_refused(self):
        """The one check between 'the brand profile' and 'any JSON the caller
        happens to have'."""
        url = "files/b/notes.json"
        self.files[url] = json.dumps(PROFILE).encode()
        result = self.attach(url)
        self.assertIn("error", result)
        self.assertIn("notes.json", result["error"])
        self.assertIn(self.PROFILE_FILE, result["error"])
        self.assertIsNone(self.store.brand_for(self.pid))
        self.assertEqual(self.downloads, [], "refused before any fetch")

    def test_the_returned_rules_let_the_agent_build_to_them(self):
        rules = self.attach()["rules"]
        self.assertEqual(rules["min_font_pt"], 14)
        self.assertEqual(rules["fonts"], ["Example Sans"])
        self.assertIn("light", rules["safe_area_in"])
        self.assertIn("Headlines must be messages.", rules["review_notes"])

    def test_a_missing_file_reports_the_url(self):
        result = self.attach("files/b/gone/brand_profile.json")
        self.assertIn("error", result)
        self.assertIn("files/b/gone/brand_profile.json", result["error"])

    def test_a_file_that_is_not_json_says_so(self):
        url = f"files/b/{self.PROFILE_FILE}"
        self.files[url] = b"{not json"
        self.assertIn("not valid JSON", self.attach(url)["error"])

    def test_a_json_array_is_not_a_profile(self):
        url = f"files/b/{self.PROFILE_FILE}"
        self.files[url] = b"[1, 2]"
        self.assertIn("error", self.attach(url))

    def test_the_shipped_example_is_accepted(self):
        root = Path(__file__).resolve().parent.parent
        url = f"files/b/{self.PROFILE_FILE}"
        self.files[url] = (root / "brand_profile.example.json").read_bytes()
        self.assertEqual(self.attach(url)["brand"], "Example Corp")

    def test_unknown_presentation_id(self):
        with self._serving():
            result = self.app.tools["attach_brand_profile"](
                "nope", f"files/b/{self.PROFILE_FILE}")
        self.assertIn("error", result)

    def test_the_description_names_the_file_to_look_for(self):
        """An agent that is not told what to look for cannot find it."""
        doc = self.app.tools["attach_brand_profile"].__doc__
        self.assertIn(self.PROFILE_FILE, doc)
        self.assertIn(self.REFERENCE_FILE, doc)


class TestAttachingTheReferenceDeck(BrandToolsTestCase):
    def setUp(self):
        super().setUp()
        self.profile_url = f"files/b/{self.PROFILE_FILE}"
        self.files[self.profile_url] = json.dumps(PROFILE).encode()

    def test_the_deck_is_opened_and_bound(self):
        url = f"files/b/{self.REFERENCE_FILE}"
        self.files[url] = self._deck_bytes()
        result = self.attach(self.profile_url, reference_deck_url=url)
        self.assertTrue(result["reference_deck"])
        self.assertEqual(
            len(self.store.brand_for(self.pid)["reference"].slides), 1)

    def test_omitting_it_is_fine_but_noted(self):
        result = self.attach(self.profile_url)
        self.assertFalse(result["reference_deck"])
        self.assertIn(self.REFERENCE_FILE, result["reference_deck_note"])

    def test_a_differently_named_deck_is_refused(self):
        url = "files/b/some_other_deck.pptx"
        self.files[url] = self._deck_bytes()
        result = self.attach(self.profile_url, reference_deck_url=url)
        self.assertIn("some_other_deck.pptx", result["error"])

    def test_a_file_that_is_not_a_pptx_says_so(self):
        url = f"files/b/{self.REFERENCE_FILE}"
        self.files[url] = b"not a pptx"
        result = self.attach(self.profile_url, reference_deck_url=url)
        self.assertIn("readable .pptx", result["error"])


class TestValidateBrandProfileTool(BrandToolsTestCase):
    def setUp(self):
        super().setUp()
        self.files[f"files/b/anywhere/{self.PROFILE_FILE}"] = \
            json.dumps(PROFILE).encode()
        self.validate = self.app.tools["validate_brand_profile"]

    def test_without_rules_attached_it_says_what_to_do(self):
        result = self.validate(self.pid)
        self.assertIn("error", result)
        self.assertIn("attach_brand_profile", result["error"])
        self.assertIn(self.PROFILE_FILE, result["error"])

    def test_with_rules_attached_it_reports(self):
        self.attach()
        result = self.validate(self.pid)
        self.assertEqual(result["brand"], "Example Corp")
        self.assertIn("problems", result)

    def test_argument_validation_still_applies(self):
        self.attach()
        self.assertIn("error", self.validate("nope"))
        self.assertIn("error", self.validate(self.pid, min_severity="loud"))

    def test_rules_do_not_leak_between_decks(self):
        """Brand context is a property of the deck, not of the server: a
        profile's review_notes become prompt text for the vision reviewer, so
        one caller's file must not reach another caller's review."""
        self.attach()
        other = self.store.new_id()
        self.store[other] = blank_deck()
        self.assertIsNone(self.store.brand_for(other))
        self.assertIn("error", self.validate(other))


class TestRegistration(unittest.TestCase):
    def setUp(self):
        self._saved = os.environ.get("BRAND_PROFILE_FILE")
        os.environ.pop("BRAND_PROFILE_FILE", None)
        self.addCleanup(self._restore)

    def _restore(self):
        if self._saved is None:
            os.environ.pop("BRAND_PROFILE_FILE", None)
        else:
            os.environ["BRAND_PROFILE_FILE"] = self._saved

    def test_no_profile_configured_no_tools(self):
        app = _FakeApp()
        register_brand_tools(app, PresentationStore())
        self.assertEqual(app.tools, {})

    def test_a_configured_name_registers_both_and_needs_no_network(self):
        os.environ["BRAND_PROFILE_FILE"] = "brand_profile.json"
        app = _FakeApp()
        register_brand_tools(app, PresentationStore())
        self.assertEqual(sorted(app.tools),
                         ["attach_brand_profile", "validate_brand_profile"])

    def test_attaching_rules_stales_an_earlier_inspection(self):
        """attach_brand_profile is not read-only, and the dirty-marking
        wrapper keys off that. The rules change what the reviewer is told to
        look for and give it a reference deck, so a deck inspected clean
        before they arrived was judged against different criteria — exporting
        it as "passed" would credit it with a review that never saw the
        brand. Checking the deck against the rules changes nothing, so
        validate_brand_profile stays read-only."""
        os.environ["BRAND_PROFILE_FILE"] = "brand_profile.json"
        app = _FakeApp()
        register_brand_tools(app, PresentationStore())
        self.assertNotEqual(
            getattr(app.annotations["attach_brand_profile"],
                    "readOnlyHint", None), True)
        self.assertTrue(
            app.annotations["validate_brand_profile"].readOnlyHint)


class TestFileNameMatching(unittest.TestCase):
    def test_the_name_is_taken_from_any_reference_form(self):
        for url, expected in (
                ("files/b/brand_profile.json", "brand_profile.json"),
                ("files/b/deep/folder/brand_profile.json",
                 "brand_profile.json"),
                ("https://dial.example.com/v1/files/b/x/brand_profile.json",
                 "brand_profile.json"),
                ("files/b/brand_profile.json?v=2", "brand_profile.json")):
            with self.subTest(url=url):
                self.assertEqual(brand_validation.file_name_of(url), expected)

    def test_matching_is_case_insensitive(self):
        brand_validation.check_file_name("files/b/Brand_Profile.JSON",
                                         "brand_profile.json", "brand profile")

    def test_review_focus_joins_the_notes(self):
        focus = brand_validation.review_focus(PROFILE)
        self.assertIn("Headlines must be messages.", focus)
        self.assertIn("No plain text on white.", focus)
        self.assertIsNone(brand_validation.review_focus(None))
        self.assertIsNone(brand_validation.review_focus({}))


if __name__ == "__main__":
    unittest.main()
