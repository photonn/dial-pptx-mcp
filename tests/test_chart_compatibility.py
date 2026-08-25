"""
Tests for the chart settings that render differently in PowerPoint and
LibreOffice when left implicit.

python-pptx omits several optional ``c:`` elements. ECMA-376 reads a missing
boolean as *true*, which PowerPoint honours and LibreOffice does not, so the
assertions here are all of the form "the element is present and says what we
mean" — an absent element is the bug, not a neutral default.
"""
import io
import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.oxml.ns import qn

import utils as ppt_utils
from state import PresentationStore
from tools.structural_tools import register_structural_tools

SINGLE_SERIES = (["Velociraptor", "T-Rex", "Brachiosaurus"], ["Weight"],
                 [[0.08, 8.0, 56.0]])


def _blank_slide():
    pres = Presentation()
    return pres, pres.slides.add_slide(pres.slide_layouts[6])


def _build(chart_type, **format_kwargs):
    """Build a chart the way the add_chart tool does, and reopen the deck."""
    pres, slide = _blank_slide()
    categories, names, values = SINGLE_SERIES
    chart = ppt_utils.add_chart(slide, chart_type, 1, 1, 8, 5,
                                categories, names, values)
    ppt_utils.format_chart(chart, **format_kwargs)
    stream = io.BytesIO()
    pres.save(stream)
    stream.seek(0)
    reopened = Presentation(stream)
    return reopened.slides[0].shapes[0].chart


def _plot_groups(chart):
    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    return [el for el in plot_area if el.tag.split("}")[-1].endswith("Chart")]


class TestVaryColors(unittest.TestCase):
    def test_series_charts_state_one_colour_per_series(self):
        # Absent, PowerPoint colours a single-series bar chart per category
        # and lists the categories in the legend; LibreOffice does not.
        for chart_type in ("column", "bar", "line", "area", "radar",
                           "stacked_column", "stacked_bar"):
            with self.subTest(chart_type=chart_type):
                chart = _build(chart_type)
                groups = _plot_groups(chart)
                self.assertTrue(groups)
                for group in groups:
                    vary = group.find(qn("c:varyColors"))
                    self.assertIsNotNone(vary, "c:varyColors missing")
                    self.assertEqual(vary.get("val"), "0")

    def test_pie_charts_keep_per_point_colours(self):
        for chart_type in ("pie", "doughnut"):
            with self.subTest(chart_type=chart_type):
                chart = _build(chart_type)
                for group in _plot_groups(chart):
                    vary = group.find(qn("c:varyColors"))
                    self.assertIsNotNone(vary)
                    self.assertEqual(vary.get("val"), "1")


class TestScatterCharts(unittest.TestCase):
    """A scatter chart plots x/y pairs and needs XyChartData.

    add_chart advertised 'scatter' as a valid type while building
    CategoryChartData for it, which raises on every call — the type had never
    worked. The x values ride in the `categories` argument because that is the
    one list the tool signature has for them.
    """

    def build(self, categories, values=((1.0, 4.0, 9.0),)):
        pres, slide = _blank_slide()
        chart = ppt_utils.add_chart(slide, "scatter", 1, 1, 8, 5,
                                    list(categories), ["S1"],
                                    [list(v) for v in values])
        stream = io.BytesIO()
        pres.save(stream)
        stream.seek(0)
        return Presentation(stream).slides[0].shapes[0].chart

    def test_a_scatter_chart_builds_and_reopens(self):
        chart = self.build([1, 2, 3])
        groups = _plot_groups(chart)
        self.assertEqual([g.tag.split("}")[-1] for g in groups],
                         ["scatterChart"])

    def test_categories_become_the_x_values(self):
        chart = self.build([10, 20, 30])
        series = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea")) \
            .find(qn("c:scatterChart")).find(qn("c:ser"))
        x_pts = series.find(qn("c:xVal")).find(qn("c:numRef")) \
            .find(qn("c:numCache")).findall(qn("c:pt"))
        self.assertEqual([p.find(qn("c:v")).text for p in x_pts],
                         ["10.0", "20.0", "30.0"])

    def test_y_values_are_the_series_values(self):
        chart = self.build([1, 2, 3], values=[(1.5, 4.5, 9.5)])
        self.assertEqual(list(chart.plots[0].series[0].values),
                         [1.5, 4.5, 9.5])

    def test_a_scatter_chart_gets_the_same_explicit_defaults(self):
        chart = self.build([1, 2, 3])
        for group in _plot_groups(chart):
            self.assertEqual(group.find(qn("c:varyColors")).get("val"), "0")

    def test_non_numeric_x_values_are_refused_with_an_instruction(self):
        with self.assertRaises(ValueError) as caught:
            ppt_utils.parse_scatter_x_values(["Q1", "Q2"])
        message = str(caught.exception)
        self.assertIn("Q1", message)
        self.assertIn("line_markers", message)

    def test_numeric_strings_are_accepted(self):
        self.assertEqual(ppt_utils.parse_scatter_x_values(["1", "2.5"]),
                         [1.0, 2.5])


class TestAddChartTool(unittest.TestCase):
    """The tool wrapper around add_chart, for the scatter argument contract."""

    def setUp(self):
        app = FastMCP(name="test")
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        register_structural_tools(
            app, self.store, lambda: None,
            lambda params: (True, None),
            lambda v: v > 0, lambda v: v >= 0,
            lambda lo, hi: (lambda v: lo <= v <= hi),
            lambda v: True, lambda *a, **k: None)
        self.add_chart = app._tool_manager._tools["add_chart"].fn
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[6])
        self.pres_id = self.store.new_id()
        self.store[self.pres_id] = pres

    def call(self, categories, chart_type="scatter"):
        return self.add_chart(
            slide_index=0, chart_type=chart_type, left=1, top=1,
            width=6, height=4, categories=categories,
            series_names=["S1"], series_values=[[1.0, 2.0, 3.0]],
            presentation_id=self.pres_id)

    def test_numeric_categories_produce_a_scatter_chart(self):
        result = self.call([1, 2, 3])
        self.assertNotIn("error", result)
        chart = self.store[self.pres_id].slides[0].shapes[0].chart
        self.assertEqual(
            [g.tag.split("}")[-1] for g in _plot_groups(chart)],
            ["scatterChart"])

    def test_label_categories_return_an_instruction_not_a_traceback(self):
        result = self.call(["Q1", "Q2", "Q3"])
        self.assertIn("error", result)
        self.assertIn("line_markers", result["error"])
        self.assertNotIn("Failed to add chart", result["error"])

    def test_label_categories_are_fine_on_other_types(self):
        self.assertNotIn("error", self.call(["Q1", "Q2", "Q3"], "column"))


class TestLegend(unittest.TestCase):
    def _legend(self, **kwargs):
        return _build("column", **kwargs).legend

    def test_overlay_is_explicit(self):
        # The absent c:overlay is why PowerPoint laid the legend over the bars.
        legend = self._legend(has_legend=True, legend_position="right")
        self.assertIsNotNone(legend._element.find(qn("c:overlay")))
        self.assertFalse(legend.include_in_layout)

    def test_requested_position_is_applied(self):
        legend = self._legend(has_legend=True, legend_position="bottom")
        self.assertEqual(legend._element.find(qn("c:legendPos")).get("val"), "b")

    def test_legend_can_still_be_suppressed(self):
        chart = _build("column", has_legend=False)
        self.assertFalse(chart.has_legend)


class TestChartElementDefaults(unittest.TestCase):
    def test_plot_vis_only_is_explicit(self):
        chart = _build("column")
        chart_el = chart._chartSpace.find(qn("c:chart"))
        self.assertEqual(
            chart_el.find(qn("c:plotVisOnly")).get("val"), "1")

    def test_untitled_chart_says_so(self):
        # PowerPoint auto-titles a single-series chart from the series name;
        # LibreOffice leaves it blank. autoTitleDeleted settles it.
        chart = _build("column", title=None)
        chart_el = chart._chartSpace.find(qn("c:chart"))
        self.assertEqual(
            chart_el.find(qn("c:autoTitleDeleted")).get("val"), "1")
        self.assertFalse(chart.has_title)

    def test_requested_title_survives(self):
        chart = _build("column", title="Maximum Estimated Weights")
        self.assertTrue(chart.has_title)
        self.assertEqual(chart.chart_title.text_frame.text,
                         "Maximum Estimated Weights")
        chart_el = chart._chartSpace.find(qn("c:chart"))
        self.assertEqual(
            chart_el.find(qn("c:autoTitleDeleted")).get("val"), "0")


if __name__ == "__main__":
    unittest.main()
