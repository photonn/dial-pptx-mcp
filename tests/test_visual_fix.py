"""Tests for the internal repair engine and the inspect-and-repair loop."""
import os
import sys
import unittest
from unittest.mock import patch
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

import visual_fix
import visual_qa

ENV = {
    "VISION_LLM_ENDPOINT": "https://example.invalid/responses",
    "VISION_LLM_API_KEY": "k",
    "VISION_LLM_MODEL": "m",
}


def make_deck():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    return pres


class TestDescribeSlides(unittest.TestCase):
    def test_describe(self):
        desc = visual_fix.describe_slides(make_deck())
        self.assertEqual(desc[0]["slide"], 1)
        shape = desc[0]["shapes"][0]
        self.assertEqual(shape["shape_index"], 0)
        self.assertEqual(shape["left_in"], 1.0)
        self.assertEqual(shape["text"], "Hello")
        self.assertEqual(shape["font_sizes_pt"], [40.0])

    def test_slide_filter(self):
        self.assertEqual(visual_fix.describe_slides(make_deck(), [2]), [])


class TestApplyRepairs(unittest.TestCase):
    def test_move_resize_font_text(self):
        pres = make_deck()
        result = visual_fix.apply_repairs(pres, [
            {"op": "move_shape", "slide": 1, "shape_index": 0,
             "left_in": 2.0, "top_in": 3.0},
            {"op": "resize_shape", "slide": 1, "shape_index": 0,
             "width_in": 5.0},
            {"op": "set_font_size", "slide": 1, "shape_index": 0,
             "size_pt": 18},
            {"op": "set_text", "slide": 1, "shape_index": 0,
             "text": "Fixed"},
        ])
        self.assertEqual(len(result["applied"]), 4)
        self.assertEqual(result["skipped"], [])
        shape = pres.slides[0].shapes[0]
        self.assertEqual(shape.left, Inches(2.0))
        self.assertEqual(shape.width, Inches(5.0))
        self.assertEqual(shape.text_frame.text, "Fixed")

    def test_delete_shape(self):
        pres = make_deck()
        result = visual_fix.apply_repairs(pres, [
            {"op": "delete_shape", "slide": 1, "shape_index": 0}])
        self.assertEqual(len(result["applied"]), 1)
        self.assertEqual(len(pres.slides[0].shapes), 0)

    def test_invalid_operations_are_skipped(self):
        pres = make_deck()
        result = visual_fix.apply_repairs(pres, [
            {"op": "move_shape", "slide": 9, "shape_index": 0,
             "left_in": 1, "top_in": 1},               # bad slide
            {"op": "move_shape", "slide": 1, "shape_index": 5,
             "left_in": 1, "top_in": 1},               # bad shape
            {"op": "set_font_size", "slide": 1, "shape_index": 0,
             "size_pt": 500},                          # out of range
            {"op": "run_shell", "slide": 1, "shape_index": 0},  # unknown op
            "garbage",                                  # not an object
        ])
        self.assertEqual(result["applied"], [])
        self.assertEqual(len(result["skipped"]), 5)
        # deck untouched
        self.assertEqual(pres.slides[0].shapes[0].left, Inches(1.0))


def make_rich_deck():
    """A deck whose text lives in a table and a chart, not a text box."""
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1),
                                   Inches(4), Inches(2)).table
    table.cell(0, 0).text = "Header"
    table.cell(1, 1).text = "A fairly long cell value"
    data = CategoryChartData()
    data.categories = ["alpha", "beta"]
    data.add_series("s", (1, 2))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1),
                                   Inches(3), Inches(4), Inches(3), data)
    frame.chart.plots[0].has_data_labels = True
    return pres, table, frame.chart


class TestDescribeNonTextShapes(unittest.TestCase):
    """Text in tables and charts must reach the planner; it used to be
    invisible, so no plan could ever address it."""

    def setUp(self):
        self.pres, self.table, self.chart = make_rich_deck()
        self.shapes = visual_fix.describe_slides(self.pres)[0]["shapes"]

    def test_table_structure_is_described(self):
        table = self.shapes[0]["table"]
        self.assertEqual((table["rows"], table["columns"]), (2, 2))
        self.assertEqual(table["column_widths_in"], [2.0, 2.0])
        self.assertEqual(table["row_heights_in"], [1.0, 1.0])
        texts = {(c["row"], c["column"]): c["text"] for c in table["cells"]}
        self.assertEqual(texts[(0, 0)], "Header")
        self.assertEqual(texts[(1, 1)], "A fairly long cell value")

    def test_chart_structure_is_described(self):
        chart = self.shapes[1]["chart"]
        self.assertIn("COLUMN_CLUSTERED", chart["chart_type"])
        self.assertEqual(chart["series_count"], 1)
        self.assertEqual(chart["categories"], ["alpha", "beta"])
        self.assertTrue(chart["has_data_labels"])
        self.assertFalse(chart["has_legend"])


class TestTableAndChartRepairs(unittest.TestCase):
    def setUp(self):
        self.pres, self.table, self.chart = make_rich_deck()

    def apply(self, *ops):
        return visual_fix.apply_repairs(self.pres, list(ops))

    def test_column_width_and_row_height(self):
        result = self.apply(
            {"op": "set_column_width", "slide": 1, "shape_index": 0,
             "column": 0, "width_in": 3.0},
            {"op": "set_row_height", "slide": 1, "shape_index": 0,
             "row": 1, "height_in": 1.5})
        self.assertEqual(len(result["applied"]), 2)
        self.assertEqual(self.table.columns[0].width.inches, 3.0)
        self.assertEqual(self.table.rows[1].height.inches, 1.5)

    def test_cell_text(self):
        self.apply({"op": "set_cell_text", "slide": 1, "shape_index": 0,
                    "row": 0, "column": 1, "text": "Q1"})
        self.assertEqual(self.table.cell(0, 1).text, "Q1")

    def test_set_axis_title_names_axes_by_role(self):
        result = self.apply(
            {"op": "set_axis_title", "slide": 1, "shape_index": 1,
             "axis": "category", "text": "Dinosaur Genus"},
            {"op": "set_axis_title", "slide": 1, "shape_index": 1,
             "axis": "value", "text": "Weight (Metric Tons)"})
        self.assertEqual(len(result["applied"]), 2, result["skipped"])
        self.assertEqual(self.chart.category_axis.axis_title.text_frame.text,
                         "Dinosaur Genus")
        self.assertEqual(self.chart.value_axis.axis_title.text_frame.text,
                         "Weight (Metric Tons)")

    def test_set_axis_title_rejects_a_bad_axis_or_target(self):
        result = self.apply(
            {"op": "set_axis_title", "slide": 1, "shape_index": 1,
             "axis": "horizontal", "text": "x"},
            {"op": "set_axis_title", "slide": 1, "shape_index": 1,
             "axis": "value", "text": 42},
            {"op": "set_axis_title", "slide": 1, "shape_index": 0,
             "axis": "value", "text": "on a table"})
        self.assertEqual(result["applied"], [])
        self.assertEqual([e["reason"] for e in result["skipped"]],
                         ["axis must be category or value", "bad text",
                          "shape is not a chart"])

    def test_describe_reports_existing_axis_titles(self):
        self.chart.category_axis.has_title = True
        self.chart.category_axis.axis_title.text_frame.text = "Genus"
        described = visual_fix.describe_slides(self.pres)[0]["shapes"][1]
        self.assertEqual(described["chart"]["category_axis_title"], "Genus")
        self.assertNotIn("value_axis_title", described["chart"])

    def test_font_size_applies_to_table_and_chart(self):
        result = self.apply(
            {"op": "set_font_size", "slide": 1, "shape_index": 0, "size_pt": 10},
            {"op": "set_font_size", "slide": 1, "shape_index": 1, "size_pt": 9})
        self.assertEqual(len(result["applied"]), 2)
        cell = self.table.cell(0, 0).text_frame.paragraphs[0]
        self.assertEqual(cell.font.size, Pt(10))
        self.assertEqual(self.chart.font.size, Pt(9))

    def test_data_labels_and_legend(self):
        self.apply(
            {"op": "set_chart_data_labels", "slide": 1, "shape_index": 1,
             "show": False},
            {"op": "set_chart_legend", "slide": 1, "shape_index": 1,
             "show": True, "position": "bottom"})
        self.assertFalse(self.chart.plots[0].has_data_labels)
        self.assertTrue(self.chart.has_legend)
        self.assertEqual(str(self.chart.legend.position), "BOTTOM (-4107)")

    def test_wrong_container_is_skipped_not_raised(self):
        result = self.apply(
            # table op aimed at the chart, chart op aimed at the table
            {"op": "set_column_width", "slide": 1, "shape_index": 1,
             "column": 0, "width_in": 2.0},
            {"op": "set_chart_legend", "slide": 1, "shape_index": 0,
             "show": True})
        self.assertEqual(result["applied"], [])
        self.assertEqual([s["reason"] for s in result["skipped"]],
                         ["shape is not a table", "shape is not a chart"])

    def test_out_of_range_table_references_are_skipped(self):
        result = self.apply(
            {"op": "set_row_height", "slide": 1, "shape_index": 0,
             "row": 9, "height_in": 1.0},
            {"op": "set_column_width", "slide": 1, "shape_index": 0,
             "column": 0, "width_in": 999.0},
            {"op": "set_cell_text", "slide": 1, "shape_index": 0,
             "row": 0, "column": 7, "text": "x"})
        self.assertEqual(result["applied"], [])
        self.assertEqual([s["reason"] for s in result["skipped"]],
                         ["bad row index", "width out of range",
                          "bad cell reference"])

    def test_bad_legend_position_is_rejected(self):
        result = self.apply({"op": "set_chart_legend", "slide": 1,
                             "shape_index": 1, "show": True,
                             "position": "middle"})
        self.assertEqual(result["skipped"][0]["reason"], "bad legend position")
        self.assertFalse(self.chart.has_legend)

    def test_show_must_be_boolean(self):
        result = self.apply({"op": "set_chart_data_labels", "slide": 1,
                             "shape_index": 1, "show": "yes"})
        self.assertEqual(result["skipped"][0]["reason"],
                         "show must be true/false")


class TestFitText(unittest.TestCase):
    """Text should fill its box: fit_text shrinks what overflows and grows
    what is starved, with growth anchored to the deck's typography."""

    def setUp(self):
        self.pres = Presentation()
        self.slide = self.pres.slides.add_slide(self.pres.slide_layouts[6])

    def box(self, w, h, text, size_pt=None, wrap=True):
        shape = self.slide.shapes.add_textbox(Inches(1), Inches(1),
                                              Inches(w), Inches(h))
        shape.text_frame.text = text
        shape.text_frame.word_wrap = wrap
        if size_pt is not None:
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(size_pt)
        return shape

    def size_of(self, shape):
        return shape.text_frame.paragraphs[0].font.size.pt

    def fit(self, index, **kwargs):
        op = {"op": "fit_text", "slide": 1, "shape_index": index}
        op.update(kwargs)
        return visual_fix.apply_repairs(self.pres, [op])

    def test_overflowing_text_is_shrunk(self):
        shape = self.box(3, 0.8, "This sentence is far too long for the "
                                 "small box it has been placed inside", 28)
        result = self.fit(0)
        self.assertEqual(result["applied"][0]["op"], "fit_text")
        self.assertLess(self.size_of(shape), 28)

    def test_starved_text_is_grown(self):
        shape = self.box(8, 2, "Key Takeaway", 12)
        self.fit(0)
        self.assertGreater(self.size_of(shape), 12)

    def test_growth_is_capped_relative_to_current_size(self):
        shape = self.box(9, 5, "Big", 12)
        self.fit(0)
        # Without the cap the estimate would reach the 96pt ceiling.
        self.assertLessEqual(self.size_of(shape),
                             12 * visual_fix.MAX_GROWTH_FACTOR)

    def test_unsized_text_uses_the_default_ceiling(self):
        shape = self.box(9, 5, "Big")
        self.fit(0)
        self.assertLessEqual(self.size_of(shape),
                             visual_fix.DEFAULT_GROWTH_CEILING_PT)

    def test_explicit_bounds_are_honoured(self):
        shape = self.box(9, 5, "Big", 12)
        self.fit(0, min_pt=20, max_pt=24)
        self.assertEqual(self.size_of(shape), 24)

    def test_result_reports_the_size_it_chose(self):
        self.box(4, 1, "Q3 Results", 18)
        applied = self.fit(0)["applied"][0]
        self.assertEqual(applied["resolved_size_pt"], self.size_of(
            self.slide.shapes[0]))

    def test_bounds_are_validated(self):
        self.box(4, 1, "Q3 Results", 18)
        self.assertEqual(self.fit(0, min_pt=200)["skipped"][0]["reason"],
                         "min_pt out of range")
        self.assertEqual(self.fit(0, max_pt=0)["skipped"][0]["reason"],
                         "max_pt out of range")

    def test_empty_and_non_text_shapes_are_skipped(self):
        self.box(4, 1, "")
        result = self.fit(0)
        self.assertEqual(result["applied"], [])
        self.assertIn(result["skipped"][0]["reason"],
                      ("shape has no text", "shape has no text frame"))

    def test_unwrapped_text_must_fit_on_one_line(self):
        shape = self.box(2, 3, "A single long unwrapped line", 28, wrap=False)
        self.fit(0)
        self.assertLess(self.size_of(shape), 28)

    def test_estimator_respects_the_box(self):
        wide = self.box(8, 2, "Same text here", 18)
        narrow = self.box(2, 0.5, "Same text here", 18)
        self.assertGreater(visual_fix.estimate_fit_font_size(wide),
                           visual_fix.estimate_fit_font_size(narrow))


class TestAutofit(unittest.TestCase):
    def setUp(self):
        self.pres = Presentation()
        slide = self.pres.slides.add_slide(self.pres.slide_layouts[6])
        self.shape = slide.shapes.add_textbox(Inches(1), Inches(1),
                                              Inches(4), Inches(1))
        self.shape.text_frame.text = "Hello"

    def apply(self, mode):
        return visual_fix.apply_repairs(self.pres, [
            {"op": "set_autofit", "slide": 1, "shape_index": 0, "mode": mode}])

    def test_modes(self):
        for mode, expected in (("shrink_text", "TEXT_TO_FIT_SHAPE"),
                               ("grow_shape", "SHAPE_TO_FIT_TEXT"),
                               ("none", "NONE")):
            self.apply(mode)
            self.assertIn(expected, str(self.shape.text_frame.auto_size))

    def test_unknown_mode_is_rejected(self):
        result = self.apply("sideways")
        self.assertEqual(result["skipped"][0]["reason"], "bad autofit mode")


class TestReviewChecklist(unittest.TestCase):
    """The reviewer must be told to look inside charts, tables and diagrams,
    not only at text boxes."""

    def test_prompt_covers_non_textbox_text(self):
        prompt = visual_qa.review_prompt(False)
        for phrase in ("axis tick labels", "data labels", "legend",
                       "Tables:", "columns too narrow", "SmartArt",
                       "grouped shapes", "overlapping"):
            self.assertIn(phrase, prompt)

    def test_prompt_covers_text_sized_badly_for_its_box(self):
        prompt = visual_qa.review_prompt(False)
        self.assertIn("mostly empty", prompt)
        self.assertIn("without crowding", prompt)

    def test_repair_prompt_advertises_the_new_operations(self):
        for op in ("set_column_width", "set_row_height", "set_cell_text",
                   "set_chart_legend", "set_chart_data_labels",
                   "fit_text", "set_autofit"):
            self.assertIn(op, visual_fix.REPAIR_PROMPT)


class TestInspectAndRepairLoop(unittest.TestCase):
    def setUp(self):
        self._saved = {k: os.environ.get(k) for k in ENV}
        os.environ.update(ENV)

    def tearDown(self):
        for k, v in self._saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v

    def _run(self, verdicts, plans):
        """Drive the loop with scripted inspection verdicts and repair plans;
        rendering and HTTP are stubbed out."""
        calls = {"review": 0, "plan": 0}

        def fake_review(self_llm, images, prompt, timeout=None):
            v = verdicts[min(calls["review"], len(verdicts) - 1)]
            calls["review"] += 1
            return dict(v)

        def fake_plan(llm, issues, pres, images, image_slides=None):
            p = plans[min(calls["plan"], len(plans) - 1)]
            calls["plan"] += 1
            return p

        with patch.object(visual_qa, "_render_deck",
                          return_value=[b"\x89PNG-fake"]), \
             patch.object(visual_qa.VisionLLM, "review", fake_review), \
             patch.object(visual_fix, "plan_repairs", fake_plan):
            outcome = visual_qa.inspect_and_repair(make_deck())
        return outcome, calls

    def test_zero_applied_round_reports_why_and_says_not_to_retry(self):
        """The 0-applied/N-skipped result an operator sees has to name the
        reason: repairs ran, the plan was rejected."""
        verdicts = [{"passed": False,
                     "issues": [{"slide": 1, "description": "layout thing"}]}]
        # shape_index 99 exists on no slide: the validator rejects both.
        plans = [[{"op": "move_shape", "slide": 1, "shape_index": 99,
                   "left_in": 1.0, "top_in": 1.0},
                  {"op": "resize_shape", "slide": 1, "shape_index": 99,
                   "width_in": 2.0}]]
        outcome, calls = self._run(verdicts, plans)
        self.assertFalse(outcome["passed"])
        round_one = outcome["repair_rounds"][0]
        self.assertEqual(round_one["operations_applied"], 0)
        self.assertEqual(round_one["skipped_reasons"], {"bad shape_index": 2})
        self.assertIn("bad shape_index x2", outcome["repair_note"])
        self.assertIn("will not help", outcome["repair_note"])
        # Stopped instead of burning the rest of the budget on a plan that
        # cannot land.
        self.assertEqual(calls["review"], 1)

    def test_successful_round_carries_no_skip_noise(self):
        verdicts = [
            {"passed": False, "issues": [{"slide": 1, "description": "x"}]},
            {"passed": True, "issues": []},
        ]
        plans = [[{"op": "move_shape", "slide": 1, "shape_index": 0,
                   "left_in": 1.0, "top_in": 1.0}]]
        outcome, _ = self._run(verdicts, plans)
        self.assertTrue(outcome["passed"])
        self.assertNotIn("skipped_reasons", outcome["repair_rounds"][0])
        self.assertNotIn("repair_note", outcome)

    def test_pass_first_time(self):
        outcome, calls = self._run([{"passed": True, "issues": []}], [[]])
        self.assertTrue(outcome["passed"])
        self.assertEqual(outcome["iterations"], 1)
        self.assertEqual(calls["plan"], 0)

    def test_repair_then_pass(self):
        verdicts = [
            {"passed": False, "issues": [{"slide": 1, "description": "x"}]},
            {"passed": True, "issues": []},
        ]
        plan = [{"op": "set_font_size", "slide": 1, "shape_index": 0,
                 "size_pt": 20}]
        outcome, calls = self._run(verdicts, [plan])
        self.assertTrue(outcome["passed"])
        self.assertEqual(outcome["iterations"], 2)
        self.assertEqual(outcome["repair_rounds"][0]["operations_applied"], 1)

    def test_gives_up_after_max_iterations(self):
        failing = {"passed": False,
                   "issues": [{"slide": 1, "description": "x"}]}
        plan = [{"op": "set_font_size", "slide": 1, "shape_index": 0,
                 "size_pt": 20}]
        os.environ["VISUAL_QA_MAX_ITERATIONS"] = "3"
        try:
            outcome, calls = self._run([failing], [plan])
        finally:
            os.environ.pop("VISUAL_QA_MAX_ITERATIONS")
        self.assertFalse(outcome["passed"])
        self.assertEqual(calls["review"], 3)
        self.assertEqual(len(outcome["repair_rounds"]), 2)
        self.assertTrue(outcome["issues"])

    def test_stops_early_when_no_repairs_apply(self):
        failing = {"passed": False,
                   "issues": [{"slide": 1, "description": "x"}]}
        outcome, calls = self._run([failing], [[]])  # planner returns nothing
        self.assertFalse(outcome["passed"])
        self.assertEqual(calls["review"], 1)  # no progress -> stop immediately


class TestSkipReasonSummary(unittest.TestCase):
    """A round that applies nothing has to explain itself: the counts are
    what tell an operator the plan was rejected, not that repair is off."""

    def test_counts_reasons_and_drops_the_detail(self):
        summary = visual_fix.skip_reason_summary([
            {"op": {}, "reason": "bad shape_index"},
            {"op": {}, "reason": "bad shape_index"},
            {"op": {}, "reason": "apply failed: something specific"},
        ])
        self.assertEqual(summary, {"bad shape_index": 2, "apply failed": 1})

    def test_empty_plan_summarizes_to_nothing(self):
        self.assertEqual(visual_fix.skip_reason_summary([]), {})


if __name__ == "__main__":
    unittest.main()
