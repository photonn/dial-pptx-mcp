"""
Tests for combo charts and per-series chart formatting.

The interesting assertions are on the reopened deck, not the in-memory object:
the whole technique is an XML rearrangement, and a chart that python-pptx is
happy to build in memory but cannot parse back is exactly the failure mode
worth guarding.
"""
import io
import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

import deck_validation
import utils as ppt_utils
from state import PresentationStore
from tools.chart_tools import register_chart_tools

CATEGORIES = ["Q1", "Q2", "Q3", "Q4"]


def validate_parameters(params):
    for name, (value, constraints) in params.items():
        for check, message in constraints:
            if not check(value):
                return False, f"{name}: {message}"
    return True, None


def is_valid_rgb(value):
    return (isinstance(value, list) and len(value) == 3
            and all(isinstance(c, int) and 0 <= c <= 255 for c in value))


def build_tools():
    app = FastMCP(name="test")
    store = PresentationStore(ttl_seconds=60, max_items=10)
    register_chart_tools(app, store, lambda: None, validate_parameters,
                         lambda v: v > 0, lambda v: v >= 0,
                         lambda lo, hi: (lambda v: lo <= v <= hi),
                         is_valid_rgb)
    return app._tool_manager._tools, store


def deck_with_slide():
    pres = Presentation()
    pres.slides.add_slide(pres.slide_layouts[6])
    return pres


def roundtrip(pres):
    buffer = io.BytesIO()
    pres.save(buffer)
    return Presentation(io.BytesIO(buffer.getvalue()))


def chart_of(pres, slide=0):
    return next(s for s in pres.slides[slide].shapes if s.has_chart).chart


def plain_chart(pres, series_count=3):
    data = CategoryChartData()
    data.categories = CATEGORIES
    for i in range(series_count):
        data.add_series(f"S{i}", tuple(range(1, len(CATEGORIES) + 1)))
    return pres.slides[0].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6),
        Inches(4), data).chart


class TestComboLayout(unittest.TestCase):
    def test_series_are_split_into_groups_by_type(self):
        pres = deck_with_slide()
        chart = plain_chart(pres)
        layout = ppt_utils.apply_combo_layout(
            chart, [{"type": "column"}, {"type": "column"},
                    {"type": "line"}])
        self.assertEqual(layout["groups"], 2)
        self.assertIsNone(layout["secondary"])
        plots = chart_of(roundtrip(pres)).plots
        self.assertEqual([[s.name for s in plot.series] for plot in plots],
                         [["S0", "S1"], ["S2"]])

    def test_a_secondary_axis_adds_a_second_axis_pair(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.apply_combo_layout(
            chart, [{"type": "column"},
                    {"type": "line", "secondary_axis": True}])
        reopened = chart_of(roundtrip(pres))
        plot_area = reopened._chartSpace.find(qn("c:chart")).find(
            qn("c:plotArea"))
        self.assertEqual(len(plot_area.findall(qn("c:valAx"))), 2)
        self.assertEqual(len(plot_area.findall(qn("c:catAx"))), 2)

    def test_the_secondary_category_axis_is_hidden(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.apply_combo_layout(
            chart, [{"type": "column"},
                    {"type": "line", "secondary_axis": True}])
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        deletes = [axis.find(qn("c:delete")).get("val")
                   for axis in plot_area.findall(qn("c:catAx"))]
        self.assertEqual(sorted(deletes), ["0", "1"])

    def test_each_group_names_its_own_axis_pair(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.apply_combo_layout(
            chart, [{"type": "column"},
                    {"type": "line", "secondary_axis": True}])
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        pairs = []
        for group in plot_area:
            if group.tag.endswith("Chart"):
                pairs.append(tuple(e.get("val")
                                   for e in group.findall(qn("c:axId"))))
        self.assertEqual(len(pairs), 2)
        self.assertNotEqual(pairs[0], pairs[1])
        declared = {axis.find(qn("c:axId")).get("val")
                    for axis in list(plot_area.findall(qn("c:valAx")))
                    + list(plot_area.findall(qn("c:catAx")))}
        for pair in pairs:
            self.assertTrue(set(pair) <= declared,
                            "a group references an axis that is not declared")

    def test_series_of_one_type_stay_in_one_group(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 3)
        layout = ppt_utils.apply_combo_layout(
            chart, [{"type": "column"}] * 3)
        self.assertEqual(layout["groups"], 1)

    def test_a_mismatched_spec_count_is_refused(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 3)
        with self.assertRaises(ValueError):
            ppt_utils.apply_combo_layout(chart, [{"type": "column"}])

    def test_the_result_is_structurally_valid(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 3)
        ppt_utils.apply_combo_layout(
            chart, [{"type": "column"}, {"type": "column"},
                    {"type": "line", "secondary_axis": True}])
        report = deck_validation.validate_presentation(pres)
        self.assertTrue(report["ok"], report["problems"])


class TestSeriesFormatting(unittest.TestCase):
    def test_a_series_colour_survives_a_roundtrip(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.set_series_color(chart.plots[0].series[0], (31, 73, 125))
        series = chart_of(roundtrip(pres)).plots[0].series[0]
        self.assertEqual(str(series.format.fill.fore_color.rgb), "1F497D")

    def test_data_labels_can_be_turned_on_for_one_series(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.set_series_data_labels(chart.plots[0].series[0], True,
                                         "outside_end", "0.0%")
        reopened = chart_of(roundtrip(pres)).plots[0]
        self.assertTrue(reopened.series[0].data_labels.show_value)
        self.assertEqual(reopened.series[0].data_labels.number_format, "0.0%")

    def test_a_trendline_is_added_to_the_series(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.add_series_trendline(chart.plots[0].series[0], "linear")
        series = chart_of(roundtrip(pres)).plots[0].series[0]._element
        trendlines = series.findall(qn("c:trendline"))
        self.assertEqual(len(trendlines), 1)
        self.assertEqual(
            trendlines[0].find(qn("c:trendlineType")).get("val"), "linear")

    def test_a_moving_average_carries_its_period(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        ppt_utils.add_series_trendline(chart.plots[0].series[0], "movingAvg", 3)
        series = chart_of(roundtrip(pres)).plots[0].series[0]._element
        self.assertEqual(
            series.find(qn("c:trendline")).find(qn("c:period")).get("val"), "3")

    def test_an_unknown_trendline_is_refused(self):
        pres = deck_with_slide()
        chart = plain_chart(pres, 2)
        with self.assertRaises(ValueError):
            ppt_utils.add_series_trendline(chart.plots[0].series[0], "magic")


class TestComboChartTool(unittest.TestCase):
    def setUp(self):
        self.tools, self.store = build_tools()
        self.pres = deck_with_slide()
        self.pid = self.store.new_id()
        self.store[self.pid] = self.pres
        self.add = self.tools["add_combo_chart"].fn
        self.format = self.tools["format_chart_series"].fn

    def series(self, **overrides):
        base = [
            {"name": "Revenue", "values": [10, 12, 15, 18], "type": "column"},
            {"name": "Margin", "values": [0.2, 0.23, 0.22, 0.25],
             "type": "line_markers", "secondary_axis": True},
        ]
        base[0].update(overrides)
        return base

    def add_default(self, **kwargs):
        return self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES, self.series(),
                        presentation_id=self.pid, **kwargs)

    def test_it_builds_a_two_group_chart(self):
        result = self.add_default()
        self.assertNotIn("error", result)
        self.assertEqual(result["plot_groups"], 2)
        self.assertTrue(result["has_secondary_axis"])
        plots = chart_of(roundtrip(self.pres)).plots
        self.assertEqual([[s.name for s in plot.series] for plot in plots],
                         [["Revenue"], ["Margin"]])

    def test_titles_are_applied_including_the_secondary_axis(self):
        self.add_default(title="Growth", y_axis_title="Revenue",
                         secondary_axis_title="Margin",
                         x_axis_title="Quarter")
        chart = chart_of(roundtrip(self.pres))
        self.assertEqual(chart.chart_title.text_frame.text, "Growth")
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        titles = [axis.find(qn("c:title")) is not None
                  for axis in plot_area.findall(qn("c:valAx"))]
        self.assertEqual(titles, [True, True])

    def test_the_left_axis_title_goes_on_the_left_axis(self):
        """python-pptx's chart.value_axis returns the *second* value axis once
        a chart has two, so a title set through it lands on the secondary
        axis."""
        self.add_default(y_axis_title="Revenue")
        chart = chart_of(roundtrip(self.pres))
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        titled = {}
        for axis in plot_area.findall(qn("c:valAx")):
            position = axis.find(qn("c:axPos")).get("val")
            titled[position] = axis.find(qn("c:title")) is not None
        self.assertEqual(titled, {"l": True, "r": False})

    def test_per_series_options_are_applied(self):
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES,
                          self.series(color=[31, 73, 125], data_labels=True,
                                      trendline="linear"),
                          presentation_id=self.pid)
        self.assertNotIn("error", result)
        series = chart_of(roundtrip(self.pres)).plots[0].series[0]
        self.assertEqual(str(series.format.fill.fore_color.rgb), "1F497D")
        self.assertEqual(len(series._element.findall(qn("c:trendline"))), 1)

    def test_the_chart_is_structurally_valid(self):
        self.add_default(title="Growth", secondary_axis_title="Margin")
        report = deck_validation.validate_presentation(self.pres)
        self.assertTrue(report["ok"], report["problems"])

    def test_a_single_series_is_sent_to_add_chart(self):
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES,
                          [{"name": "x", "values": [1, 2, 3, 4],
                            "type": "column"}], presentation_id=self.pid)
        self.assertIn("add_chart", result["error"])

    def test_an_unsupported_type_names_the_alternatives(self):
        series = self.series()
        series[0]["type"] = "pie"
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES, series,
                          presentation_id=self.pid)
        self.assertIn("cannot share a plot area", result["error"])
        self.assertIn("add_chart", result["error"])

    def test_a_short_series_is_refused_before_anything_is_built(self):
        series = self.series()
        series[0]["values"] = [1, 2]
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES, series,
                          presentation_id=self.pid)
        self.assertIn("value(s) for 4 categories", result["error"])
        self.assertEqual(len(self.pres.slides[0].shapes), 0)

    def test_a_missing_key_is_named(self):
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES,
                          [{"name": "a", "values": [1, 2, 3, 4]},
                           {"name": "b", "values": [1, 2, 3, 4],
                            "type": "line"}], presentation_id=self.pid)
        self.assertIn("'type'", result["error"])

    def test_an_unknown_presentation_is_refused(self):
        result = self.add(0, 0.5, 0.5, 8, 4.5, CATEGORIES, self.series(),
                          presentation_id="nope")
        self.assertIn("expired presentation_id", result["error"])


class TestFormatSeriesTool(TestComboChartTool):
    def test_it_restyles_an_existing_series(self):
        self.add_default()
        result = self.format(0, 0, 1, color=[200, 30, 30],
                             presentation_id=self.pid)
        self.assertEqual(result["applied"], ["color"])
        series = chart_of(roundtrip(self.pres)).plots[1].series[0]
        self.assertEqual(str(series.format.fill.fore_color.rgb), "C81E1E")

    def test_series_index_spans_every_group(self):
        self.add_default()
        self.assertIn("error", self.format(0, 0, 2, color=[1, 2, 3],
                                           presentation_id=self.pid))

    def test_a_shape_that_is_not_a_chart_is_refused(self):
        self.pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1))
        result = self.format(0, 0, 0, color=[1, 2, 3],
                             presentation_id=self.pid)
        self.assertIn("not a chart", result["error"])

    def test_doing_nothing_is_reported_rather_than_claimed(self):
        self.add_default()
        result = self.format(0, 0, 0, presentation_id=self.pid)
        self.assertIn("at least one", result["error"])

    def test_a_bad_colour_is_refused(self):
        self.add_default()
        result = self.format(0, 0, 0, color=[300, 0, 0],
                             presentation_id=self.pid)
        self.assertIn("0-255", result["error"])


if __name__ == "__main__":
    unittest.main()
