"""
Tests for the preview composers, the summary-card tool and the design-guidance
document.

Rendering needs LibreOffice, so the tests that would drive it self-skip the way
the rest of the suite does; the composer and the guidance parser are pure and
are tested directly.
"""
import io
import sys
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import previews
import visual_qa
from tools.guidance_tools import GUIDANCE_PATH, _load, _slug


def tile(width=320, height=180, colour=(120, 140, 160)):
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), colour).save(buffer, "PNG")
    return buffer.getvalue()


def open_sheet(data):
    return Image.open(io.BytesIO(data))


class FakeDialClient:
    """Stands in for DialFileClient: returns a file URL, uploads nothing."""

    def __init__(self, *args, **kwargs):
        self.uploaded = []

    def upload(self, data, filename, folder=None, content_type=None):
        self.uploaded.append((filename, len(data), content_type))
        return f"files/test-bucket/{filename}"


class TestContactSheet(unittest.TestCase):
    def test_a_single_slide_makes_a_one_cell_sheet(self):
        sheet = open_sheet(previews.compose_contact_sheet([tile()], [1]))
        self.assertGreater(sheet.width, 320)
        self.assertGreater(sheet.height, 180)

    def test_the_grid_wraps_at_the_column_count(self):
        tiles = [tile() for _ in range(4)]
        two_across = open_sheet(
            previews.compose_contact_sheet(tiles, [1, 2, 3, 4], columns=2))
        four_across = open_sheet(
            previews.compose_contact_sheet(tiles, [1, 2, 3, 4], columns=4))
        self.assertGreater(four_across.width, two_across.width)
        self.assertGreater(two_across.height, four_across.height)

    def test_columns_never_exceed_the_number_of_slides(self):
        sheet = open_sheet(
            previews.compose_contact_sheet([tile()], [1], columns=6))
        # One tile plus padding, not six cells' worth of empty grid.
        self.assertLess(sheet.width, 320 * 2)

    def test_slides_of_different_sizes_are_not_stretched(self):
        tiles = [tile(320, 180), tile(160, 90)]
        sheet = previews.compose_contact_sheet(tiles, [1, 2], columns=2)
        self.assertTrue(open_sheet(sheet).width > 320)

    def test_an_empty_selection_is_rejected(self):
        with self.assertRaises(ValueError):
            previews.compose_contact_sheet([], [])

    def test_the_sheet_is_jpeg(self):
        sheet = previews.compose_contact_sheet([tile()], [1])
        self.assertEqual(open_sheet(sheet).format, "JPEG")

    def test_the_describe_prompt_maps_images_to_absolute_slide_numbers(self):
        prompt = previews.DESCRIBE_PROMPT.format(
            mapping="image 1 = slide 7, image 2 = slide 9")
        self.assertIn("image 1 = slide 7", prompt)
        self.assertIn("best_for", prompt)


class TestSummaryCard(unittest.TestCase):
    """One image of the whole deck, whatever the deck's length."""

    def card(self, count, **kwargs):
        data, size, columns = previews.compose_summary_card(
            [tile() for _ in range(count)], list(range(1, count + 1)),
            **kwargs)
        return open_sheet(data), size, columns

    def test_a_long_deck_still_makes_exactly_one_image(self):
        image, size, _ = self.card(60)
        self.assertEqual(image.size, size)

    def test_the_card_stays_within_its_width_budget(self):
        for count in (1, 5, 34, 150):
            with self.subTest(slides=count):
                _, (width, _), _ = self.card(count)
                self.assertLessEqual(width, previews.CARD_MAX_WIDTH)

    def test_the_grid_stays_roughly_landscape_as_the_deck_grows(self):
        """A fixed column count would turn a 60-slide deck into a stripe."""
        for count in (12, 34, 60, 150):
            with self.subTest(slides=count):
                _, (width, height), _ = self.card(count)
                self.assertGreater(width / height, 0.6)
                self.assertLess(width / height, 2.5)

    def test_cells_stay_legible_even_if_the_card_must_grow_taller(self):
        """The height cap yields to CARD_MIN_CELL_WIDTH: a card that scrolls
        beats one whose thumbnails are unreadable."""
        _, (width, _), columns = self.card(150)
        self.assertGreaterEqual(width / columns, previews.CARD_MIN_CELL_WIDTH)

    def test_an_explicit_column_count_is_honoured(self):
        _, _, columns = self.card(12, columns=3)
        self.assertEqual(columns, 3)

    def test_columns_never_exceed_the_number_of_slides(self):
        _, _, columns = self.card(2, columns=6)
        self.assertEqual(columns, 2)

    def test_the_title_adds_a_header_band(self):
        _, (_, plain), _ = self.card(4)
        _, (_, titled), _ = self.card(4, title="Q3 Review")
        self.assertEqual(titled - plain, previews._HEADER_HEIGHT)

    def test_an_empty_deck_is_rejected(self):
        with self.assertRaises(ValueError):
            previews.compose_summary_card([], [])

    def test_the_card_is_jpeg(self):
        image, _, _ = self.card(3)
        self.assertEqual(image.format, "JPEG")

    def test_auto_columns_degenerates_safely(self):
        self.assertEqual(previews.auto_columns(0), 1)
        self.assertEqual(previews.auto_columns(1), 1)


class TestSummaryCardTool(unittest.TestCase):
    """The tool layer: guards, and the shape the agent is handed."""

    def setUp(self):
        from mcp.server.fastmcp import FastMCP
        from pptx import Presentation

        from state import PresentationStore
        from tools.preview_tools import register_preview_tools

        app = FastMCP(name="test")
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        with patch.object(visual_qa, "_soffice_binary", return_value="soffice"):
            register_preview_tools(app, self.store)
        self.tool = app._tool_manager._tools["render_deck_summary_card"]
        self.card = self.tool.fn

        pres = Presentation()
        for _ in range(3):
            pres.slides.add_slide(pres.slide_layouts[6])
        self.pid = self.store.new_id()
        self.store[self.pid] = pres

    def run_card(self, **kwargs):
        with patch("dial_client.DialFileClient", FakeDialClient), \
             patch.object(visual_qa, "_render_deck",
                          return_value=[tile() for _ in range(3)]):
            return self.card(presentation_id=self.pid, **kwargs)

    def test_it_is_read_only(self):
        """Otherwise summarising a finished deck marks it visually stale."""
        self.assertTrue(self.tool.annotations.readOnlyHint)

    def test_it_returns_one_image_url_for_the_whole_deck(self):
        result = self.run_card(title="Q3 Review")
        self.assertEqual(result["mime_type"], previews.JPEG_MIME)
        self.assertEqual(result["slide_count"], 3)
        self.assertTrue(result["image_url"].startswith("files/"))
        self.assertGreater(result["size_bytes"], 0)
        self.assertIn("attachment", result["message"])

    def test_an_unknown_handle_is_refused(self):
        self.assertIn("error", self.card(presentation_id="nope"))

    def test_an_empty_deck_is_refused(self):
        from pptx import Presentation

        pid = self.store.new_id()
        self.store[pid] = Presentation()
        self.assertIn("error", self.card(presentation_id=pid))

    def test_a_nonsense_column_count_is_refused(self):
        self.assertIn("error", self.run_card(columns=0))
        self.assertIn("error", self.run_card(columns=99))

    def test_a_failed_upload_is_an_error_not_a_half_answer(self):
        """The stored image is the whole deliverable — there is nothing
        useful to return without it."""
        class Broken(FakeDialClient):
            def upload(self, *args, **kwargs):
                raise RuntimeError("no bucket")

        with patch("dial_client.DialFileClient", Broken), \
             patch.object(visual_qa, "_render_deck",
                          return_value=[tile() for _ in range(3)]):
            result = self.card(presentation_id=self.pid)
        self.assertIn("no bucket", result["error"])
        self.assertNotIn("image_url", result)


class TestDesignGuidance(unittest.TestCase):
    def test_the_document_ships_with_the_server(self):
        self.assertTrue(GUIDANCE_PATH.exists(), GUIDANCE_PATH)

    def test_it_parses_into_sections(self):
        text, sections = _load()
        self.assertIsNotNone(text)
        self.assertGreaterEqual(len(sections), 8)
        for slug, meta in sections.items():
            with self.subTest(section=slug):
                self.assertTrue(meta["body"].startswith("## "))
                self.assertGreater(len(meta["body"]), 100)

    def test_sections_are_numbered_in_document_order(self):
        _, sections = _load()
        numbers = [meta["number"] for meta in
                   sorted(sections.values(), key=lambda m: m["number"])]
        self.assertEqual(numbers, list(range(1, len(numbers) + 1)))

    def test_it_covers_the_topics_the_tool_promises(self):
        text, _ = _load()
        for topic in ("duplicate_slide", "visual_inspect_slides",
                      "validate_presentation", "manage_speaker_notes",
                      "Liberation Sans", "Aptos", "Margins"):
            with self.subTest(topic=topic):
                self.assertIn(topic, text)

    def test_slugs_are_stable(self):
        self.assertEqual(_slug("Charts and tables"), "charts-and-tables")
        self.assertEqual(_slug("Type"), "type")


if __name__ == "__main__":
    unittest.main()
