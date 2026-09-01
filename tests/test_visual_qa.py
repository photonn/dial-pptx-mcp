"""Tests for the visual-inspection capability (rendering + vision LLM client)."""
import base64
import os
import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from visual_qa import VisionLLM, VisionLLMConfigError, review_prompt

ENV = {
    "VISION_LLM_ENDPOINT": "https://example.invalid/openai/responses?api-version=x",
    "VISION_LLM_API_KEY": "k",
    "VISION_LLM_MODEL": "vision-model",
}


def _slide_has_text(slide, needle):
    return any(shape.has_text_frame and needle in shape.text_frame.text
              for shape in slide.shapes)


class WithEnv(unittest.TestCase):
    def setUp(self):
        self._saved = {k: os.environ.get(k) for k in ENV}
        os.environ.update(ENV)

    def tearDown(self):
        for k, v in self._saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v


class TestConfig(unittest.TestCase):
    def test_unconfigured_raises(self):
        saved = {k: os.environ.pop(k, None) for k in ENV}
        try:
            with self.assertRaises(VisionLLMConfigError):
                VisionLLM()
        finally:
            for k, v in saved.items():
                if v is not None:
                    os.environ[k] = v


class TestPayload(WithEnv):
    def test_payload_shape(self):
        llm = VisionLLM()
        png = b"\x89PNG-fake"
        payload = llm.build_payload([png, png], "check this")
        self.assertEqual(payload["model"], "vision-model")
        content = payload["input"][0]["content"]
        self.assertEqual(content[0], {"type": "input_text", "text": "check this"})
        self.assertEqual(len(content), 3)
        expected = "data:image/png;base64," + base64.b64encode(png).decode()
        self.assertEqual(content[1]["image_url"], expected)
        self.assertEqual(content[1]["type"], "input_image")


class TestResponseParsing(WithEnv):
    def test_extract_output_text_field(self):
        self.assertEqual(VisionLLM().extract_text({"output_text": "hi"}), "hi")

    def test_extract_from_output_list(self):
        resp = {"output": [
            {"type": "reasoning", "content": []},
            {"type": "message", "content": [
                {"type": "output_text", "text": '{"passed": true, "issues": []}'}]},
        ]}
        self.assertIn("passed", VisionLLM().extract_text(resp))

    def test_parse_verdict_plain(self):
        v = VisionLLM.parse_verdict('{"passed": false, "issues": [{"slide": 2}]}')
        self.assertFalse(v["passed"])
        self.assertEqual(v["issues"][0]["slide"], 2)

    def test_parse_verdict_fenced(self):
        v = VisionLLM.parse_verdict('```json\n{"passed": true, "issues": []}\n```')
        self.assertTrue(v["passed"])

    def test_parse_verdict_garbage_keeps_raw(self):
        v = VisionLLM.parse_verdict("The deck looks fine to me!")
        self.assertIsNone(v["passed"])
        self.assertIn("raw_review", v)


class TestDialProvider(unittest.TestCase):
    """VISION_LLM served through DIAL Core as a deployment."""

    def setUp(self):
        self._saved = {k: os.environ.get(k) for k in (
            *ENV, "DIAL_CORE_URL", "DIAL_API_KEY", "VISION_LLM_PROVIDER")}
        for k in self._saved:
            os.environ.pop(k, None)
        os.environ["VISION_LLM_MODEL"] = "vision-deployment"
        os.environ["DIAL_CORE_URL"] = "https://dial.example.invalid"
        os.environ["DIAL_API_KEY"] = "dial-key"

    def tearDown(self):
        for k, v in self._saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v

    def test_provider_autodetect(self):
        llm = VisionLLM()  # no VISION_LLM_ENDPOINT -> dial
        self.assertEqual(llm.provider, "dial")
        os.environ["VISION_LLM_ENDPOINT"] = "https://x.invalid/responses"
        os.environ["VISION_LLM_API_KEY"] = "k"
        self.assertEqual(VisionLLM().provider, "direct")
        os.environ["VISION_LLM_PROVIDER"] = "dial"
        self.assertEqual(VisionLLM().provider, "dial")

    def test_dial_payload_and_target(self):
        llm = VisionLLM()
        payload = llm.build_payload([b"\x89PNG-x"], "check")
        content = payload["messages"][0]["content"]
        self.assertEqual(content[0], {"type": "text", "text": "check"})
        self.assertEqual(content[1]["type"], "image_url")
        self.assertTrue(content[1]["image_url"]["url"].startswith("data:image/png;base64,"))
        self.assertNotIn("model", payload)  # deployment is in the URL
        url, headers = llm._request_target()
        # api-version is always sent: Azure (and DIAL's Azure upstream)
        # reject requests without it.
        from visual_qa import DEFAULT_API_VERSION
        self.assertEqual(url, "https://dial.example.invalid/openai/"
                              "deployments/vision-deployment/chat/completions"
                              f"?api-version={DEFAULT_API_VERSION}")
        self.assertEqual(headers["Api-Key"], "dial-key")

    def test_dial_api_version(self):
        os.environ["VISION_LLM_API_VERSION"] = "2025-01-01-preview"
        try:
            url, _ = VisionLLM()._request_target()
            self.assertTrue(url.endswith("?api-version=2025-01-01-preview"))
        finally:
            os.environ.pop("VISION_LLM_API_VERSION")

    def test_direct_endpoint_gets_api_version_when_missing(self):
        from visual_qa import DEFAULT_API_VERSION
        os.environ["VISION_LLM_PROVIDER"] = "direct"
        os.environ["VISION_LLM_ENDPOINT"] = "https://x.invalid/openai/responses"
        os.environ["VISION_LLM_API_KEY"] = "k"
        url, _ = VisionLLM()._request_target()
        self.assertEqual(url, "https://x.invalid/openai/responses"
                              f"?api-version={DEFAULT_API_VERSION}")

    def test_configured_api_version_in_url_wins(self):
        os.environ["VISION_LLM_PROVIDER"] = "direct"
        os.environ["VISION_LLM_ENDPOINT"] = (
            "https://x.invalid/openai/responses?api-version=2025-04-01-preview")
        os.environ["VISION_LLM_API_KEY"] = "k"
        os.environ["VISION_LLM_API_VERSION"] = "2020-01-01"
        try:
            url, _ = VisionLLM()._request_target()
            self.assertEqual(
                url,
                "https://x.invalid/openai/responses?api-version=2025-04-01-preview")
        finally:
            os.environ.pop("VISION_LLM_API_VERSION")

    def test_dial_extract_text(self):
        llm = VisionLLM()
        resp = {"choices": [{"message": {"content": '{"passed": true}'}}]}
        self.assertEqual(llm.extract_text(resp), '{"passed": true}')
        parts = {"choices": [{"message": {"content": [
            {"type": "text", "text": "a"}, {"type": "text", "text": "b"}]}}]}
        self.assertEqual(llm.extract_text(parts), "a\nb")

    def test_dial_requires_core_url(self):
        os.environ.pop("DIAL_CORE_URL")
        from visual_qa import VisionLLMConfigError
        with self.assertRaises(VisionLLMConfigError):
            VisionLLM()

    def test_enforcement_via_dial(self):
        import visual_qa
        self.assertTrue(visual_qa.enforcement_enabled())
        os.environ.pop("DIAL_CORE_URL")
        self.assertFalse(visual_qa.enforcement_enabled())


class TestPrompt(unittest.TestCase):
    def test_reference_and_focus_wording(self):
        p = review_prompt(True, focus="check logo size")
        self.assertIn("reference template", p)
        self.assertIn("check logo size", p)
        self.assertIn('"passed"', p)
        p2 = review_prompt(False)
        self.assertNotIn("FIRST images", p2)


class TestRendering(unittest.TestCase):
    def test_render_demo_deck(self):
        import shutil
        if not (os.environ.get("SOFFICE_PATH") or shutil.which("soffice")):
            self.skipTest("LibreOffice not installed")
        from visual_qa import render_pptx_bytes_to_pngs
        data = (Path(__file__).resolve().parent.parent /
                "mcp_all_tools_templates_effects_demo.pptx").read_bytes()
        images = render_pptx_bytes_to_pngs(data, max_slides=2)
        self.assertEqual(len(images), 2)
        for png in images:
            self.assertTrue(png.startswith(b"\x89PNG"))

    def test_render_selected_slides_only(self):
        import shutil
        if not (os.environ.get("SOFFICE_PATH") or shutil.which("soffice")):
            self.skipTest("LibreOffice not installed")
        from visual_qa import render_pptx_bytes_to_pngs
        data = (Path(__file__).resolve().parent.parent /
                "mcp_all_tools_templates_effects_demo.pptx").read_bytes()
        both = render_pptx_bytes_to_pngs(data, max_slides=2)
        only_second = render_pptx_bytes_to_pngs(data, slides=[2])
        self.assertEqual(len(only_second), 1)
        self.assertTrue(only_second[0].startswith(b"\x89PNG"))
        # It is slide 2 that came back, not the first page again.
        self.assertNotEqual(only_second[0], both[0])
        # Out-of-range selections are dropped, not fatal.
        self.assertEqual(render_pptx_bytes_to_pngs(data, slides=[9999]), [])

    def test_render_deck_with_slides_converts_a_stripped_copy(self):
        import shutil
        from pptx import Presentation
        if not (os.environ.get("SOFFICE_PATH") or shutil.which("soffice")):
            self.skipTest("LibreOffice not installed")
        import visual_qa
        pres = Presentation(str(Path(__file__).resolve().parent.parent /
                                "mcp_all_tools_templates_effects_demo.pptx"))
        images = visual_qa._render_deck(pres, slides=[2])
        self.assertEqual(len(images), 1)
        self.assertTrue(images[0].startswith(b"\x89PNG"))
        # pres itself must be untouched: still all 34 slides, in order.
        self.assertEqual(len(pres.slides), 34)
        self.assertTrue(_slide_has_text(pres.slides[1], "Transform Your Business"))


class TestSubsetDeckBytes(unittest.TestCase):
    def test_keeps_only_requested_slides_in_order(self):
        from pptx import Presentation
        import io
        import visual_qa
        pres = Presentation(str(Path(__file__).resolve().parent.parent /
                                "mcp_all_tools_templates_effects_demo.pptx"))
        subset_bytes = visual_qa._subset_deck_bytes(pres, [3, 8])
        subset = Presentation(io.BytesIO(subset_bytes))
        self.assertEqual(len(subset.slides), 2)
        self.assertTrue(_slide_has_text(subset.slides[0], "Agenda"))
        self.assertTrue(_slide_has_text(subset.slides[1], "Performance Dashboard"))
        # Original presentation object is untouched.
        self.assertEqual(len(pres.slides), 34)


if __name__ == "__main__":
    unittest.main()
