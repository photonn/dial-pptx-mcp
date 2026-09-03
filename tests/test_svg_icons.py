"""
Tests for SVG icon rendering, its guards, and the tools the agent sees.

Rasterizing is PyMuPDF, a hard dependency, so these run everywhere — no
LibreOffice skip. The vision review is the only call that leaves the process,
and it is patched; nothing here touches DIAL, which is the point of the icon
store.
"""
import io
import sys
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import svg_icons
import visual_qa
from tools.icon_tools import GUIDANCE_PATH, register_icon_tools

ICON = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200">
  <circle cx="100" cy="100" r="88" fill="none" stroke="#005DB9" stroke-width="4"/>
  <g fill="none" stroke="#005DB9" stroke-width="4" stroke-linecap="round">
    <circle cx="100" cy="100" r="45"/>
    <polyline points="100,70 100,100 125,110"/>
  </g>
</svg>"""


class TestValidation(unittest.TestCase):
    """The SVG is written by a model and parsed in-process: it is input."""

    def reject(self, svg):
        with self.assertRaises(svg_icons.SvgIconError) as caught:
            svg_icons.validate_svg(svg)
        return str(caught.exception)

    def test_plain_line_art_is_accepted(self):
        self.assertIn("<svg", svg_icons.validate_svg(ICON))

    def test_empty_source_is_refused(self):
        self.reject("   ")

    def test_something_that_is_not_svg_is_refused(self):
        self.assertIn("SVG source", self.reject("/tmp/icon.svg"))

    def test_entity_declarations_are_refused(self):
        """The XXE shape: an entity expanding to a local file."""
        self.reject('<!DOCTYPE svg [<!ENTITY x SYSTEM "file:///etc/passwd">]>'
                    '<svg viewBox="0 0 10 10"><path d="M0 0"/></svg>')

    def test_scripts_and_foreign_objects_are_refused(self):
        self.reject('<svg viewBox="0 0 10 10"><script>x()</script></svg>')
        self.reject('<svg viewBox="0 0 10 10"><foreignObject/></svg>')
        self.reject('<svg viewBox="0 0 10 10" onload="x()"/>')

    def test_external_references_are_refused(self):
        """An SVG that fetches is an SSRF primitive, like a URL parameter."""
        self.reject('<svg viewBox="0 0 10 10">'
                    '<use href="https://evil.test/a.svg#i"/></svg>')
        self.reject('<svg viewBox="0 0 10 10">'
                    '<path fill="url(file:///etc/passwd)" d="M0 0"/></svg>')

    def test_an_internal_reference_is_still_fine(self):
        svg_icons.validate_svg('<svg viewBox="0 0 10 10">'
                               '<use href="#dot"/></svg>')

    def test_text_is_refused_with_a_reason(self):
        """Glyphs would come from a substituted font — the artifact class this
        whole module exists to catch."""
        message = self.reject('<svg viewBox="0 0 10 10"><text>KPI</text></svg>')
        self.assertIn("line art", message)

    def test_oversized_source_is_refused(self):
        with patch.dict("os.environ", {"SVG_ICON_MAX_KB": "1"}):
            self.reject('<svg viewBox="0 0 10 10">' + "<path d='M0 0'/>" * 200
                        + "</svg>")

    def test_an_unparsable_limit_falls_back(self):
        with patch.dict("os.environ", {"SVG_ICON_MAX_KB": "many"}):
            self.assertEqual(svg_icons.max_svg_bytes(),
                             int(svg_icons.DEFAULT_MAX_KB * 1024))


class TestColours(unittest.TestCase):
    def test_hex_forms(self):
        self.assertEqual(svg_icons.parse_color("#FFFFFF"), (255, 255, 255))
        self.assertEqual(svg_icons.parse_color("005db9"), (0, 93, 185))
        self.assertEqual(svg_icons.parse_color("#0AF"), (0, 170, 255))

    def test_transparent_is_none(self):
        self.assertIsNone(svg_icons.parse_color("transparent"))
        self.assertIsNone(svg_icons.parse_color(None))

    def test_a_colour_name_is_refused(self):
        with self.assertRaises(svg_icons.SvgIconError):
            svg_icons.parse_color("cornflowerblue")


class TestRender(unittest.TestCase):
    def render(self, svg=ICON, size=200, background=None):
        png, info = svg_icons.render_svg(svg, size, background)
        return Image.open(io.BytesIO(png)), info

    def test_it_renders_a_transparent_png_at_the_asked_size(self):
        image, info = self.render()
        self.assertEqual(image.format, "PNG")
        self.assertEqual((image.width, image.height), (200, 200))
        self.assertEqual(info["width"], 200)
        self.assertEqual(image.mode, "RGBA")
        self.assertLess(min(image.getchannel("A").getextrema()), 255,
                        "an icon with no background must stay transparent")

    def test_a_background_colour_is_composited_behind_it(self):
        image, _ = self.render(background=(255, 255, 255))
        self.assertEqual(image.getchannel("A").getextrema(), (255, 255))

    def test_the_drawings_aspect_ratio_is_kept(self):
        wide = ('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
                '<rect x="10" y="10" width="180" height="80" fill="none" '
                'stroke="#000" stroke-width="4"/></svg>')
        image, _ = self.render(wide, size=400)
        self.assertEqual((image.width, image.height), (400, 200))

    def test_a_drawing_outside_the_viewbox_is_reported_as_empty(self):
        """The classic hand-written-path failure: plausible XML, no picture."""
        blank = ('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200">'
                 '<circle cx="9000" cy="9000" r="20" fill="none" '
                 'stroke="#000" stroke-width="4"/></svg>')
        with self.assertRaises(svg_icons.SvgIconError) as caught:
            self.render(blank)
        self.assertIn("empty image", str(caught.exception))

    def test_malformed_svg_is_reported_as_such(self):
        with self.assertRaises(svg_icons.SvgIconError):
            self.render('<svg viewBox="0 0 200 200"><path d=')

    def test_coverage_notes_name_the_two_failure_modes(self):
        self.assertIn("Almost nothing", svg_icons.coverage_note(0.0001))
        self.assertIn("filled", svg_icons.coverage_note(0.99))
        self.assertIsNone(svg_icons.coverage_note(0.1))

    def test_a_filled_canvas_is_flagged_not_rejected(self):
        """Only a note: a solid disc is a legitimate icon background."""
        filled = ('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200">'
                  '<rect x="0" y="0" width="200" height="200" fill="#005DB9"/>'
                  '</svg>')
        _, info = self.render(filled)
        self.assertGreater(info["ink_coverage"], svg_icons.MAX_INK_COVERAGE)


class TestReviewPreview(unittest.TestCase):
    def test_the_second_image_is_the_icon_at_slide_size(self):
        png, _ = svg_icons.render_svg(ICON, 800, None)
        small = Image.open(io.BytesIO(svg_icons._preview_size(
            png, 96, (0, 93, 185))))
        self.assertEqual(max(small.width, small.height), 96)
        self.assertEqual(small.mode, "RGB",
                         "the reviewer must see it on the slide's colour, not "
                         "on an alpha channel")

    def test_review_sends_both_images_and_returns_the_verdict(self):
        png, _ = svg_icons.render_svg(ICON, 200, None)
        seen = {}

        def fake_review(self, images, prompt, timeout=300.0):
            seen["images"] = images
            seen["prompt"] = prompt
            return {"passed": True, "issues": [], "reads_as": "a clock"}

        with patch.object(visual_qa.VisionLLM, "__init__", lambda s: None), \
             patch.object(visual_qa.VisionLLM, "review", fake_review):
            verdict = svg_icons.review_icon(png, "time to market")
        self.assertTrue(verdict["passed"])
        self.assertEqual(len(seen["images"]), 2)
        self.assertIn("time to market", seen["prompt"])


class TestIconTool(unittest.TestCase):
    """The tool layer: guards, and the shape the agent is handed."""

    def setUp(self):
        from mcp.server.fastmcp import FastMCP
        from pptx import Presentation

        from state import PresentationStore

        app = FastMCP(name="test")
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        self.icons = svg_icons.IconStore(ttl_seconds=60, max_items=10)
        register_icon_tools(app, self.store, self.icons)
        self.tool = app._tool_manager._tools["render_svg_icon"]
        self.render = self.tool.fn
        self.place_tool = app._tool_manager._tools["add_icon_to_slide"]
        self.place = self.place_tool.fn
        self.guidance = app._tool_manager._tools["get_icon_guidance"].fn

        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[6])
        self.pid = self.store.new_id()
        self.store[self.pid] = pres

    def run_tool(self, verdict=None, **kwargs):
        kwargs.setdefault("svg", ICON)
        kwargs.setdefault("size", 200)
        with patch.object(visual_qa, "vision_configured",
                          return_value=verdict is not None), \
             patch.object(svg_icons, "review_icon", return_value=verdict):
            return self.render(**kwargs)

    def test_it_is_read_only(self):
        """It touches no deck; without the hint every call would mark one
        visually stale."""
        self.assertTrue(self.tool.annotations.readOnlyHint)

    def test_a_passing_icon_comes_back_placeable(self):
        result = self.run_tool(verdict={"passed": True, "issues": [],
                                        "reads_as": "a clock"},
                               concept="Time to market")
        self.assertIn(result["icon_id"], self.icons)
        self.assertEqual(result["mime_type"], svg_icons.PNG_MIME)
        self.assertEqual(result["review"]["passed"], True)
        self.assertIn("add_icon_to_slide", result["message"])

    def test_the_icon_never_goes_through_dial_storage(self):
        """A file this server writes lands in its own appdata, which the
        orchestrator cannot get granted back to itself — the 403 that made
        the handle necessary."""
        with patch("dial_client.DialFileClient") as client:
            self.run_tool(verdict={"passed": True, "issues": []})
        self.assertFalse(client.called)

    def test_a_failed_review_keeps_nothing_and_says_what_to_do(self):
        result = self.run_tool(verdict={
            "passed": False,
            "issues": [{"severity": "major", "description": "a blob",
                        "suggested_fix": "close the path"}]})
        self.assertNotIn("icon_id", result)
        self.assertEqual(len(self.icons), 0)
        self.assertIn("render_svg_icon again", result["message"])
        self.assertEqual(len(result["review"]["issues"]), 1)

    def test_review_can_be_skipped(self):
        result = self.run_tool(verdict={"passed": True, "issues": []},
                               review=False)
        self.assertNotIn("review", result)
        self.assertIn("icon_id", result)

    def test_without_a_vision_model_the_icon_still_ships_with_a_caveat(self):
        result = self.run_tool(verdict=None)
        self.assertIn("icon_id", result)
        self.assertIn("No vision model", result["review_note"])

    def test_the_slide_qa_switch_does_not_turn_off_the_icon_review(self):
        """VISUAL_QA_ENFORCE hides the slide inspect/repair tools; whether one
        icon can be looked at depends on the model existing, not on that."""
        with patch.dict("os.environ", {"VISUAL_QA_ENFORCE": "false",
                                       "VISION_LLM_MODEL": "gpt-4o",
                                       "VISION_LLM_PROVIDER": "dial",
                                       "DIAL_CORE_URL": "https://dial.test"}), \
             patch.object(svg_icons, "review_icon",
                          return_value={"passed": True, "issues": []}) as review:
            result = self.render(svg=ICON, size=200)
        self.assertTrue(review.called)
        self.assertEqual(result["review"]["passed"], True)

    def test_an_unreachable_reviewer_is_a_note_not_a_failure(self):
        with patch.object(visual_qa, "vision_configured", return_value=True), \
             patch.object(svg_icons, "review_icon",
                          side_effect=visual_qa.VisualQAError("endpoint down")):
            result = self.render(svg=ICON, size=200)
        self.assertIn("icon_id", result)
        self.assertIn("endpoint down", result["review_note"])

    def test_a_bad_svg_is_an_error_dict_not_an_exception(self):
        result = self.run_tool(svg="not svg at all")
        self.assertIn("error", result)

    def test_sizes_outside_the_range_are_refused(self):
        self.assertIn("error", self.run_tool(size=8))
        self.assertIn("error", self.run_tool(size=99999))

    def test_a_bad_colour_is_refused_before_rendering(self):
        self.assertIn("error", self.run_tool(background="chartreuse"))
        self.assertIn("error", self.run_tool(slide_background="#GGG"))

    def place_one(self, **kwargs):
        icon = self.run_tool(verdict={"passed": True, "issues": []})
        kwargs.setdefault("presentation_id", self.pid)
        kwargs.setdefault("slide_index", 0)
        kwargs.setdefault("icon_id", icon["icon_id"])
        return self.place(**kwargs)

    def test_placing_edits_the_deck_so_it_is_not_read_only(self):
        """The wrapper reads this hint to decide whether the deck now needs
        re-inspection — an icon on a slide certainly does."""
        self.assertFalse(getattr(self.place_tool.annotations, "readOnlyHint",
                                 False))

    def test_an_icon_is_placed_square_and_centred_in_its_box(self):
        result = self.place_one(left=1.0, top=2.0, size=0.8)
        self.assertEqual(result["shape_index"], 0)
        self.assertEqual(result["placed"]["width"], 0.8)
        self.assertEqual(result["placed"]["height"], 0.8)
        self.assertEqual(result["placed"]["left"], 1.0)
        self.assertEqual(len(self.store[self.pid].slides[0].shapes), 1)

    def test_one_icon_can_be_placed_repeatedly(self):
        icon = self.run_tool(verdict={"passed": True, "issues": []})
        for left in (1.0, 2.0, 3.0):
            self.assertIn("shape_index",
                          self.place(presentation_id=self.pid, slide_index=0,
                                     icon_id=icon["icon_id"], left=left))
        self.assertEqual(len(self.store[self.pid].slides[0].shapes), 3)

    def test_an_unknown_icon_handle_is_refused_with_the_way_out(self):
        result = self.place(presentation_id=self.pid, slide_index=0,
                            icon_id="deadbeef")
        self.assertIn("render the icon again", result["error"])

    def test_placement_guards(self):
        self.assertIn("error", self.place_one(presentation_id="nope"))
        self.assertIn("error", self.place_one(slide_index=7))
        self.assertIn("error", self.place_one(left=-1))
        self.assertIn("error", self.place_one(size=0))
        self.assertIn("error", self.place_one(size=99))

    def test_the_guidance_document_is_served(self):
        result = self.guidance()
        self.assertIn("viewBox", result["guidance"])
        self.assertGreater(result["characters"], 500)

    def test_a_missing_guidance_document_degrades_to_advice(self):
        with patch.object(Path, "stat", side_effect=OSError("gone")):
            result = self.guidance()
        self.assertIn("error", result)
        self.assertIn("line art", result["error"])


class TestIconStore(unittest.TestCase):
    """Handles are capabilities and memory is bounded, as for decks."""

    def test_a_handle_returns_the_bytes_and_the_metadata(self):
        store = svg_icons.IconStore()
        icon_id = store.put(b"png-bytes", {"concept": "gear"})
        png, meta = store.get(icon_id)
        self.assertEqual(png, b"png-bytes")
        self.assertEqual(meta["concept"], "gear")
        self.assertIn(icon_id, store)

    def test_handles_are_unguessable_and_unique(self):
        store = svg_icons.IconStore()
        ids = {store.put(b"x") for _ in range(5)}
        self.assertEqual(len(ids), 5)
        self.assertTrue(all(len(i) == 32 for i in ids))

    def test_an_unknown_handle_is_none_not_an_error(self):
        self.assertIsNone(svg_icons.IconStore().get("nope"))

    def test_the_oldest_icon_is_evicted_at_the_cap(self):
        store = svg_icons.IconStore(max_items=2)
        first = store.put(b"a")
        store.put(b"b")
        store.put(b"c")
        self.assertEqual(len(store), 2)
        self.assertNotIn(first, store)

    def test_icons_expire(self):
        store = svg_icons.IconStore(ttl_seconds=0)
        icon_id = store.put(b"a")
        self.assertIsNone(store.get(icon_id))


class TestGuidanceDocument(unittest.TestCase):
    def test_it_exists_and_carries_both_variants(self):
        text = GUIDANCE_PATH.read_text(encoding="utf-8")
        self.assertIn("viewBox=\"0 0 200 200\"", text)
        self.assertIn("render_svg_icon", text)
        self.assertIn("add_image_from_dial_url", text)


if __name__ == "__main__":
    unittest.main()
