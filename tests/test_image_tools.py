"""Tests for add_image_from_dial_url: DIAL download, validation, and the
aspect-ratio-aware placement modes."""
import io
import os
import sys
import unittest
from unittest.mock import MagicMock, patch
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from state import PresentationStore
from tools.image_tools import register_image_tools

EMU_PER_INCH = 914400


def png_bytes(width, height, color=(200, 30, 30)):
    buf = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buf, format="PNG")
    return buf.getvalue()


class _FakeApp:
    """Collects the functions registered with @app.tool()."""

    def __init__(self):
        self.tools = {}

    def tool(self, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn
        return decorator


class ImageToolTestCase(unittest.TestCase):
    def setUp(self):
        self._saved = os.environ.get("DIAL_IMAGE_MAX_MB")
        os.environ.pop("DIAL_IMAGE_MAX_MB", None)
        self.store = PresentationStore(ttl_seconds=60, max_items=10)
        self.pid = self.store.new_id()
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[6])
        self.store[self.pid] = pres
        app = _FakeApp()
        register_image_tools(app, self.store)
        self.add_image = app.tools["add_image_from_dial_url"]

    def tearDown(self):
        if self._saved is None:
            os.environ.pop("DIAL_IMAGE_MAX_MB", None)
        else:
            os.environ["DIAL_IMAGE_MAX_MB"] = self._saved

    def call(self, data=None, **kwargs):
        """Invoke the tool with the DIAL download stubbed out."""
        data = png_bytes(800, 400) if data is None else data
        client = MagicMock()
        client.download.return_value = data
        with patch("dial_client.DialFileClient", return_value=client):
            kwargs.setdefault("presentation_id", self.pid)
            kwargs.setdefault("slide_index", 0)
            kwargs.setdefault("image_url", "files/BUCKET/img/x.png")
            return self.add_image(**kwargs)

    def add_image_raw(self, **kwargs):
        """Invoke the tool with the real DialFileClient (HTTP stubbed out)."""
        kwargs.setdefault("presentation_id", self.pid)
        kwargs.setdefault("slide_index", 0)
        return self.add_image(**kwargs)

    def picture(self):
        return self.store[self.pid].slides[0].shapes[-1]


class TestValidation(ImageToolTestCase):
    def test_unknown_presentation(self):
        result = self.call(presentation_id="nope")
        self.assertIn("error", result)
        self.assertNotIn("shape_index", result)

    def test_bad_slide_index(self):
        self.assertIn("error", self.call(slide_index=4))
        self.assertIn("error", self.call(slide_index=-1))

    def test_bad_fit_mode(self):
        result = self.call(fit="squash")
        self.assertIn("error", result)
        self.assertIn("contain", result["error"])

    def test_non_positive_size(self):
        self.assertIn("error", self.call(width=0, height=2))
        self.assertIn("error", self.call(width=2, height=-1))
        self.assertIn("error", self.call(left=-1))

    def test_download_failure_is_reported_not_raised(self):
        client = MagicMock()
        client.download.side_effect = RuntimeError("404 not found")
        with patch("dial_client.DialFileClient", return_value=client):
            result = self.add_image(presentation_id=self.pid, slide_index=0,
                                    image_url="files/B/missing.png")
        self.assertIn("error", result)
        self.assertIn("404", result["error"])

    def test_non_image_bytes_get_a_useful_error(self):
        result = self.call(data=b"<svg xmlns='http://www.w3.org/2000/svg'/>")
        self.assertIn("error", result)
        self.assertIn("SVG", result["error"])

    def test_oversize_image_refused(self):
        os.environ["DIAL_IMAGE_MAX_MB"] = "0.001"
        result = self.call()
        self.assertIn("error", result)
        self.assertIn("limit", result["error"])
        self.assertEqual(len(self.store[self.pid].slides[0].shapes), 0)

    def test_invalid_max_mb_falls_back_to_default(self):
        os.environ["DIAL_IMAGE_MAX_MB"] = "not-a-number"
        self.assertNotIn("error", self.call())


class TestPlacement(ImageToolTestCase):
    def test_contain_preserves_aspect_ratio_and_centres(self):
        # 2:1 image into a 4x4in box -> 4x2in, centred vertically.
        result = self.call(left=1, top=1, width=4, height=4, fit="contain")
        self.assertEqual(result["fit"], "contain")
        self.assertAlmostEqual(result["placed"]["width"], 4.0, places=2)
        self.assertAlmostEqual(result["placed"]["height"], 2.0, places=2)
        self.assertAlmostEqual(result["placed"]["left"], 1.0, places=2)
        self.assertAlmostEqual(result["placed"]["top"], 2.0, places=2)

    def test_stretch_fills_the_box_exactly(self):
        result = self.call(left=1, top=1, width=4, height=4, fit="stretch")
        self.assertAlmostEqual(result["placed"]["width"], 4.0, places=2)
        self.assertAlmostEqual(result["placed"]["height"], 4.0, places=2)

    def test_cover_fills_the_box_and_crops_the_long_axis(self):
        result = self.call(left=1, top=1, width=4, height=4, fit="cover")
        self.assertAlmostEqual(result["placed"]["width"], 4.0, places=2)
        self.assertAlmostEqual(result["placed"]["height"], 4.0, places=2)
        pic = self.picture()
        self.assertAlmostEqual(pic.crop_left, 0.25, places=3)
        self.assertAlmostEqual(pic.crop_right, 0.25, places=3)
        self.assertEqual(pic.crop_top, 0)

    def test_cover_crops_the_other_axis_for_a_tall_image(self):
        result = self.call(data=png_bytes(400, 800), width=4, height=2,
                           fit="cover")
        self.assertEqual(result["fit"], "cover")
        pic = self.picture()
        self.assertAlmostEqual(pic.crop_top, 0.375, places=3)
        self.assertAlmostEqual(pic.crop_bottom, 0.375, places=3)
        self.assertEqual(pic.crop_left, 0)

    def test_single_dimension_scales_proportionally(self):
        result = self.call(width=6, fit="contain")
        self.assertEqual(result["fit"], "proportional")
        self.assertAlmostEqual(result["placed"]["height"], 3.0, places=2)

        result = self.call(height=1, fit="contain")
        self.assertEqual(result["fit"], "proportional")
        self.assertAlmostEqual(result["placed"]["width"], 2.0, places=2)

    def test_native_size_is_kept_when_it_fits(self):
        # 96x96 px at the default 72 dpi assumption -> well inside the slide.
        result = self.call(data=png_bytes(96, 96))
        self.assertEqual(result["fit"], "native")
        self.assertEqual(result["native"], {"width_px": 96, "height_px": 96})

    def test_oversized_native_image_is_clamped_to_the_slide(self):
        result = self.call(data=png_bytes(4000, 4000))
        self.assertEqual(result["fit"], "native_clamped")
        pres = self.store[self.pid]
        self.assertLessEqual(self.picture().width, pres.slide_width)
        self.assertLessEqual(self.picture().height, pres.slide_height)

    def test_reports_the_shape_it_added(self):
        result = self.call(width=3, height=3)
        self.assertEqual(result["shape_index"],
                         len(self.store[self.pid].slides[0].shapes) - 1)
        self.assertEqual(result["size_bytes"], len(png_bytes(800, 400)))


class TestUrlAcceptance(ImageToolTestCase):
    """No stubbing of DialFileClient here: a web URL must be refused before
    any HTTP request is made, and reported as an error dict, not raised."""

    def setUp(self):
        super().setUp()
        self._core = os.environ.get("DIAL_CORE_URL")
        os.environ["DIAL_CORE_URL"] = "https://dial.example.com"

    def tearDown(self):
        super().tearDown()
        if self._core is None:
            os.environ.pop("DIAL_CORE_URL", None)
        else:
            os.environ["DIAL_CORE_URL"] = self._core

    def test_public_web_url_is_refused_with_instructions(self):
        result = self.add_image(
            presentation_id=self.pid, slide_index=0,
            image_url="https://upload.wikimedia.org/wikipedia/commons/x.jpg")
        self.assertIn("error", result)
        self.assertIn("DIAL file storage", result["error"])
        self.assertEqual(len(self.store[self.pid].slides[0].shapes), 0)

    def test_dial_chat_link_is_accepted(self):
        """The .../api/files/... form an image deployment returns reaches the
        Files API instead of 404ing on a mangled path."""
        seen = {}

        def fake_get(url, headers=None, timeout=None):
            seen["url"] = url
            return MagicMock(status_code=200, content=png_bytes(64, 64),
                             raise_for_status=lambda: None)

        with patch("dial_client.httpx.get", fake_get), \
                patch("dial_client.resolve_dial_auth_headers",
                      return_value={"Api-Key": "k"}):
            result = self.add_image_raw(
                image_url="https://dial.example.com/api/files/BUCKET/img/x.png")
        self.assertNotIn("error", result)
        self.assertEqual(seen["url"],
                         "https://dial.example.com/v1/files/BUCKET/img/x.png")


if __name__ == "__main__":
    unittest.main()
