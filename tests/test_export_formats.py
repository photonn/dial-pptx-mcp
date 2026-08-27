"""
Tests for PDF export and legacy .ppt input.

Both go through LibreOffice, which is not installed everywhere the suite runs,
so the converter is stubbed: what is under test is the export tool's handling —
which files it produces, what it names them, and what it does when the renderer
is missing — not LibreOffice itself.
"""
import json
import os
import sys
import unittest
from pathlib import Path
from unittest.mock import patch

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from mcp.server.fastmcp import FastMCP
from pptx import Presentation

from state import PresentationStore
from tools.presentation_tools import register_presentation_tools, OLE_MAGIC
import visual_qa


class FakeDialClient:
    """Stands in for DialFileClient: records uploads, returns a file URL, and
    serves whatever `files` holds on the download side (the export path reads
    the brand profile through the same client)."""

    instances = []
    files = {}

    def __init__(self, *args, **kwargs):
        self.uploaded = []
        FakeDialClient.instances.append(self)

    def upload(self, data, filename, folder=None, content_type=None):
        self.uploaded.append({"filename": filename, "bytes": len(data),
                              "content_type": content_type})
        return f"files/test-bucket/{filename}"

    def download(self, ref):
        if ref not in FakeDialClient.files:
            from dial_client import DialConfigError
            raise DialConfigError("DIAL Core refused this file (404).")
        return FakeDialClient.files[ref]


def build_tools():
    """Register the presentation tools against a throwaway app and store."""
    app = FastMCP(name="test")
    store = PresentationStore(ttl_seconds=60, max_items=10)
    register_presentation_tools(app, store, lambda: None, lambda: ["."])
    return app._tool_manager._tools, store


class ExportTestCase(unittest.TestCase):
    def setUp(self):
        FakeDialClient.instances = []
        FakeDialClient.files = {}
        self.tools, self.store = build_tools()
        self.export = self.tools["export_presentation"].fn
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[6])
        self.pid = self.store.new_id()
        self.store[self.pid] = pres

    def run_export(self, pdf_bytes=b"%PDF-1.4 fake", **kwargs):
        with patch("dial_client.DialFileClient", FakeDialClient), \
             patch.object(visual_qa, "render_pptx_bytes_to_pdf",
                          return_value=pdf_bytes):
            return self.export(self.pid, **kwargs)

    @property
    def uploads(self):
        return [u for client in FakeDialClient.instances
                for u in client.uploaded]


class TestExportFormats(ExportTestCase):
    def test_pptx_is_the_default_and_uploads_one_file(self):
        result = self.run_export()
        self.assertNotIn("error", result)
        self.assertEqual([u["filename"] for u in self.uploads],
                         ["presentation.pptx"])
        self.assertEqual(result["files"][0]["format"], "pptx")

    def test_pdf_only(self):
        result = self.run_export(format="pdf")
        self.assertEqual([u["filename"] for u in self.uploads],
                         ["presentation.pdf"])
        self.assertEqual(result["mime_type"], "application/pdf")
        self.assertTrue(result["file_url"].endswith(".pdf"))

    def test_both_delivers_the_pair_with_the_pptx_first(self):
        result = self.run_export(format="both")
        self.assertEqual([f["format"] for f in result["files"]],
                         ["pptx", "pdf"])
        self.assertEqual([u["filename"] for u in self.uploads],
                         ["presentation.pptx", "presentation.pdf"])
        # file_url stays the editable deck, for callers reading only that key.
        self.assertTrue(result["file_url"].endswith(".pptx"))

    def test_the_extension_is_normalised_not_appended(self):
        for given, expected in (("deck", "deck.pptx"),
                                ("deck.pptx", "deck.pptx"),
                                ("deck.pdf", "deck.pptx")):
            with self.subTest(filename=given):
                FakeDialClient.instances = []
                self.run_export(filename=given)
                self.assertEqual(self.uploads[0]["filename"], expected)

    def test_pdf_carries_its_own_content_type(self):
        self.run_export(format="pdf")
        self.assertEqual(self.uploads[0]["content_type"], "application/pdf")

    def test_an_unknown_format_is_refused(self):
        result = self.run_export(format="xps")
        self.assertIn("error", result)
        self.assertIn("xps", result["error"])
        self.assertEqual(self.uploads, [])

    def test_an_empty_filename_is_refused(self):
        self.assertIn("error", self.run_export(filename=".pptx"))


class TestBrandSummaryOnExport(ExportTestCase):
    """Export reports where the deck stands against the rules attached to it;
    it never blocks on them, and says nothing at all where the server has no
    brand profile."""

    PROFILE = {"name": "Example Corp", "min_font_pt": 14,
               "families": {"light": {"require": ["slide_number"]}},
               "chrome_shapes": {"slide_number": {"type": "placeholder",
                                                  "idx": 12}}}

    def setUp(self):
        super().setUp()
        self._saved = os.environ.get("BRAND_PROFILE_FILE")
        os.environ.pop("BRAND_PROFILE_FILE", None)
        self.addCleanup(self._restore)

    def _restore(self):
        if self._saved is None:
            os.environ.pop("BRAND_PROFILE_FILE", None)
        else:
            os.environ["BRAND_PROFILE_FILE"] = self._saved

    def test_no_profile_configured_no_brand_key(self):
        result = self.run_export()
        self.assertNotIn("brand", result)
        self.assertNotIn("brand_note", result)

    def test_attached_rules_are_reported_and_do_not_block(self):
        os.environ["BRAND_PROFILE_FILE"] = "brand_profile.json"
        self.store.set_brand(self.pid, {"profile": self.PROFILE})
        result = self.run_export()
        self.assertNotIn("error", result)          # delivered regardless
        self.assertIs(result["brand"]["validated"], True)
        self.assertEqual(result["brand"]["brand"], "Example Corp")
        self.assertEqual(result["brand"]["warnings"], 1)  # no page number
        self.assertIn("validate_brand_profile", result["brand_note"])

    def test_a_deck_with_nothing_attached_is_told_so(self):
        """Delivered unchecked must not look the same as delivered clean."""
        os.environ["BRAND_PROFILE_FILE"] = "brand_profile.json"
        result = self.run_export()
        self.assertNotIn("error", result)
        self.assertTrue(result["file_url"])
        self.assertIs(result["brand"]["validated"], False)
        self.assertIn("attach_brand_profile", result["brand_note"])
        self.assertIn("not been checked", result["brand_note"])


class TestPdfFailureHandling(ExportTestCase):
    def failing_export(self, **kwargs):
        with patch("dial_client.DialFileClient", FakeDialClient), \
             patch.object(visual_qa, "render_pptx_bytes_to_pdf",
                          side_effect=visual_qa.VisualQAError("no soffice")):
            return self.export(self.pid, **kwargs)

    def test_a_failed_pdf_does_not_cost_the_user_the_deck(self):
        result = self.failing_export(format="both")
        self.assertNotIn("error", result)
        self.assertEqual([f["format"] for f in result["files"]], ["pptx"])
        self.assertIn("no soffice", result["pdf_note"])

    def test_pdf_only_reports_the_failure(self):
        result = self.failing_export(format="pdf")
        self.assertIn("error", result)
        self.assertIn("LibreOffice", result["error"])


class TestLegacyPptInput(unittest.TestCase):
    def setUp(self):
        self.tools, self.store = build_tools()
        self.create = self.tools["create_presentation_from_template_content"].fn

    @staticmethod
    def encoded(raw):
        import base64
        return base64.b64encode(raw).decode()

    def legacy_bytes(self):
        return OLE_MAGIC + b"legacy powerpoint payload"

    def real_pptx_bytes(self):
        import io
        buffer = io.BytesIO()
        Presentation().save(buffer)
        return buffer.getvalue()

    def test_a_legacy_file_is_converted_on_the_way_in(self):
        with patch.object(visual_qa, "convert_legacy_ppt",
                          return_value=self.real_pptx_bytes()) as convert:
            result = self.create(self.encoded(self.legacy_bytes()))
        convert.assert_called_once()
        self.assertNotIn("error", result)
        self.assertEqual(result["converted_from"], "ppt")
        self.assertIn("approximate", result["conversion_note"])

    def test_a_failed_conversion_says_what_to_do(self):
        with patch.object(visual_qa, "convert_legacy_ppt",
                          side_effect=visual_qa.VisualQAError("no soffice")):
            result = self.create(self.encoded(self.legacy_bytes()))
        self.assertIn("error", result)
        self.assertIn("97-2003", result["error"])
        self.assertIn(".pptx", result["error"])

    def test_a_normal_pptx_is_not_run_through_the_converter(self):
        with patch.object(visual_qa, "convert_legacy_ppt") as convert:
            result = self.create(self.encoded(self.real_pptx_bytes()))
        convert.assert_not_called()
        self.assertNotIn("error", result)
        self.assertNotIn("converted_from", result)

    def test_junk_is_still_rejected(self):
        result = self.create(self.encoded(b"not a presentation at all"))
        self.assertIn("error", result)


class TestConversionHelpers(unittest.TestCase):
    """The soffice wrapper's contract, without invoking LibreOffice."""

    def test_pdf_conversion_asks_for_a_pdf_from_a_pptx(self):
        with patch.object(visual_qa, "convert_with_soffice",
                          return_value=b"%PDF") as convert:
            visual_qa.render_pptx_bytes_to_pdf(b"deck")
        convert.assert_called_once_with(b"deck", ".pptx", "pdf")

    def test_legacy_conversion_asks_for_a_pptx_from_a_ppt(self):
        with patch.object(visual_qa, "convert_with_soffice",
                          return_value=b"PK") as convert:
            visual_qa.convert_legacy_ppt(b"legacy")
        convert.assert_called_once_with(b"legacy", ".ppt", "pptx")

    def test_a_missing_renderer_raises_the_documented_error(self):
        with patch.object(visual_qa, "_soffice_binary",
                          side_effect=visual_qa.VisualQAError("not found")):
            with self.assertRaises(visual_qa.VisualQAError):
                visual_qa.convert_with_soffice(b"deck", ".pptx", "pdf")


if __name__ == "__main__":
    unittest.main()
