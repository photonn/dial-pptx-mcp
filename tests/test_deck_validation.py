"""
Tests for structural validation and the font-metrics advisory.

Each check is driven by a deck deliberately broken in one way, because a
validator that reports nothing passes trivially. The demo fixture doubles as
the "clean deck" case: it must produce no errors.
"""
import io
import unittest

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import deck_validation
import fonts

DEMO = "mcp_all_tools_templates_effects_demo.pptx"


def codes(report, severity=None):
    return {p["code"] for p in report["problems"]
            if severity is None or p["severity"] == severity}


def blank_deck(slides=1):
    pres = Presentation()
    for _ in range(slides):
        pres.slides.add_slide(pres.slide_layouts[6])  # blank
    return pres


class TestCleanDeck(unittest.TestCase):
    def test_the_committed_fixture_has_no_errors(self):
        report = deck_validation.validate_presentation(Presentation(DEMO))
        self.assertTrue(report["ok"], codes(report, "error"))
        self.assertEqual(report["counts"]["error"], 0)

    def test_a_blank_deck_is_clean(self):
        report = deck_validation.validate_presentation(blank_deck())
        self.assertTrue(report["ok"])
        self.assertEqual(codes(report, "warning"), set())

    def test_report_shape(self):
        report = deck_validation.validate_presentation(blank_deck())
        self.assertEqual(
            set(report), {"ok", "slides", "size_bytes", "counts", "problems"})
        self.assertEqual(report["slides"], 1)
        self.assertGreater(report["size_bytes"], 0)


class TestGeometryChecks(unittest.TestCase):
    def test_shape_past_the_edge_is_a_warning(self):
        pres = blank_deck()
        slide = pres.slides[0]
        # Slide is 10in wide by default; this box runs 2in past the right edge.
        slide.shapes.add_textbox(Inches(9), Inches(1), Inches(3), Inches(1))
        report = deck_validation.validate_presentation(pres)
        self.assertIn("shape_past_edge", codes(report, "warning"))
        self.assertTrue(report["ok"], "a shape past the edge is not an error")

    def test_shape_entirely_off_the_slide_is_an_error(self):
        pres = blank_deck()
        pres.slides[0].shapes.add_textbox(
            Inches(20), Inches(1), Inches(2), Inches(1))
        report = deck_validation.validate_presentation(pres)
        self.assertIn("shape_off_slide", codes(report, "error"))
        self.assertFalse(report["ok"])

    def test_a_shape_inside_the_canvas_is_not_reported(self):
        pres = blank_deck()
        pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(2))
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("shape_past_edge", codes(report))
        self.assertNotIn("shape_off_slide", codes(report))

    def test_the_problem_names_the_slide_and_shape(self):
        pres = blank_deck(2)
        pres.slides[1].shapes.add_textbox(
            Inches(20), Inches(1), Inches(2), Inches(1))
        report = deck_validation.validate_presentation(pres)
        problem = next(p for p in report["problems"]
                       if p["code"] == "shape_off_slide")
        self.assertEqual(problem["slide_index"], 1)
        self.assertEqual(problem["shape_index"], 0)
        self.assertIn("visual_repair_slides", problem["fix"])


class TestPartialTransform(unittest.TestCase):
    """A half-written a:xfrm — the LibreOffice-only-looks-fine defect.

    python-pptx reports the inherited value for the half that is missing, so
    the shape reads as correctly placed through the API and renders correctly
    in LibreOffice. Only PowerPoint moves it to the corner, which is why the
    geometry check cannot catch this and a separate one has to.
    """

    def broken_deck(self, *, drop):
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[1])  # title + content
        shape = slide.placeholders[1]
        xfrm = shape._element.spPr.find(qn("a:xfrm"))
        if xfrm is None:
            shape.left, shape.top = shape.left, shape.top
            shape.width, shape.height = shape.width, shape.height
            xfrm = shape._element.spPr.find(qn("a:xfrm"))
        xfrm.remove(xfrm.find(qn("a:" + drop)))
        return pres, shape

    def test_missing_offset_is_a_warning(self):
        pres, shape = self.broken_deck(drop="off")
        report = deck_validation.validate_presentation(pres)
        self.assertIn("partial_transform", codes(report, "warning"))
        # The geometry check sees nothing wrong — that is the whole point.
        self.assertNotIn("shape_off_slide", codes(report))
        self.assertIsNotNone(shape.left)

    def test_missing_extent_is_a_warning(self):
        pres, _ = self.broken_deck(drop="ext")
        report = deck_validation.validate_presentation(pres)
        self.assertIn("partial_transform", codes(report, "warning"))

    def test_the_problem_names_the_shape_and_a_fix(self):
        pres, _ = self.broken_deck(drop="off")
        report = deck_validation.validate_presentation(pres)
        problem = next(p for p in report["problems"]
                       if p["code"] == "partial_transform")
        self.assertEqual(problem["slide_index"], 0)
        self.assertEqual(problem["shape_index"], 1)
        self.assertIn("visual_repair_slides", problem["fix"])
        self.assertIn("PowerPoint", problem["message"])

    def test_a_complete_transform_is_not_reported(self):
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        shape = slide.placeholders[1]
        shape.left, shape.top = shape.left, shape.top
        shape.width, shape.height = shape.width, shape.height
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("partial_transform", codes(report))

    def test_an_inheriting_placeholder_is_not_reported(self):
        # No a:xfrm at all is correct: the placeholder inherits all four
        # values and every renderer agrees on them.
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[1])
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("partial_transform", codes(report))


class TestPlaceholderText(unittest.TestCase):
    def assert_flags(self, text):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = text
        report = deck_validation.validate_presentation(pres)
        self.assertIn("leftover_placeholder_text", codes(report, "warning"),
                      f"{text!r} was not flagged")

    def test_common_filler_is_flagged(self):
        for text in ("Lorem ipsum dolor sit amet",
                     "Click to add title",
                     "XXXX",
                     "TODO: write this",
                     "[insert customer name]",
                     "Your title here",
                     "Sample text"):
            with self.subTest(text=text):
                self.assert_flags(text)

    def test_real_content_is_not_flagged(self):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "Revenue grew 18% year on year"
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("leftover_placeholder_text", codes(report))

    def test_placeholder_text_in_a_table_cell_is_flagged(self):
        pres = blank_deck()
        table = pres.slides[0].shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(1, 1).text_frame.text = "TODO"
        report = deck_validation.validate_presentation(pres)
        self.assertIn("leftover_placeholder_text", codes(report, "warning"))

    def test_an_empty_placeholder_is_not_reported(self):
        """Its prompt text shows only in edit view, so it is not a defect —
        and a template deck has dozens."""
        pres = Presentation()
        pres.slides.add_slide(pres.slide_layouts[1])  # title and content
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("empty_placeholder", codes(report))


class TestChartAndTableChecks(unittest.TestCase):
    def test_chart_with_data_is_clean(self):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        pres = blank_deck()
        data = CategoryChartData()
        data.categories = ["a", "b"]
        data.add_series("s", (1, 2))
        pres.slides[0].shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
            Inches(4), Inches(3), data)
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("chart_without_series", codes(report))
        self.assertNotIn("chart_series_length_mismatch", codes(report))

    def test_series_shorter_than_its_categories_is_reported(self):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        pres = blank_deck()
        data = CategoryChartData()
        data.categories = ["a", "b", "c"]
        data.add_series("s", (1, 2, None))
        chart = pres.slides[0].shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
            Inches(4), Inches(3), data).chart
        # Drop a point so the series is genuinely shorter than the categories.
        chart_ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
        series = chart.plots[0].series[0]
        points = list(series._element.iter(f"{chart_ns}val"))
        if not points:
            self.skipTest("chart XML shape differs in this python-pptx build")
        last_point = list(points[0].iter(f"{chart_ns}pt"))[-1]
        last_point.getparent().remove(last_point)
        for count in points[0].iter(f"{chart_ns}ptCount"):
            count.set("val", "2")
        report = deck_validation.validate_presentation(pres)
        self.assertIn("chart_series_length_mismatch", codes(report, "warning"))

    def test_empty_table_is_an_error(self):
        pres = blank_deck()
        graphic_frame = pres.slides[0].shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = graphic_frame.table
        for row in list(table.rows._tbl.tr_lst):
            table._tbl.remove(row)
        report = deck_validation.validate_presentation(pres)
        self.assertIn("empty_table", codes(report, "error"))


class TestPackageChecks(unittest.TestCase):
    def test_unreadable_package_is_reported(self):
        report = deck_validation._Report()
        deck_validation._check_package(b"not a zip file at all", report)
        self.assertIn("package_unreadable", {p["code"] for p in report.problems})

    def test_dangling_relationship_is_an_error(self):
        pres = blank_deck()
        picture = pres.slides[0].shapes.add_picture(
            _one_pixel_png(), Inches(1), Inches(1), Inches(1), Inches(1))
        # Point the picture at a relationship the slide does not have.
        blip = picture._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
        blip.set("{http://schemas.openxmlformats.org/officeDocument/2006/"
                 "relationships}embed", "rId999")
        report = deck_validation.validate_presentation(pres)
        self.assertIn("dangling_relationship", codes(report, "error"))
        self.assertFalse(report["ok"])

    def test_slide_ids_are_checked_for_duplicates(self):
        pres = blank_deck(2)
        entries = list(pres.slides._sldIdLst)
        entries[1].set("id", entries[0].get("id"))
        report = deck_validation.validate_presentation(pres)
        self.assertIn("duplicate_slide_id", codes(report, "error"))

    def test_slide_id_out_of_range_is_reported(self):
        pres = blank_deck()
        list(pres.slides._sldIdLst)[0].set("id", "12")
        report = deck_validation.validate_presentation(pres)
        self.assertIn("slide_id_out_of_range", codes(report, "error"))


class TestPictureDistortion(unittest.TestCase):
    def test_a_stretched_picture_is_reported(self):
        pres = blank_deck()
        pres.slides[0].shapes.add_picture(
            _one_pixel_png(), Inches(1), Inches(1), Inches(4), Inches(1))
        report = deck_validation.validate_presentation(pres)
        self.assertIn("distorted_picture", codes(report, "warning"))

    def test_a_proportional_picture_is_not_reported(self):
        pres = blank_deck()
        pres.slides[0].shapes.add_picture(
            _one_pixel_png(), Inches(1), Inches(1), Inches(2), Inches(2))
        report = deck_validation.validate_presentation(pres)
        self.assertNotIn("distorted_picture", codes(report))


class TestFontAdvisory(unittest.TestCase):
    def test_metric_safe_fonts(self):
        for name in ("Arial", "calibri", "Times New Roman", "Cambria"):
            with self.subTest(name=name):
                self.assertTrue(fonts.is_metric_safe(name))
        for name in ("Georgia", "Aptos", "Segoe UI", "", None):
            with self.subTest(name=name):
                self.assertFalse(fonts.is_metric_safe(name))

    def test_fonts_used_by_a_slide_are_found(self):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "hello"
        run.font.name = "Georgia"
        run.font.size = Pt(18)
        self.assertIn("Georgia", fonts.fonts_in(pres))
        self.assertIn("Georgia", fonts.unreliable_fonts_in(pres))

    def test_a_calibri_deck_reports_no_risky_fonts(self):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "hello"
        run.font.name = "Calibri"
        # The default template's theme fonts are Calibri/Calibri Light; only
        # the slide-level choice is under test here.
        self.assertNotIn("Calibri", fonts.unreliable_fonts_in(pres))

    def test_aptos_is_discouraged(self):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "hello"
        run.font.name = "Aptos"
        self.assertEqual(fonts.discouraged_fonts_in(pres), {"Aptos"})

    def test_risky_fonts_surface_as_an_info_problem(self):
        pres = blank_deck()
        box = pres.slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "hello"
        run.font.name = "Trebuchet MS"
        report = deck_validation.validate_presentation(pres)
        self.assertIn("qa_unreliable_fonts", codes(report, "info"))
        self.assertTrue(report["ok"], "an advisory is never an error")

    def test_the_review_caveat_names_the_fonts(self):
        caveat = fonts.qa_font_caveat({"Georgia", "Impact"})
        self.assertIn("Georgia", caveat)
        self.assertIn("Impact", caveat)
        self.assertEqual(fonts.qa_font_caveat(set()), "")

    def test_the_caveat_reaches_the_review_prompt(self):
        import visual_qa
        prompt = visual_qa.review_prompt(False, None, None, {"Georgia"})
        self.assertIn("Georgia", prompt)
        self.assertNotIn("Georgia", visual_qa.review_prompt(False))


def _one_pixel_png():
    """A 1x1 PNG, so any non-square placement is a distortion."""
    import base64
    return io.BytesIO(base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQ"
        "DwAEhQGAhKmMIQAAAABJRU5ErkJggg=="))


if __name__ == "__main__":
    unittest.main()
