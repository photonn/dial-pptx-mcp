"""
Put the template's slide-number placeholder on the slides.

A slide number is not inherited. A template can define where the number sits
and how it looks — position, font, colour — on the master and on every layout,
and PowerPoint will still render nothing until each *slide* carries its own
`<p:ph type="sldNum">` shape containing a `slidenum` field. python-pptx's
`add_slide` deliberately does not clone it: `iter_cloneable_placeholders`
excludes date, footer and slide-number placeholders, because PowerPoint treats
those as opt-in per deck rather than per layout. `duplicate_slide` copies
whatever the source slide has, so a deck built the intended duplicate-then-fill
way inherits the gap from the first slide it filled.

So an agent building a deck from a corporate template gets every other piece of
the template's chrome for free and silently loses the page numbers. This closes
that gap by copying the layout's own sldNum placeholder — position, formatting
and field intact — onto each slide, which is what PowerPoint does when you tick
"Slide number" in Insert > Header & Footer.

Entirely generic: it copies whatever the template defined and knows nothing
about any particular brand. Slides whose layout has no sldNum placeholder are
left alone, which is how a template says "no number on this layout" for its
title and closing slides.
"""
import copy
from typing import Dict, List, Optional

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pptx.oxml.ns import qn

from logging_utils import get_logger
from state import short_id

logger = get_logger("tools.slide_number")

UNKNOWN_ID = (
    "Unknown or expired presentation_id. Pass the presentation_id returned "
    "by create_presentation, create_presentation_from_template, or "
    "open_presentation"
)

_SLDNUM = "sldNum"


def _placeholder_type(shape):
    """The `type` attribute of a shape's `<p:ph>`, or None if it has none.

    Read off the XML rather than through `shape.placeholder_format`: a shape
    copied in at the OPC level is not always bound to a placeholder object,
    and the layout's shapes are walked here as elements too.
    """
    nvSpPr = shape.element.find(qn("p:nvSpPr"))
    if nvSpPr is None:
        return None
    nvPr = nvSpPr.find(qn("p:nvPr"))
    if nvPr is None:
        return None
    ph = nvPr.find(qn("p:ph"))
    return None if ph is None else ph.get("type")


def _layout_sldnum(slide):
    """The sldNum placeholder element on this slide's layout, if it has one."""
    for shape in slide.slide_layout.placeholders:
        if _placeholder_type(shape) == _SLDNUM:
            return shape.element
    return None


def _has_sldnum(slide):
    for shape in slide.shapes:
        if _placeholder_type(shape) == _SLDNUM:
            return True
    return False


def register_slide_number_tools(app: FastMCP, presentations):
    """Register the slide-number tool."""

    @app.tool(
        annotations=ToolAnnotations(
            title="Add Slide Numbers",
        ),
    )
    def add_slide_numbers(presentation_id: str,
                          skip_slides: Optional[List[int]] = None) -> Dict:
        """Render the template's slide numbers on the deck's slides.

        Copies each slide's layout slide-number placeholder onto the slide, so
        the number appears where and how the template defines it — python-pptx
        does not do this when a slide is added, which is why a deck built from
        a template that styles its page numbers still shows none.

        Call it once after the last slide is in place. Slides whose layout
        defines no slide-number placeholder are skipped, as are slides that
        already have one, so calling it twice is safe and calling it again
        after adding more slides only fills the gaps.

        skip_slides: 0-based indexes to leave unnumbered. Usually unnecessary —
        a template that wants no number on its title and closing slides says so
        by omitting the placeholder from those layouts.
        """
        if presentation_id not in presentations:
            return {"error": UNKNOWN_ID}

        pres = presentations[presentation_id]
        skip = set(skip_slides or [])
        added, already, no_layout_ph = 0, 0, 0

        for index, slide in enumerate(pres.slides):
            if index in skip:
                continue
            if _has_sldnum(slide):
                already += 1
                continue
            element = _layout_sldnum(slide)
            if element is None:
                no_layout_ph += 1
                continue
            slide.shapes._spTree.append(copy.deepcopy(element))
            added += 1

        logger.info("slide_numbers_added presentation_id=%s added=%d "
                    "already_present=%d no_layout_placeholder=%d",
                    short_id(presentation_id), added, already, no_layout_ph)
        result = {
            "message": f"Added slide numbers to {added} slide(s).",
            "added": added,
            "already_present": already,
            "layout_has_no_slide_number_placeholder": no_layout_ph,
        }
        if no_layout_ph and not added:
            result["note"] = (
                "No slide got a number: none of the layouts in use define a "
                "slide-number placeholder. This template does not carry page "
                "numbers — add text boxes yourself if the deck needs them."
            )
        elif no_layout_ph:
            result["note"] = (
                f"{no_layout_ph} slide(s) use a layout with no slide-number "
                "placeholder and were left unnumbered — that is how a template "
                "marks its title and closing slides."
            )
        return result
