"""
Custom icons drawn as SVG by the orchestrator, rendered and vetted here.

The icon problem a deck runs into is not "which icon", it is "which icon on
which colour": a stock PNG carries an opaque background, so the moment a card
is tinted or a panel is brand blue, a library icon shows as a white square.
Drawing the icon per slide solves that — an SVG with no background renders to a
transparent PNG — and it also covers the concepts no library has.

The split is the same one `add_image_from_dial_url` makes: the orchestrator
composes (it has the style guide, served by `get_icon_guidance`), the server
renders. What is added here is the check in between. The agent cannot see the
icon it just drew, and once a bad one is on a slide the visual-repair
whitelist can only move, resize or delete it — so the icon is reviewed on its
own, before placement, while regenerating it is one cheap call.

Unlike a generated illustration, the icon never goes through DIAL file storage.
It was rendered in this process, so a round trip would only push it across an
identity boundary it cannot come back over: a file this server writes lands in
`{user}/appdata/{this-deployment}/`, and placing it again means the
*orchestrator* asking Core to grant that file to the toolset key — which Core
refuses, because the orchestrator is neither the user who owns the bucket nor
the deployment that owns the appdata folder. So the PNG stays in an `IconStore`
under an unguessable handle and `add_icon_to_slide` places it from there.
"""
import io
from pathlib import Path
from typing import Dict, Optional

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pptx.util import Inches

from logging_utils import get_logger
from state import short_id

logger = get_logger("tools.icon")

GUIDANCE_PATH = (Path(__file__).resolve().parent.parent / "docs"
                 / "ICON_GUIDANCE.md")

UNKNOWN_ID = (
    "Unknown or expired presentation_id. Pass the presentation_id returned "
    "by create_presentation, create_presentation_from_template, or "
    "open_presentation"
)

DEFAULT_ICON_SIZE = 0.8      # inches — an icon leading a card
MAX_ICON_SIZE = 4.0

_cache = {}


def _guidance():
    """The icon style guide, cached on the file's mtime (as
    get_design_guidance does) so an edit needs no restart."""
    try:
        stamp = GUIDANCE_PATH.stat().st_mtime
    except OSError as e:
        logger.warning("icon_guidance_unreadable path=%s error=%s",
                       GUIDANCE_PATH, e)
        return None
    if _cache.get("stamp") != stamp:
        _cache.update(stamp=stamp,
                      text=GUIDANCE_PATH.read_text(encoding="utf-8"))
    return _cache["text"]


def register_icon_tools(app: FastMCP, presentations, icons=None):
    """Register the SVG icon tools. No renderer dependency: PyMuPDF rasterizes
    the SVG, so these work even where LibreOffice (and therefore visual QA) is
    absent — the review step is the only part that needs a vision model.

    `icons` is the rendered-icon store; one is created per registration, which
    is per server process."""
    import svg_icons
    import visual_qa

    icons = icons if icons is not None else svg_icons.IconStore()
    logger.info("icon_tools_registered review=%s", visual_qa.vision_configured())

    @app.tool(
        annotations=ToolAnnotations(
            title="Get Icon Guidance",
            readOnlyHint=True,
        ),
    )
    def get_icon_guidance() -> Dict:
        """The style guide for drawing an icon as SVG, and how to use it.

        Read this before writing SVG for render_svg_icon: it holds the two
        icon templates (line art for a light surface, white-on-filled for a
        coloured one), the rules that keep a hand-drawn icon looking like the
        rest of the set, worked path examples you can adapt, and the decision
        of when a custom icon is the right answer at all.

        Returns the document as "guidance".
        """
        text = _guidance()
        if text is None:
            return {"error": f"The icon guidance document is missing from this "
                             f"server ({GUIDANCE_PATH.name}). Draw icons as "
                             f"simple line art on a 200x200 viewBox: one "
                             f"stroke weight, fill=\"none\", rounded caps and "
                             f"joins, no text."}
        return {"guidance": text, "characters": len(text)}

    @app.tool(
        annotations=ToolAnnotations(
            title="Render SVG Icon",
            readOnlyHint=True,
        ),
    )
    def render_svg_icon(
        svg: str,
        concept: Optional[str] = None,
        size: int = svg_icons.DEFAULT_PX,
        background: str = "transparent",
        slide_background: str = "#FFFFFF",
        review: bool = True,
    ) -> Dict:
        """Render an SVG icon you have written to a PNG, and check it.

        Use this when a slide needs an icon: call get_icon_guidance, write the
        SVG yourself following it, pass it here, and then place the returned
        icon_id with add_icon_to_slide. PowerPoint cannot embed SVG, which is
        why the rasterizing happens here rather than in your answer.

        The icon is kept on the server, not in DIAL file storage, so the
        handle is all you need — and add_icon_to_slide is what places it, not
        add_image_from_dial_url (that tool reads DIAL storage, which is for
        images your image model generated).

        The returned icon has been looked at. You cannot see what you drew, so
        the vision model is asked whether the render is free of artifacts and
        still readable at slide size. **If "passed" is false, do not place the
        icon**: fix the SVG as the issues say and call this again. An icon_id
        is returned only for an icon that is fit to use.

        svg: the SVG source, self-contained line art — <path>, <circle>,
          <rect>, <line>, <polyline>, <polygon>. No text elements, no embedded
          images, no external references; a text label belongs in a textbox on
          the slide, not in the icon.
        concept: what the icon depicts ("supply chain", "regulatory
          approval"). Passed to the reviewer, which is the difference between
          "this is a clean pictogram" and "this reads as the thing you meant".
        size: longest side of the PNG in pixels (default 800). An icon is
          displayed at about an inch, so this is deliberate headroom; there is
          no reason to raise it.
        background: the icon's own background. Leave "transparent" — that is
          the whole advantage over a stock icon, and it is what lets the same
          icon sit on a white slide and on a brand-coloured panel. A hex
          colour bakes a solid rectangle in behind it.
        slide_background: hex colour of the surface the icon will sit on. The
          review judges legibility and contrast against it, so set it to the
          panel colour for an icon going on a coloured card.
        review: set false only to skip the vision check (faster, unverified).

        Returns "icon_id" (pass it straight to add_icon_to_slide), the pixel
        size, and "review" with the verdict and any issues. The handle stays
        valid for the session, so reuse it wherever the icon repeats instead
        of rendering the same drawing again.
        """
        if size < svg_icons.MIN_PX or size > svg_icons.MAX_PX:
            return {"error": f"size must be between {svg_icons.MIN_PX} and "
                             f"{svg_icons.MAX_PX} pixels."}
        try:
            canvas = svg_icons.parse_color(background, "background")
            surface = svg_icons.parse_color(slide_background,
                                            "slide_background")
        except svg_icons.SvgIconError as e:
            return {"error": str(e)}

        try:
            png, info = svg_icons.render_svg(svg, size, canvas)
        except svg_icons.SvgIconError as e:
            logger.info("icon_render_rejected concept=%s reason=%s",
                        (concept or "-")[:40], e)
            return {"error": str(e)}

        result = {
            "image_size": {"width": info["width"], "height": info["height"]},
            "size_bytes": len(png),
            "mime_type": svg_icons.PNG_MIME,
            "ink_coverage": info["ink_coverage"],
        }
        # Named before any model call: a blank or solid-black render has a
        # specific cause in the SVG, and saying so beats a vision verdict.
        coverage_note = svg_icons.coverage_note(info["ink_coverage"])
        if coverage_note:
            result["render_note"] = coverage_note

        verdict = None
        if review:
            # vision_configured(), not enforcement_enabled(): VISUAL_QA_ENFORCE
            # switches off slide inspection and repair, which is a different
            # decision from whether a model exists to look at one icon.
            if not visual_qa.vision_configured():
                result["review_note"] = (
                    "No vision model is configured on this server, so the "
                    "icon could not be checked. Place it if you are confident "
                    "in the SVG, and check the deck with the user."
                )
            else:
                try:
                    verdict = svg_icons.review_icon(png, concept, surface)
                except visual_qa.VisualQAError as e:
                    logger.warning("icon_review_failed concept=%s error=%s",
                                   (concept or "-")[:40], e)
                    result["review_note"] = (
                        f"The icon was rendered but could not be reviewed "
                        f"({e}). Its handle is below, unchecked."
                    )
                else:
                    result["review"] = verdict

        if verdict is not None and verdict.get("passed") is False:
            # Nothing to place, so nothing is kept: the next call carries a
            # corrected SVG, not this handle.
            logger.info("icon_rejected concept=%s issues=%d",
                        (concept or "-")[:40], len(verdict.get("issues", [])))
            result["message"] = (
                "This icon did not pass review, so it was not kept and there "
                "is no icon_id to place. Rewrite the SVG following the issues "
                "above and call render_svg_icon again. If two attempts do not "
                "fix it, simplify the drawing — fewer, larger shapes — or "
                "leave the icon out rather than placing a broken one."
            )
            return result

        icon_id = icons.put(png, {"concept": concept,
                                  "width": info["width"],
                                  "height": info["height"]})
        result["icon_id"] = icon_id
        result["message"] = (
            f"Icon ready (icon_id {icon_id}). Place it with add_icon_to_slide, "
            "and reuse the same icon_id wherever this icon repeats."
        )
        logger.info("icon_ready icon_id=%s concept=%s bytes=%d reviewed=%s "
                    "held=%d", short_id(icon_id), (concept or "-")[:40],
                    len(png), bool(verdict) and verdict.get("passed"),
                    len(icons))
        return result

    @app.tool(
        annotations=ToolAnnotations(
            title="Add Icon to Slide",
        ),
    )
    def add_icon_to_slide(
        presentation_id: str,
        slide_index: int,
        icon_id: str,
        left: float = 1.0,
        top: float = 1.0,
        size: float = DEFAULT_ICON_SIZE,
    ) -> Dict:
        """Place an icon you rendered with render_svg_icon onto a slide.

        The icon lives on this server, so this takes the icon_id that
        render_svg_icon returned — not a DIAL file URL. Place the same
        icon_id as many times as the icon repeats; the deck stores the image
        once however often it appears.

        presentation_id: the deck handle.
        slide_index: 0-based, like the other content tools (visual QA slide
          numbers are 1-based; slide_index 0 is slide 1).
        left, top: inches, the top-left corner of the icon's box.
        size: the box's side in inches — icons are square. About 0.7 for an
          icon inside a card, 1.0-1.2 for one leading a section. A non-square
          drawing is centred in that box rather than stretched to fill it.

        Returns the shape index and the geometry actually applied.
        """
        # The same geometry helpers the image tool uses: an icon is a picture,
        # and "contain" placement should mean the same thing for both.
        from tools.image_tools import _inches, _place

        if presentation_id not in presentations:
            return {"error": UNKNOWN_ID}
        pres = presentations[presentation_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Available "
                             f"slides: 0-{len(pres.slides) - 1}"}
        if left < 0 or top < 0:
            return {"error": "left and top must be zero or positive inches."}
        if size <= 0 or size > MAX_ICON_SIZE:
            return {"error": f"size must be between 0 and {MAX_ICON_SIZE} "
                             f"inches — an icon is roughly 0.7-1.2in."}

        entry = icons.get(icon_id)
        if entry is None:
            return {"error": "Unknown or expired icon_id. Pass the icon_id "
                             "returned by render_svg_icon in this session; "
                             "if it has expired, render the icon again."}
        png, meta = entry

        slide = pres.slides[slide_index]
        pic = slide.shapes.add_picture(io.BytesIO(png), Inches(left),
                                       Inches(top))
        _place(pic, int(Inches(left)), int(Inches(top)), int(Inches(size)),
               int(Inches(size)), "contain", int(pres.slide_width),
               int(pres.slide_height))

        logger.info("icon_placed presentation_id=%s slide=%d icon_id=%s "
                    "concept=%s", short_id(presentation_id), slide_index,
                    short_id(icon_id), (meta.get("concept") or "-")[:40])
        return {
            "message": f"Added icon to slide {slide_index}.",
            "shape_index": len(slide.shapes) - 1,
            "placed": {
                "left": _inches(pic.left),
                "top": _inches(pic.top),
                "width": _inches(pic.width),
                "height": _inches(pic.height),
            },
        }
