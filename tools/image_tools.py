"""
DIAL-hosted image insertion for the PowerPoint MCP Server.

The server does not generate images: the orchestrator calls whatever image
deployment it likes, stores the result in DIAL file storage, and passes the
short `files/{bucket}/{path}` URL here. The bytes are then fetched
server-side with the caller's own DIAL credentials, so a multi-megabyte PNG
never travels through the agent's context window (which is what
`manage_image(source_type="base64")` costs).

Placement is aspect-ratio aware by default: `fit="contain"` scales the
picture into the box you give it and centres it there, so the visual QA pass
never has to report distortion it cannot repair.
"""
import io
import os
from typing import Dict, Optional

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pptx.util import Emu, Inches

from logging_utils import get_logger
from state import short_id

logger = get_logger("tools.image")

UNKNOWN_ID = (
    "Unknown or expired presentation_id. Pass the presentation_id returned "
    "by create_presentation, create_presentation_from_template, or "
    "open_presentation"
)

FIT_MODES = ("contain", "cover", "stretch")
DEFAULT_MAX_MB = 20.0
EMU_PER_INCH = 914400


def _max_bytes():
    raw = os.environ.get("DIAL_IMAGE_MAX_MB", DEFAULT_MAX_MB)
    try:
        return int(float(raw) * 1024 * 1024)
    except (TypeError, ValueError):
        logger.warning("image_max_mb_invalid value=%s using=%s", raw,
                       DEFAULT_MAX_MB)
        return int(DEFAULT_MAX_MB * 1024 * 1024)


def _inches(emu):
    return round(emu / EMU_PER_INCH, 3)


def _place(pic, left, top, box_w, box_h, fit, slide_w, slide_h):
    """Size and position an already-inserted picture. Returns the fit mode
    actually applied."""
    native_w, native_h = pic.width, pic.height

    if box_w is None and box_h is None:
        # Keep the native size, but never let it hang off the slide.
        scale = min(1.0, slide_w / native_w, slide_h / native_h)
        pic.width = Emu(int(native_w * scale))
        pic.height = Emu(int(native_h * scale))
        return "native" if scale == 1.0 else "native_clamped"

    if box_w is None or box_h is None:
        # One dimension given: python-pptx's own proportional scaling.
        if box_w is None:
            pic.height = Emu(box_h)
            pic.width = Emu(int(native_w * box_h / native_h))
        else:
            pic.width = Emu(box_w)
            pic.height = Emu(int(native_h * box_w / native_w))
        return "proportional"

    if fit == "stretch":
        pic.width, pic.height = Emu(box_w), Emu(box_h)
        return "stretch"

    aspect_img = native_w / native_h
    aspect_box = box_w / box_h

    if fit == "cover":
        # Fill the box exactly and crop the overflowing axis symmetrically.
        if aspect_img > aspect_box:
            trim = (1 - aspect_box / aspect_img) / 2
            pic.crop_left = pic.crop_right = trim
        else:
            trim = (1 - aspect_img / aspect_box) / 2
            pic.crop_top = pic.crop_bottom = trim
        pic.left, pic.top = Emu(left), Emu(top)
        pic.width, pic.height = Emu(box_w), Emu(box_h)
        return "cover"

    # contain: largest undistorted size that fits, centred in the box.
    scale = min(box_w / native_w, box_h / native_h)
    pic.width = Emu(int(native_w * scale))
    pic.height = Emu(int(native_h * scale))
    pic.left = Emu(left + (box_w - pic.width) // 2)
    pic.top = Emu(top + (box_h - pic.height) // 2)
    return "contain"


def register_image_tools(app: FastMCP, presentations):
    """Register the DIAL-hosted image tools with the FastMCP app."""

    @app.tool(
        annotations=ToolAnnotations(
            title="Add Image from DIAL File URL",
        ),
    )
    def add_image_from_dial_url(
        presentation_id: str,
        slide_index: int,
        image_url: str,
        left: float = 1.0,
        top: float = 1.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
        fit: str = "contain",
    ) -> Dict:
        """Place an image that already lives in DIAL file storage onto a slide.

        Use this for generated imagery: call your image model, save the result
        to DIAL files, then pass the file URL here. The server downloads the
        bytes itself, so nothing large passes through your context — prefer
        this over manage_image(source_type="base64") for anything bigger than
        a small icon.

        image_url: a DIAL file reference — the "files/{bucket}/{path}" URL an
        upload returns, or the full https URL of that file on this DIAL
        installation (the ".../api/files/{bucket}/{path}" link an image
        deployment hands back works as-is). Arbitrary web URLs are NOT
        downloaded: to use a picture from the web, fetch it yourself and store
        it in DIAL file storage first.

        slide_index: 0-based, like the other content tools (visual QA slide
        numbers are 1-based; slide_index 0 is slide 1).
        left, top, width, height: inches. Give width and height to define the
        box the picture should occupy — for a half-and-half slide on a 13.33in
        deck, text on the left and left=6.9, top=1.2, width=5.6, height=4.5
        for the image. Omit both to keep the image's natural size (clamped to
        the slide).
        fit: how the picture relates to that box.
          "contain" (default) largest undistorted size that fits, centred in
                    the box — safe for photos and illustrations;
          "cover"   fills the box exactly, cropping the overflowing edges;
          "stretch" forces the exact box and will distort the image.

        Returns the shape index and the geometry actually applied
        ("placed": left/top/width/height in inches), which may be smaller than
        the box you asked for under "contain".
        """
        from dial_client import DialFileClient, DialConfigError

        if presentation_id not in presentations:
            return {"error": UNKNOWN_ID}
        pres = presentations[presentation_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Available "
                             f"slides: 0-{len(pres.slides) - 1}"}
        if fit not in FIT_MODES:
            return {"error": f"Invalid fit: {fit}. Must be one of "
                             f"{', '.join(FIT_MODES)}."}
        if left < 0 or top < 0:
            return {"error": "left and top must be zero or positive inches."}
        for name, value in (("width", width), ("height", height)):
            if value is not None and value <= 0:
                return {"error": f"{name} must be a positive number of inches."}

        try:
            data = DialFileClient().download(image_url)
        except DialConfigError as e:
            logger.error("image_download_failed presentation_id=%s reason="
                         "dial_config error=%s", short_id(presentation_id), e)
            return {"error": str(e)}
        except Exception as e:
            logger.error("image_download_failed presentation_id=%s reason=%s "
                         "error=%s", short_id(presentation_id),
                         type(e).__name__, e)
            return {"error": f"Failed to download the image from DIAL file "
                             f"storage: {str(e)}. Pass the file URL exactly as "
                             f"the upload returned it, e.g. "
                             f"files/{{bucket}}/{{path}}, and check the file "
                             f"still exists."}

        limit = _max_bytes()
        if len(data) > limit:
            logger.warning("image_too_large presentation_id=%s bytes=%d "
                           "limit=%d", short_id(presentation_id), len(data),
                           limit)
            return {"error": f"Image is {len(data) / 1048576:.1f} MB, over the "
                             f"{limit / 1048576:.1f} MB limit. Generate or "
                             f"store a smaller image."}

        slide = pres.slides[slide_index]
        try:
            pic = slide.shapes.add_picture(io.BytesIO(data), Inches(left),
                                           Inches(top))
        except Exception as e:
            logger.error("image_insert_failed presentation_id=%s slide=%d "
                         "reason=%s error=%s", short_id(presentation_id),
                         slide_index, type(e).__name__, e)
            return {"error": f"The downloaded file is not an image PowerPoint "
                             f"can embed ({str(e)}). PNG, JPEG, GIF, BMP and "
                             f"TIFF work; SVG does not — ask the image model "
                             f"for a raster format."}

        native = {"width_px": pic.image.size[0], "height_px": pic.image.size[1]}
        applied = _place(
            pic,
            int(Inches(left)), int(Inches(top)),
            int(Inches(width)) if width is not None else None,
            int(Inches(height)) if height is not None else None,
            fit,
            int(pres.slide_width), int(pres.slide_height),
        )

        logger.info("image_added presentation_id=%s slide=%d bytes=%d fit=%s",
                    short_id(presentation_id), slide_index, len(data), applied)
        return {
            "message": f"Added image to slide {slide_index} ({applied}).",
            "shape_index": len(slide.shapes) - 1,
            "fit": applied,
            "native": native,
            "placed": {
                "left": _inches(pic.left),
                "top": _inches(pic.top),
                "width": _inches(pic.width),
                "height": _inches(pic.height),
            },
            "size_bytes": len(data),
        }
