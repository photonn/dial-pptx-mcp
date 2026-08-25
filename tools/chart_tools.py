"""
Chart data management tools for PowerPoint MCP Server.
Implements advanced chart data manipulation capabilities.
"""

from typing import Dict, List, Optional, Any
from mcp.types import ToolAnnotations
from pptx.chart.data import ChartData

from logging_utils import get_logger

logger = get_logger("tools.chart")


def _base_chart_types():
    """Which single-type chart to build before the series are regrouped.

    Any of them would do — every series is moved afterwards — but starting
    from the first series' own type leaves a chart whose series happen to
    share a type exactly as python-pptx built it. Imported lazily, like the
    rest of the pptx enums in this module.
    """
    from pptx.enum.chart import XL_CHART_TYPE
    return {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "stacked_area": XL_CHART_TYPE.AREA_STACKED,
    }


def _validate_combo_series(series, categories):
    """Check a combo chart's series list; return an error string or None."""
    import utils as ppt_utils

    if not categories:
        return "categories must not be empty"
    if not series:
        return "series must not be empty"
    if len(series) < 2:
        return ("A combo chart needs at least two series — with one, use "
                "add_chart instead.")

    known = ", ".join(sorted(ppt_utils.COMBO_SERIES_TYPES))
    for index, entry in enumerate(series):
        if not isinstance(entry, dict):
            return f"series[{index}] must be an object"
        for key in ("name", "values", "type"):
            if key not in entry:
                return f"series[{index}] is missing '{key}'"
        if entry["type"] not in ppt_utils.COMBO_SERIES_TYPES:
            return (f"series[{index}] has type '{entry['type']}', which "
                    f"cannot share a plot area. Use one of: {known}. Pie and "
                    f"scatter charts must be their own chart (add_chart).")
        values = entry["values"]
        if not isinstance(values, list) or len(values) != len(categories):
            return (f"series[{index}] ('{entry['name']}') has "
                    f"{len(values) if isinstance(values, list) else 'non-list'} "
                    f"value(s) for {len(categories)} categories; every series "
                    f"must supply one value per category (use null for a gap).")
    return None


def _format_combo_chart(chart, layout, title, has_legend, legend_position,
                        x_axis_title, y_axis_title, secondary_axis_title):
    """Chart-level formatting, including the secondary axis python-pptx does
    not expose as an object.

    Both value axes are addressed by the ids apply_combo_layout reports:
    chart.value_axis returns the second valAx when there are two, so using it
    here would put the left axis' title on the right-hand axis.
    """
    import utils as ppt_utils
    from pptx.enum.chart import XL_LEGEND_POSITION

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    chart.has_legend = bool(has_legend)
    if has_legend:
        position = getattr(XL_LEGEND_POSITION, (legend_position or "bottom").upper(),
                           XL_LEGEND_POSITION.BOTTOM)
        chart.legend.position = position
        chart.legend.include_in_layout = False

    if x_axis_title:
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.text = x_axis_title
    _title_axis(ppt_utils.value_axis_with_id(chart, layout["primary"][1]),
                y_axis_title)
    if layout["secondary"]:
        _title_axis(ppt_utils.value_axis_with_id(chart, layout["secondary"][1]),
                    secondary_axis_title)
    elif secondary_axis_title:
        logger.debug("secondary_axis_title_skipped reason=no_secondary_axis")


def _title_axis(axis, text):
    if axis is None or not text:
        return
    axis.has_title = True
    axis.axis_title.text_frame.text = text


def _apply_series_options(chart, series):
    """Per-series colour, labels and trendlines, as given to add_combo_chart."""
    import utils as ppt_utils

    ordered = [s for plot in chart.plots for s in plot.series]
    by_name = {s.name: s for s in ordered}
    for entry in series:
        target = by_name.get(entry["name"])
        if target is None:
            continue
        try:
            if entry.get("color"):
                ppt_utils.set_series_color(target, tuple(entry["color"]))
            if entry.get("data_labels") or entry.get("label_position"):
                ppt_utils.set_series_data_labels(
                    target, entry.get("data_labels", True),
                    entry.get("label_position"), entry.get("number_format"))
            if entry.get("trendline"):
                ppt_utils.add_series_trendline(target, entry["trendline"],
                                               entry.get("trendline_period"))
        except Exception as e:
            # One series' styling is not worth losing the chart over.
            logger.warning("series_option_skipped series=%s error=%s",
                           entry.get("name"), e)


def register_chart_tools(app, presentations, get_current_presentation_id, validate_parameters, 
                          is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register chart data management tools with the FastMCP app."""
    
    @app.tool(
        annotations=ToolAnnotations(
            title="Update Chart Data",
        ),
    )
    def update_chart_data(
        slide_index: int,
        shape_index: int,
        categories: List[str],
        series_data: List[Dict],
        presentation_id: str = None
    ) -> Dict:
        """
        Replace existing chart data with new categories and series.
        
        Args:
            slide_index: Index of the slide (0-based)
            shape_index: Index of the chart shape (0-based)
            categories: List of category names
            series_data: List of dictionaries with 'name' and 'values' keys
            presentation_id: Optional presentation ID (uses current if not provided)
            
        Returns:
            Dictionary with operation results
        """
        try:
            # Get presentation
            pres_id = presentation_id or get_current_presentation_id()
            if pres_id not in presentations:
                return {"error": "Presentation not found"}
            
            pres = presentations[pres_id]
            
            # Validate slide index
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Slide index {slide_index} out of range"}
            
            slide = pres.slides[slide_index]
            
            # Validate shape index
            if not (0 <= shape_index < len(slide.shapes)):
                return {"error": f"Shape index {shape_index} out of range"}
            
            shape = slide.shapes[shape_index]
            
            # Check if shape is a chart
            if not hasattr(shape, 'has_chart') or not shape.has_chart:
                return {"error": "Shape is not a chart"}
            
            chart = shape.chart
            
            # Create new ChartData
            chart_data = ChartData()
            chart_data.categories = categories
            
            # Add series data
            for series in series_data:
                if 'name' not in series or 'values' not in series:
                    return {"error": "Each series must have 'name' and 'values' keys"}
                
                chart_data.add_series(series['name'], series['values'])
            
            # Replace chart data
            chart.replace_data(chart_data)
            logger.debug("chart_data_replaced slide=%s shape=%s categories=%d "
                         "series=%d", slide_index, shape_index, len(categories),
                         len(series_data))
            
            return {
                "message": f"Updated chart data on slide {slide_index}, shape {shape_index}",
                "categories": categories,
                "series_count": len(series_data),
                "series_names": [s['name'] for s in series_data]
            }
            
        except Exception as e:
            return {"error": f"Failed to update chart data: {str(e)}"}
    @app.tool(
        annotations=ToolAnnotations(
            title="Add Combo Chart",
        ),
    )
    def add_combo_chart(
        slide_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        categories: List[str],
        series: List[Dict],
        title: Optional[str] = None,
        has_legend: bool = True,
        legend_position: str = "bottom",
        x_axis_title: Optional[str] = None,
        y_axis_title: Optional[str] = None,
        secondary_axis_title: Optional[str] = None,
        presentation_id: str = None
    ) -> Dict:
        """Add a chart whose series are not all the same type, and may use a
        second value axis.

        This is what add_chart cannot do. Reach for it when:
        - a measure needs a reference line across it — actuals as columns with
          a target as a line;
        - two measures share categories but not units — revenue in millions
          against a margin percentage. On one axis the percentage flattens to
          nothing along the bottom; on a secondary axis both are readable.

        series: one entry per series, plotted in the order given:
          {"name": "Revenue", "values": [10, 12, 15],
           "type": "column", "secondary_axis": false,
           "color": [31, 73, 125], "data_labels": true,
           "label_position": "outside_end", "trendline": "linear"}
          "type" is required — column, stacked_column, bar, stacked_bar, line,
          line_markers, area, stacked_area. Pie and scatter cannot share a plot
          area with other types; use add_chart for those.
          "color" is [r, g, b]; omit to take the template's own chart palette,
          which is usually the right choice.
          "label_position" on a stacked series must be "center", "inside_end"
          or "inside_base" — a stacked segment has no outside.
          "trendline" is linear, movingAvg, exp, log, poly or power.

        secondary_axis_title labels the right-hand axis; y_axis_title labels
        the left one. x_axis_title labels the category axis (see add_chart on
        why that is not the same as "the horizontal one").

        The result is a real, editable PowerPoint chart, not a picture.
        """
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "Unknown or expired presentation_id. Pass the presentation_id returned by create_presentation, create_presentation_from_template, or open_presentation"
            }

        pres = presentations[pres_id]
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        problem = _validate_combo_series(series, categories)
        if problem:
            return {"error": problem}

        valid, error = validate_parameters({
            'left': (left, [(is_non_negative, "must be non-negative")]),
            'top': (top, [(is_non_negative, "must be non-negative")]),
            'width': (width, [(is_positive, "must be positive")]),
            'height': (height, [(is_positive, "must be positive")]),
        })
        if not valid:
            return {"error": error}

        from pptx.util import Inches
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        import utils as ppt_utils

        slide = pres.slides[slide_index]
        try:
            # Build the whole thing as one ordinary chart first: that is what
            # writes a correct embedded workbook and the category caches every
            # series shares. The series are only regrouped afterwards.
            data = CategoryChartData()
            data.categories = categories
            for entry in series:
                data.add_series(entry["name"], tuple(entry["values"]))

            base = _base_chart_types().get(series[0]["type"],
                                           XL_CHART_TYPE.COLUMN_CLUSTERED)
            graphic_frame = slide.shapes.add_chart(
                base, Inches(left), Inches(top), Inches(width),
                Inches(height), data)
            chart = graphic_frame.chart

            layout = ppt_utils.apply_combo_layout(
                chart, [{"type": e["type"],
                         "secondary_axis": bool(e.get("secondary_axis"))}
                        for e in series])

            _format_combo_chart(chart, layout, title, has_legend,
                                legend_position, x_axis_title, y_axis_title,
                                secondary_axis_title)
            _apply_series_options(chart, series)

            shape_index = len(slide.shapes) - 1
            logger.info("combo_chart_added slide=%d shape=%d series=%d "
                        "groups=%d secondary=%s", slide_index, shape_index,
                        len(series), layout["groups"],
                        layout["secondary"] is not None)
            return {
                "message": f"Added a combo chart to slide {slide_index}",
                "shape_index": shape_index,
                "plot_groups": layout["groups"],
                "series_names": [e["name"] for e in series],
                "has_secondary_axis": layout["secondary"] is not None,
            }
        except Exception as e:
            logger.error("combo_chart_failed slide=%d error=%s", slide_index, e)
            return {"error": f"Failed to add combo chart: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Format Chart Series",
        ),
    )
    def format_chart_series(
        slide_index: int,
        shape_index: int,
        series_index: int,
        color: Optional[List[int]] = None,
        data_labels: Optional[bool] = None,
        label_position: Optional[str] = None,
        number_format: Optional[str] = None,
        trendline: Optional[str] = None,
        trendline_period: Optional[int] = None,
        presentation_id: str = None
    ) -> Dict:
        """Restyle one series of an existing chart.

        Use it to pin a series to a brand colour, to label just the series
        that carries the point rather than all of them, or to add a trendline.

        series_index is 0-based across the whole chart, in the order the
        series were added.
        label_position: center, inside_end, inside_base, outside_end (bars and
          columns), or above / below / left / right (line charts). A stacked
          series rejects outside_end.
        number_format: an Excel format string, e.g. "0.0%" or "#,##0".
        trendline: linear, movingAvg, exp, log, poly or power;
          trendline_period is the window for movingAvg or the order for poly.
        """
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "Unknown or expired presentation_id. Pass the presentation_id returned by create_presentation, create_presentation_from_template, or open_presentation"
            }

        pres = presentations[pres_id]
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        slide = pres.slides[slide_index]
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Slide {slide_index} has {len(slide.shapes)} shape(s)"
            }
        shape = slide.shapes[shape_index]
        if not getattr(shape, "has_chart", False):
            return {"error": f"Shape {shape_index} on slide {slide_index} is "
                             f"not a chart"}

        all_series = [s for plot in shape.chart.plots for s in plot.series]
        if series_index < 0 or series_index >= len(all_series):
            return {"error": f"Invalid series_index: {series_index}. This "
                             f"chart has {len(all_series)} series, indexed "
                             f"0-{len(all_series) - 1}"}
        if color is not None and not is_valid_rgb(color):
            return {"error": "color must be [r, g, b] with each value 0-255"}

        import utils as ppt_utils
        if trendline is not None and trendline not in ppt_utils.TRENDLINE_TYPES:
            return {"error": f"Invalid trendline: {trendline}. Must be one of "
                             f"{', '.join(ppt_utils.TRENDLINE_TYPES)}."}

        target = all_series[series_index]
        applied = []
        try:
            if color is not None:
                ppt_utils.set_series_color(target, tuple(color))
                applied.append("color")
            if data_labels is not None or label_position or number_format:
                ppt_utils.set_series_data_labels(
                    target,
                    True if data_labels is None else data_labels,
                    label_position, number_format)
                applied.append("data_labels")
            if trendline:
                ppt_utils.add_series_trendline(target, trendline,
                                               trendline_period)
                applied.append("trendline")
        except (AttributeError, ValueError) as e:
            return {"error": f"Could not apply that formatting to series "
                             f"{series_index}: {e}. A label position must "
                             f"suit the chart type — a stacked or line series "
                             f"rejects some positions."}
        except Exception as e:
            logger.error("series_format_failed slide=%d shape=%d series=%d "
                         "error=%s", slide_index, shape_index, series_index, e)
            return {"error": f"Failed to format series {series_index}: {e}"}

        if not applied:
            return {"error": "Nothing to do: pass at least one of color, "
                             "data_labels, label_position, number_format or "
                             "trendline."}
        return {
            "message": f"Formatted series {series_index} ('{target.name}') on "
                       f"slide {slide_index}: {', '.join(applied)}.",
            "series_index": series_index,
            "applied": applied,
        }
