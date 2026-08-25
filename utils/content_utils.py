"""
Content management utilities for PowerPoint MCP Server.
Functions for slides, text, images, tables, charts, and shapes.
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import os
import base64

from logging_utils import get_logger, flatten

logger = get_logger("utils.content")


def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use
        
    Returns:
        A tuple containing the slide and its layout
    """
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    return slide, layout


def get_slide_info(slide, slide_index: int) -> Dict:
    """
    Get information about a specific slide.
    
    Args:
        slide: The slide object
        slide_index: Index of the slide
        
    Returns:
        Dictionary containing slide information
    """
    try:
        placeholders = []
        for placeholder in slide.placeholders:
            placeholder_info = {
                "idx": placeholder.placeholder_format.idx,
                "type": str(placeholder.placeholder_format.type),
                "name": placeholder.name
            }
            placeholders.append(placeholder_info)
        
        shapes = []
        for i, shape in enumerate(slide.shapes):
            shape_info = {
                "index": i,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            shapes.append(shape_info)
        
        return {
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name,
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
            "shape_count": len(shapes),
            "shapes": shapes
        }
    except Exception as e:
        raise Exception(f"Failed to get slide info: {str(e)}")


def set_title(slide, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
    """
    if slide.shapes.title:
        slide.shapes.title.text = title


def populate_placeholder(slide, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
    """
    placeholder = slide.placeholders[placeholder_idx]
    placeholder.text = text


def add_bullet_points(placeholder, bullet_points: List[str]) -> None:
    """
    Add bullet points to a placeholder.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts
    """
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    for i, point in enumerate(bullet_points):
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int = None, font_name: str = None, bold: bool = None,
                italic: bool = None, underline: bool = None, 
                color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                alignment: str = None, vertical_alignment: str = None, 
                auto_fit: bool = True) -> Any:
    """
    Add a textbox to a slide with formatting options.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
        auto_fit: Whether to auto-fit text
        
    Returns:
        The created textbox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    textbox.text_frame.text = text
    
    # Apply formatting if provided
    if any([font_size, font_name, bold, italic, underline, color, bg_color, alignment, vertical_alignment]):
        format_text_advanced(
            textbox.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            underline=underline,
            color=color,
            bg_color=bg_color,
            alignment=alignment,
            vertical_alignment=vertical_alignment
        )
    
    return textbox


def format_text(text_frame, font_size: int = None, font_name: str = None, 
                bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None) -> None:
    """
    Format text in a text frame.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            if font_size is not None:
                font.size = Pt(font_size)
            if font_name is not None:
                font.name = font_name
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)


def format_text_advanced(text_frame, font_size: int = None, font_name: str = None, 
                        bold: bool = None, italic: bool = None, underline: bool = None,
                        color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                        alignment: str = None, vertical_alignment: str = None) -> Dict:
    """
    Advanced text formatting with comprehensive options.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
    
    Returns:
        Dictionary with formatting results
    """
    result = {
        'success': True,
        'warnings': []
    }
    
    try:
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        # Enable text wrapping
        text_frame.word_wrap = True
        
        # Apply formatting to all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            if alignment and alignment in alignment_map:
                paragraph.alignment = alignment_map[alignment]
            
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                if font_name is not None:
                    font.name = font_name
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['error'] = str(e)
        return result


def add_image(slide, image_path: str, left: float, top: float, width: float = None, height: float = None) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created image shape
    """
    if width is not None and height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width)
        )
    elif height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top)
        )


def add_table(slide, rows: int, cols: int, left: float, top: float, width: float, height: float) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table shape
    """
    return slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def format_table_cell(cell, font_size: int = None, font_name: str = None, 
                     bold: bool = None, italic: bool = None, 
                     color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                     alignment: str = None, vertical_alignment: str = None) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell object
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment
        vertical_alignment: Vertical alignment
    """
    # Format text
    if any([font_size, font_name, bold, italic, color, alignment]):
        format_text_advanced(
            cell.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            alignment=alignment
        )
    
    # Set background color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


# Chart-group elements that colour each data point separately; every other
# plot type gets one colour per series.
_PER_POINT_COLOUR_PLOTS = ('pieChart', 'pie3DChart', 'ofPieChart',
                           'doughnutChart')

_LEGEND_POSITIONS = {
    'right': XL_LEGEND_POSITION.RIGHT,
    'left': XL_LEGEND_POSITION.LEFT,
    'top': XL_LEGEND_POSITION.TOP,
    'bottom': XL_LEGEND_POSITION.BOTTOM,
    'corner': XL_LEGEND_POSITION.CORNER,
}

_CHART_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'


def _c(tag: str) -> str:
    return '{%s}%s' % (_CHART_NS, tag)


def normalize_chart_defaults(chart) -> None:
    """Make explicit the chart settings whose XML default is renderer-dependent.

    python-pptx leaves several optional ``c:`` elements out of the charts it
    generates. Where ECMA-376 says a missing boolean means *true*, PowerPoint
    applies it and LibreOffice does not, so one file renders as two different
    charts:

    * no ``c:varyColors`` — PowerPoint colours a single-series bar chart one
      colour per *category* and lists the categories in the legend, where
      LibreOffice draws a single series colour and a single legend entry.
    * a bare ``c:legend`` — PowerPoint reads the absent ``c:overlay`` as true
      and lays the legend over the plot; LibreOffice reserves space beside it.

    Writing both out is what makes the two agree, and it also means visual QA
    (which renders through LibreOffice) is judging what PowerPoint will show.
    """
    plot_area = chart._chartSpace.find(_c('chart')).find(_c('plotArea'))
    for group in plot_area:
        tag = group.tag.split('}')[-1]
        if not tag.endswith('Chart'):
            continue
        if group.find(_c('varyColors')) is not None:
            continue
        vary = '1' if tag in _PER_POINT_COLOUR_PLOTS else '0'
        if hasattr(group, 'get_or_add_varyColors'):
            group.get_or_add_varyColors().val = vary == '1'
        else:
            # varyColors sits immediately before the series in every group.
            element = group.makeelement(_c('varyColors'), {'val': vary})
            first_ser = group.find(_c('ser'))
            if first_ser is not None:
                first_ser.addprevious(element)
            else:
                group.insert(0, element)

    if chart.has_legend:
        legend = chart.legend
        if legend._element.find(_c('legendPos')) is None:
            legend.position = XL_LEGEND_POSITION.RIGHT
        # The one that matters: absent, PowerPoint overlays the legend.
        legend.include_in_layout = False

    # Missing c:plotVisOnly also defaults to true; say so rather than rely on it.
    chart_el = chart._chartSpace.find(_c('chart'))
    if chart_el.find(_c('plotVisOnly')) is None:
        element = chart_el.makeelement(_c('plotVisOnly'), {'val': '1'})
        chart_el.insert_element_before(element, 'c:dispBlanksAs',
                                       'c:showDLblsOverMax', 'c:extLst')


# Chart types plotted from x/y pairs rather than from categories. They need a
# different chart-data class, and their x values have to be numbers.
SCATTER_CHART_TYPES = frozenset({'scatter'})


def parse_scatter_x_values(categories) -> List[float]:
    """Read a scatter chart's x values out of the `categories` argument.

    A scatter chart has no categories: both axes are numeric. The tool takes
    one list of labels for every chart type, so on a scatter chart that list
    is the x values and has to parse as numbers. Raises ValueError naming what
    to send instead — the caller turns that into the tool's error string.
    """
    x_values = []
    for value in categories:
        try:
            x_values.append(float(value))
        except (TypeError, ValueError):
            raise ValueError(
                f"A scatter chart plots numbers against numbers, but "
                f"categories contains '{value}'. Pass the x values in "
                f"categories (as numbers) and the matching y values in "
                f"series_values, or use 'line_markers' if the horizontal "
                f"axis really is a list of labels.")
    return x_values


def add_chart(slide, chart_type: str, left: float, top: float, width: float, height: float,
              categories: List[str], series_names: List[str], series_values: List[List[float]]) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of value lists for each series
        
    Returns:
        The created chart object
    """
    # Map chart type names to enum values
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'line': XL_CHART_TYPE.LINE,
        'line_markers': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
        'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
    }
    
    chart_type = chart_type.lower()
    xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data. A scatter chart carries its own x values per point
    # instead of a shared category axis, and CategoryChartData cannot write
    # the c:xVal element its series need.
    if chart_type in SCATTER_CHART_TYPES:
        x_values = parse_scatter_x_values(categories)
        chart_data = XyChartData()
        for i, series_name in enumerate(series_names):
            if i >= len(series_values):
                continue
            series = chart_data.add_series(series_name)
            for x, y in zip(x_values, series_values[i]):
                series.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for i, series_name in enumerate(series_names):
            if i < len(series_values):
                chart_data.add_series(series_name, series_values[i])
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    
    chart = chart_shape.chart
    normalize_chart_defaults(chart)
    return chart


def format_chart(chart, has_legend: bool = True, legend_position: str = 'right',
                has_data_labels: bool = False, title: str = None,
                x_axis_title: str = None, y_axis_title: str = None,
                color_scheme: str = None) -> None:
    """
    Format a chart with various options.
    
    Args:
        chart: The chart object
        has_legend: Whether to show legend
        legend_position: Position of legend ('right', 'top', 'bottom', 'left')
        has_data_labels: Whether to show data labels
        title: Chart title
        x_axis_title: Title for the CATEGORY axis (horizontal on column/
            line charts, vertical on bar charts — see the mapping below)
        y_axis_title: Title for the VALUE axis (vertical on column/line
            charts, horizontal on bar charts)
        color_scheme: Color scheme to apply
    """
    try:
        # Set chart title. An absent title is stated as autoTitleDeleted rather
        # than left out: PowerPoint auto-titles a single-series chart from the
        # series name, LibreOffice does not, so silence is two different charts.
        if title:
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend. legendPos and overlay are both written out — see
        # normalize_chart_defaults for why leaving them implicit diverges.
        chart.has_legend = bool(has_legend)
        if has_legend:
            position = _LEGEND_POSITIONS.get((legend_position or 'right').lower())
            if position is not None:
                chart.legend.position = position
            chart.legend.include_in_layout = False
        
        # Configure data labels
        if has_data_labels:
            for series in chart.series:
                series.has_data_labels = True
        
        # Set axis titles if available
        try:
            if x_axis_title and hasattr(chart, 'category_axis'):
                chart.category_axis.axis_title.text_frame.text = x_axis_title
            if y_axis_title and hasattr(chart, 'value_axis'):
                chart.value_axis.axis_title.text_frame.text = y_axis_title
        except:
            pass  # Axis titles may not be available for all chart types
            
    except Exception:
        pass  # Graceful degradation for chart formatting


def extract_slide_text_content(slide) -> Dict:
    """
    Extract all text content from a slide including placeholders and text shapes.
    
    Args:
        slide: The slide object to extract text from
        
    Returns:
        Dictionary containing all text content organized by source type
    """
    try:
        text_content = {
            "slide_title": "",
            "placeholders": [],
            "text_shapes": [],
            "table_text": [],
            "speaker_notes": "",
            "all_text_combined": ""
        }
        
        all_texts = []
        
        # Extract title from slide if available
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            try:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    text_content["slide_title"] = title_text
                    all_texts.append(title_text)
            except Exception as e:
                logger.debug("title_extract_skipped error=%s", flatten(str(e)))
        
        # Extract text from all shapes
        for i, shape in enumerate(slide.shapes):
            shape_text_info = {
                "shape_index": i,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "text": ""
            }
            
            try:
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shape_text_info["text"] = text
                        all_texts.append(text)
                        
                        # Categorize by shape type
                        if hasattr(shape, 'placeholder_format'):
                            # This is a placeholder
                            placeholder_info = shape_text_info.copy()
                            placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                            placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
                            text_content["placeholders"].append(placeholder_info)
                        else:
                            # This is a regular text shape
                            text_content["text_shapes"].append(shape_text_info)
                
                # Extract text from tables
                elif hasattr(shape, 'table'):
                    table_texts = []
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_texts = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                                all_texts.append(cell_text)
                        if row_texts:
                            table_texts.append({
                                "row": row_idx,
                                "cells": row_texts
                            })
                    
                    if table_texts:
                        text_content["table_text"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "table_content": table_texts
                        })
                        
            except Exception as e:
                # Skip shapes that can't be processed
                continue
        
        # Speaker notes are kept out of all_text_combined: that string is what
        # the caller reads to check what the audience sees, and notes are not
        # on the slide.
        try:
            if slide.has_notes_slide:
                text_content["speaker_notes"] = \
                    slide.notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            logger.debug("notes_extract_skipped error=%s", flatten(str(e)))

        # Combine all text
        text_content["all_text_combined"] = "\n".join(all_texts)
        
        return {
            "success": True,
            "text_content": text_content,
            "total_text_shapes": len(text_content["placeholders"]) + len(text_content["text_shapes"]),
            "has_title": bool(text_content["slide_title"]),
            "has_tables": len(text_content["table_text"]) > 0,
            "has_notes": bool(text_content["speaker_notes"])
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to extract text content: {str(e)}",
            "text_content": None
        }