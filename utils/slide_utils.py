"""
Slide-level structural operations: delete, reorder, duplicate, copy between decks.

python-pptx has no API for any of these. Its only entry point is
``slides.add_slide(layout)``, which appends a blank slide built from a layout —
so a template's fully designed slide cannot be repeated, an unwanted slide
cannot be removed, and slide order is fixed at creation time. Everything here
works at the OPC package level (the ``<p:sldIdLst>`` in presentation.xml plus
the slide part's relationships), which is where those operations actually live.

Part-sharing policy for duplication: image, media and font parts are immutable
binaries, so copies reference the same part (that is what PowerPoint itself
does). Charts, SmartArt data and embedded objects are *edited* through this
server, so they are cloned — otherwise ``update_chart_data`` on the copy would
silently rewrite the original's chart too.
"""
import copy
import re

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import _Relationship  # noqa: PLC2701 - no public equivalent
from pptx.oxml.ns import qn

from logging_utils import get_logger

logger = get_logger("utils.slide")

__all__ = [
    "check_index",
    "delete_slide",
    "move_slide",
    "duplicate_slide",
    "copy_slide_to_presentation",
    "get_speaker_notes",
    "set_speaker_notes",
    "pin_inherited_geometry",
]

# Relationship targets shared (referenced) rather than cloned by a duplicate.
# Immutable payloads, or parts that live above the slide in the package.
_SHARED_RELTYPES = frozenset({
    RT.IMAGE,
    RT.MEDIA,
    RT.VIDEO,
    RT.AUDIO,
    RT.FONT,
    RT.SLIDE_LAYOUT,
    RT.SLIDE_MASTER,
    RT.THEME,
    RT.SLIDE,
})

# A notes slide belongs to exactly one slide; its rel is never copied verbatim
# (two slides pointing at one notesSlide part is an invalid package). The notes
# text is transferred separately.
_SKIPPED_RELTYPES = frozenset({RT.NOTES_SLIDE})


def _sld_id_lst(pres):
    return pres.slides._sldIdLst


def _slide_count(pres):
    return len(pres.slides)


def check_index(pres, index, what="slide_index"):
    """Return an error string for an out-of-range 0-based index, else None."""
    total = _slide_count(pres)
    if not isinstance(index, int) or isinstance(index, bool):
        return f"{what} must be an integer"
    if total == 0:
        return "This presentation has no slides"
    if index < 0 or index >= total:
        return (f"Invalid {what}: {index}. This presentation has {total} "
                f"slide(s), indexed 0-{total - 1}")
    return None


# ---- Delete / reorder ----

def delete_slide(pres, index):
    """Remove the slide at `index` from the presentation.

    Drops the presentation's relationship to the slide part as well as its
    <p:sldId> entry, so the part is garbage-collected on save instead of
    lingering as an orphan in the package.
    """
    sld_id_lst = _sld_id_lst(pres)
    sld_id = list(sld_id_lst)[index]
    rId = sld_id.rId
    sld_id_lst.remove(sld_id)
    pres.part.drop_rel(rId)


def move_slide(pres, index, new_index):
    """Move the slide at `index` so that it ends up at position `new_index`."""
    sld_id_lst = _sld_id_lst(pres)
    entries = list(sld_id_lst)
    sld_id = entries[index]
    sld_id_lst.remove(sld_id)
    remaining = list(sld_id_lst)
    if new_index >= len(remaining):
        sld_id_lst.append(sld_id)
    else:
        remaining[new_index].addprevious(sld_id)


def _move_to(pres, from_index, to_index):
    if from_index != to_index:
        move_slide(pres, from_index, to_index)


# ---- Part cloning ----

def _partname_template(partname):
    """"/ppt/charts/chart1.xml" -> "/ppt/charts/chart%d.xml"."""
    name = str(partname)
    if re.search(r"\d+(?=\.[^./]+$)", name):
        return re.sub(r"\d+(?=\.[^./]+$)", "%d", name)
    return re.sub(r"(\.[^./]+)$", r"%d\1", name)


def _set_rel(part, rId, reltype, target, external=False):
    """Add a relationship under a caller-chosen rId.

    Duplicated slide XML still carries the source's r:embed / r:id attributes,
    so the copy's relationships have to keep the same ids. The public
    `rels.get_or_add` allocates its own rId, which would not match.
    """
    part.rels._rels[rId] = _Relationship(
        part.partname.baseURI, rId, reltype,
        target_mode="External" if external else "Internal", target=target)


def _clone_part(package, part, memo):
    """Deep-copy a package part and everything it references (except shared
    payloads), returning the new part."""
    key = id(part)
    if key in memo:
        return memo[key]

    partname = package.next_partname(_partname_template(part.partname))
    new_part = type(part).load(partname, part.content_type, package, part.blob)
    memo[key] = new_part

    for rel in part.rels.values():
        if rel.is_external:
            _set_rel(new_part, rel.rId, rel.reltype, rel.target_ref, external=True)
        elif rel.reltype in _SHARED_RELTYPES:
            _set_rel(new_part, rel.rId, rel.reltype, rel.target_part)
        else:
            _set_rel(new_part, rel.rId, rel.reltype,
                     _clone_part(package, rel.target_part, memo))
    return new_part


def _copy_rels(source_part, dest_part, package, clone_all=False):
    """Rebuild dest_part's relationships from source_part's, preserving rIds.

    `clone_all` copies even the normally shared payloads — required when the
    destination lives in a different package, where the source's parts are not
    reachable.
    """
    dest_part.rels._rels.clear()
    memo = {}
    for rel in source_part.rels.values():
        if rel.reltype in _SKIPPED_RELTYPES:
            continue
        if rel.is_external:
            _set_rel(dest_part, rel.rId, rel.reltype, rel.target_ref, external=True)
            continue
        shared = rel.reltype in _SHARED_RELTYPES and not clone_all
        target = (rel.target_part if shared
                  else _clone_part(package, rel.target_part, memo))
        _set_rel(dest_part, rel.rId, rel.reltype, target)
    return memo


_CSLD = qn("p:cSld")
_SP_TREE = qn("p:spTree")
# The group properties every spTree opens with; they describe the tree itself,
# not a shape, and are the two children a cleared tree keeps.
_TREE_PROPS = frozenset({qn("p:nvGrpSpPr"), qn("p:grpSpPr")})


def _replace_children(dest_el, source_el, keep=frozenset()):
    for child in list(dest_el):
        if child.tag not in keep:
            dest_el.remove(child)
    for child in source_el:
        if child.tag not in keep:
            dest_el.append(copy.deepcopy(child))


def _replace_slide_xml(dest_slide, source_slide):
    """Overwrite the destination slide's XML with a copy of the source's.

    <p:cSld> and its <p:spTree> are emptied and refilled rather than replaced:
    python-pptx binds a slide's `shapes` collection to the spTree *element* the
    first time it is touched, and add_slide touches it while cloning layout
    placeholders. Swapping the element out would leave that collection — and so
    every later edit to this slide — attached to a tree no longer in the deck.
    """
    src_sld, dst_sld = source_slide._element, dest_slide._element
    src_csld, dst_csld = src_sld.find(_CSLD), dst_sld.find(_CSLD)

    # Siblings of <p:cSld>: colour-map override, transition, timing.
    _replace_children(dst_sld, src_sld, keep={_CSLD})

    name = src_csld.get("name")
    if name is None:
        dst_csld.attrib.pop("name", None)
    else:
        dst_csld.set("name", name)

    # cSld's own children other than the shape tree — background first
    # (schema order is bg, spTree, custDataLst, controls, extLst).
    dst_tree = dst_csld.find(_SP_TREE)
    for child in list(dst_csld):
        if child.tag != _SP_TREE:
            dst_csld.remove(child)
    seen_tree = False
    for child in src_csld:
        if child.tag == _SP_TREE:
            seen_tree = True
            continue
        clone = copy.deepcopy(child)
        if seen_tree:
            dst_csld.append(clone)
        else:
            dst_tree.addprevious(clone)

    _replace_children(dst_tree, src_csld.find(_SP_TREE), keep=_TREE_PROPS)


def _copy_notes(source_slide, dest_slide):
    """Transfer speaker notes, if the source has any, into a fresh notes slide."""
    if not source_slide.has_notes_slide:
        return
    text = source_slide.notes_slide.notes_text_frame.text
    if text:
        dest_slide.notes_slide.notes_text_frame.text = text


# ---- Duplicate / copy ----

def duplicate_slide(pres, index, insert_after=None):
    """Duplicate the slide at `index`; return the new slide's 0-based index.

    `insert_after` is the 0-based index the copy is placed after; None appends
    it to the end of the deck.
    """
    source = pres.slides[index]
    dest = pres.slides.add_slide(source.slide_layout)

    _replace_slide_xml(dest, source)
    _copy_rels(source.part, dest.part, pres.part.package)
    _copy_notes(source, dest)

    new_index = len(pres.slides) - 1
    if insert_after is not None:
        target = insert_after + 1 if insert_after < new_index else new_index
        _move_to(pres, new_index, target)
        new_index = target
    return new_index


def _matching_layout(target_pres, layout):
    """Find a layout in the target deck equivalent to `layout`.

    Name first (corporate templates name their layouts, and the same template
    on both sides is the common case), then the layout's placeholder-type
    signature, so a copy still lands on something structurally comparable.
    """
    name = getattr(layout, "name", None)
    if name:
        for candidate in target_pres.slide_layouts:
            if getattr(candidate, "name", None) == name:
                return candidate

    def signature(lay):
        return tuple(sorted(str(ph.placeholder_format.type)
                            for ph in lay.placeholders))

    wanted = signature(layout)
    for candidate in target_pres.slide_layouts:
        if signature(candidate) == wanted:
            return candidate
    return None


def copy_slide_to_presentation(source_pres, index, target_pres,
                               insert_after=None, layout_index=None):
    """Copy a slide from one presentation into another.

    Returns ``(new_index, layout_name, layout_matched)``. Everything the slide
    references is cloned into the target package, since the source package's
    parts are not reachable from it. The slide keeps its own formatting, but
    inherited theme colours and fonts resolve against the *target* deck's
    master — which is normally what merging decks is for, and the reason the
    result is worth a visual check.
    """
    source = source_pres.slides[index]

    if layout_index is not None:
        layout = target_pres.slide_layouts[layout_index]
        matched = True
    else:
        layout = _matching_layout(target_pres, source.slide_layout)
        matched = layout is not None
        if layout is None:
            layout = target_pres.slide_layouts[
                min(1, len(target_pres.slide_layouts) - 1)]

    dest = target_pres.slides.add_slide(layout)
    layout_rel_target = dest.part.part_related_by(RT.SLIDE_LAYOUT)

    _replace_slide_xml(dest, source)
    _copy_rels(source.part, dest.part, target_pres.part.package, clone_all=True)
    # The source's slideLayout rel was cloned into the target package by
    # clone_all; point it back at the target deck's own layout instead, so the
    # slide inherits the destination theme rather than dragging a duplicate
    # master along.
    for rel in list(dest.part.rels.values()):
        if rel.reltype == RT.SLIDE_LAYOUT:
            _set_rel(dest.part, rel.rId, rel.reltype, layout_rel_target)
    _copy_notes(source, dest)

    new_index = len(target_pres.slides) - 1
    if insert_after is not None:
        target = insert_after + 1 if insert_after < new_index else new_index
        _move_to(target_pres, new_index, target)
        new_index = target
    return new_index, getattr(layout, "name", None), matched


# ---- Speaker notes ----

def get_speaker_notes(slide):
    """Notes text for a slide, or "" when it has no notes slide."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def set_speaker_notes(slide, text):
    """Set (or clear, with "") a slide's speaker notes.

    Clearing a slide that has no notes is a no-op. `slide.notes_slide`
    *creates* the part on access, so clearing would otherwise add an empty
    notesSlide — and pull a notesMaster into the package with it — to every
    slide it was pointed at.
    """
    if not (text or slide.has_notes_slide):
        return
    slide.notes_slide.notes_text_frame.text = text or ""


# ---- Placeholder geometry ----

def pin_inherited_geometry(shape):
    """Write a placeholder's inherited position and size into its own `a:xfrm`.

    A placeholder with no ``<a:xfrm>`` takes all four values from the layout.
    Setting one of them — ``shape.width = ...`` — makes python-pptx create an
    ``<a:xfrm>`` holding only what was set: an ``<a:ext>`` with a zero for the
    extent nobody supplied, and no ``<a:off>`` at all. That partial transform
    stops the inheritance dead. LibreOffice falls back to the layout anyway and
    still draws the shape in place, so the render looks right; PowerPoint reads
    it literally and parks the shape at the top-left corner of the slide. It is
    the classic "renders fine in LibreOffice, lands in the corner in
    PowerPoint" defect, and it is invisible to visual QA for exactly that
    reason.

    Copying all four values across first turns any later geometry edit into an
    edit of a fully specified shape. Returns True when it materialized a
    transform, False when the shape already had one or is not a placeholder.
    """
    if not getattr(shape, "is_placeholder", False):
        return False
    # A graphic-frame placeholder (chart, table) keeps its transform in a
    # p:xfrm that python-pptx always writes in full, and has no spPr at all.
    spPr = getattr(shape._element, "spPr", None)
    if spPr is None or spPr.find(qn("a:xfrm")) is not None:
        return False
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    if None in (left, top, width, height):
        return False
    shape.left, shape.top = left, top
    shape.width, shape.height = width, height
    return True
